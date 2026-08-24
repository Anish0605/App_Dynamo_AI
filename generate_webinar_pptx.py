"""
Dynamo AI Webinar Presentation Generator
23-slide story-mode deck — light academic theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy, os

# ─── DESIGN TOKENS ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2E, 0x4B)   # primary / headings
GOLD    = RGBColor(0xEA, 0xB3, 0x08)   # Dynamo accent
DARK    = RGBColor(0x1F, 0x29, 0x37)   # body text
MED     = RGBColor(0x4B, 0x55, 0x63)   # supporting text
LIGHT   = RGBColor(0xF0, 0xF4, 0xFA)   # card / strip bg
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xF9, 0xFB)   # slide background
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)

LOGO_PATH = "frontend/assets/dynamo-logo.png"

# Slide dimensions: 16:9 widescreen (13.33" × 7.5")
W = Inches(13.33)
H = Inches(7.5)


# ─── HELPERS ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor = None, line: RGBColor = None, line_pt=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,   # rectangle
        x, y, w, h
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_pt:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_name="Arial", font_size=24, bold=False, italic=False,
                color: RGBColor = None, align=PP_ALIGN.LEFT,
                wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as _Pt
        p.line_spacing = line_spacing
    return txBox


def add_para(tf, text, font_name="Arial", font_size=18, bold=False,
             color: RGBColor = None, align=PP_ALIGN.LEFT, space_before=0, indent=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = _Pt(space_before)
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def add_logo(slide, size=Inches(0.5)):
    if not os.path.exists(LOGO_PATH):
        return
    x = Inches(0.25)
    y = H - Inches(0.72)
    slide.shapes.add_picture(LOGO_PATH, x, y, height=size)


def add_footer_bar(slide, text="Dynamo AI   |   Webinar 2026"):
    """Thin gold rule + footer text at bottom of slide."""
    bar_h = Inches(0.05)
    bar_y = H - Inches(0.6)
    add_rect(slide, 0, bar_y, W, bar_h, fill=GOLD)
    # footer text — no logo, just text from left edge
    add_textbox(slide, text,
                Inches(0.35), H - Inches(0.56), Inches(12.5), Inches(0.45),
                font_size=9, color=MED, align=PP_ALIGN.LEFT)


def slide_number_tag(slide, num):
    add_textbox(slide, str(num),
                W - Inches(0.55), H - Inches(0.52), Inches(0.4), Inches(0.35),
                font_size=9, color=MED, align=PP_ALIGN.RIGHT)


def section_label(slide, text):
    """Small navy capsule label top-right."""
    lw = Inches(2.2)
    lh = Inches(0.3)
    lx = W - lw - Inches(0.35)
    ly = Inches(0.22)
    r = add_rect(slide, lx, ly, lw, lh, fill=NAVY)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Arial"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = WHITE


def heading(slide, text, x, y, w, h, size=36, color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or NAVY
    return add_textbox(slide, text, x, y, w, h,
                       font_size=size, bold=True, color=color, align=align, wrap=wrap)


def gold_accent_line(slide, x, y, w):
    add_rect(slide, x, y, w, Inches(0.045), fill=GOLD)


def number_circle(slide, num_text, cx, cy, radius=Inches(0.42)):
    """Draw a filled navy circle with white number."""
    x = cx - radius
    y = cy - radius
    d = radius * 2
    shape = slide.shapes.add_shape(9, x, y, d, d)   # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    from pptx.oxml.ns import qn
    from pptx.util import Pt as _Pt
    # vertical centering
    tf.auto_size = None
    shape.text_frame.margin_top    = Inches(0.0)
    shape.text_frame.margin_bottom = Inches(0.0)
    run = p.add_run()
    run.text = num_text
    run.font.name  = "Arial"
    run.font.size  = Pt(18)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    return shape


def card_box(slide, x, y, w, h, fill=None):
    fill = fill or LIGHT
    r = add_rect(slide, x, y, w, h, fill=fill)
    r.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    r.line.width = Pt(0.5)
    return r


# ═══════════════════════════════════════════════════════════════════════════════
#   SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide01_welcome(prs):
    """Title slide — navy left panel + light right panel."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)

    # Left navy panel
    panel_w = Inches(5.6)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)

    # Gold accent strip on far left edge
    add_rect(slide, 0, 0, Inches(0.08), H, fill=GOLD)

    # Logo on left panel
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.5), height=Inches(0.85))

    # "WEBINAR 2026" label
    add_textbox(slide, "WEBINAR  2026",
                Inches(0.5), Inches(1.6), panel_w - Inches(0.8), Inches(0.35),
                font_size=9, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Main title on left
    add_textbox(slide, "Research Smarter,\nPublish Faster",
                Inches(0.5), Inches(2.0), panel_w - Inches(0.7), Inches(2.4),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT, wrap=True)

    # Gold rule
    add_rect(slide, Inches(0.5), Inches(4.55), Inches(2.6), Inches(0.05), fill=GOLD)

    # Subtitle on left
    add_textbox(slide,
                "The AI-Powered Future of Academic Research\nA practical workshop for researchers, PhD scholars & faculty",
                Inches(0.5), Inches(4.7), panel_w - Inches(0.7), Inches(1.0),
                font_size=13, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.LEFT, wrap=True)

    # Speaker info bottom-left
    add_textbox(slide, "Anish Krisna\nFounder, Dynamo AI",
                Inches(0.5), H - Inches(1.5), panel_w - Inches(0.7), Inches(0.9),
                font_size=12, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT, wrap=True)

    # Right side — decorative content
    rx = panel_w + Inches(0.5)
    rw = W - panel_w - Inches(0.7)

    # Big decorative quote
    add_textbox(slide,
                '"Researchers don\'t need another AI chatbot.\nThey need one intelligent\nresearch workspace."',
                rx, Inches(1.5), rw, Inches(2.4),
                font_size=22, italic=True, color=NAVY, align=PP_ALIGN.LEFT, wrap=True)

    gold_accent_line(slide, rx, Inches(1.45), Inches(3.2))

    # Three feature pills
    features = ["AI-Powered Research", "Cite & Verify Sources", "Gap Discovery"]
    for i, feat in enumerate(features):
        fy = Inches(4.1) + i * Inches(0.75)
        c = add_rect(slide, rx, fy, Inches(3.2), Inches(0.55), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        tf = c.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "  " + feat
        run.font.name = "Arial"
        run.font.size = Pt(13)
        run.font.color.rgb = NAVY

    return slide


def slide02_why_here(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Opening")
    slide_number_tag(slide, 2)

    heading(slide, "Why Are We Here Today?",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.85), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.45), Inches(3.0))

    add_textbox(slide, "How many of you have experienced at least one of these?",
                Inches(0.7), Inches(1.6), Inches(8), Inches(0.5),
                font_size=16, color=MED)

    checks = [
        "Spending weeks on literature review",
        "Not knowing if your research idea is truly novel",
        "Struggling with citations and references",
        "Fear of accidental plagiarism before submission",
        "Drowning in information overload",
        "Pressure to publish — limited time",
    ]
    col1 = checks[:3]
    col2 = checks[3:]
    cx1 = Inches(0.7)
    cx2 = Inches(7.0)
    cy_start = Inches(2.3)
    gap = Inches(0.78)

    for i, item in enumerate(col1):
        y = cy_start + i * gap
        c = add_rect(slide, cx1, y, Inches(5.8), Inches(0.62), fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        # gold tick bar
        add_rect(slide, cx1, y, Inches(0.07), Inches(0.62), fill=GOLD)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "    " + item
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = DARK

    for i, item in enumerate(col2):
        y = cy_start + i * gap
        c = add_rect(slide, cx2, y, Inches(5.8), Inches(0.62), fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, cx2, y, Inches(0.07), Inches(0.62), fill=GOLD)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "    " + item
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = DARK

    add_textbox(slide, "You are not alone.",
                Inches(0.7), H - Inches(1.4), Inches(8), Inches(0.5),
                font_size=18, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide


def slide03_reality(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "The Problem")
    slide_number_tag(slide, 3)

    heading(slide, "The Reality of Research Today",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    # Big stat
    add_textbox(slide, "80%",
                Inches(0.5), Inches(1.6), Inches(4.0), Inches(2.2),
                font_size=110, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "of research time is spent\nsearching, organizing & rewriting",
                Inches(0.5), Inches(3.7), Inches(4.2), Inches(0.9),
                font_size=15, color=MED, align=PP_ALIGN.CENTER, wrap=True)

    # Visual blocks — right side
    labels = [
        ("SEARCHING", Inches(6.0), Inches(1.6), Inches(2.8), Inches(1.0), NAVY),
        ("READING", Inches(9.1), Inches(1.6), Inches(2.0), Inches(1.0), NAVY),
        ("FORMATTING", Inches(6.0), Inches(2.8), Inches(2.0), Inches(1.0), NAVY),
        ("CITATIONS", Inches(8.2), Inches(2.8), Inches(2.9), Inches(1.0), NAVY),
    ]
    for label, lx, ly, lw, lh, lc in labels:
        b = add_rect(slide, lx, ly, lw, lh, fill=lc)
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Discovery block — small, gold
    b = add_rect(slide, Inches(6.0), Inches(4.0), Inches(5.1), Inches(0.75), fill=GOLD)
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ACTUAL DISCOVERY  —  only 20%"
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NAVY

    add_textbox(slide, "Researchers deserve to spend more time on discovery, less on management.",
                Inches(0.7), H - Inches(1.45), Inches(11.5), Inches(0.5),
                font_size=14, bold=True, italic=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide


def slide04_story(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "My Story")
    slide_number_tag(slide, 4)

    # Left accent
    add_rect(slide, 0, 0, Inches(0.1), H, fill=GOLD)

    heading(slide, "I Didn't Set Out to Build\nAnother AI Tool",
            Inches(0.55), Inches(0.4), Inches(7), Inches(1.5), size=34, wrap=True)
    gold_accent_line(slide, Inches(0.55), Inches(2.0), Inches(3.2))

    # Personal story — 2 columns
    story_left = [
        ("1 Conference Paper", "Took me 6 months — juggling tools, papers, and citations."),
        ("After Dynamo AI", "The same quality paper — done in 1 week. Published."),
    ]
    for i, (title, body) in enumerate(story_left):
        cy = Inches(2.2) + i * Inches(1.55)
        c = add_rect(slide, Inches(0.55), cy, Inches(5.6), Inches(1.3), fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, Inches(0.55), cy, Inches(0.1), Inches(1.3), fill=NAVY)
        add_textbox(slide, title, Inches(0.85), cy + Inches(0.12), Inches(5.1), Inches(0.38),
                    font_size=14, bold=True, color=NAVY)
        add_textbox(slide, body, Inches(0.85), cy + Inches(0.5), Inches(5.1), Inches(0.65),
                    font_size=13, color=MED, wrap=True)

    # Right — what I kept seeing
    add_textbox(slide, "I kept seeing the same pattern:",
                Inches(7.0), Inches(2.15), Inches(5.5), Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    patterns = [
        "Brilliant researchers spending hours finding papers",
        "Struggling to organize notes and citations",
        "Unsure whether their idea was truly novel",
        "Pressure to publish faster with limited time",
        "Jumping between multiple disconnected tools",
    ]
    for i, p in enumerate(patterns):
        py = Inches(2.65) + i * Inches(0.65)
        add_rect(slide, Inches(7.0), py + Inches(0.16), Inches(0.18), Inches(0.18), fill=GOLD)
        add_textbox(slide, p, Inches(7.35), py, Inches(5.15), Inches(0.55),
                    font_size=13, color=DARK, wrap=True)

    add_textbox(slide, "I built Dynamo AI so researchers spend more time discovering\nand less time managing.",
                Inches(0.55), H - Inches(1.4), Inches(12.0), Inches(0.65),
                font_size=14, bold=True, italic=True, color=NAVY, wrap=True)

    return slide


def slide05_toolkit(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "The Landscape")
    slide_number_tag(slide, 5)

    heading(slide, "The Current Research Toolkit",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    tools = [
        ("ChatGPT", "AI Writing"), ("Claude", "AI Research"), ("Gemini", "AI Analysis"),
        ("Perplexity", "Search AI"), ("Scite", "Citation Check"), ("Consensus", "Paper Search"),
        ("Zotero", "Reference Mgmt"), ("Mendeley", "PDF Library"), ("Turnitin", "Plagiarism"),
        ("Grammarly", "Writing"), ("EndNote", "Citations"), ("ResearchRabbit", "Discovery"),
    ]

    cols, rows = 4, 3
    tw = Inches(2.9)
    th = Inches(1.05)
    gx = Inches(0.22)
    gy = Inches(0.18)
    start_x = Inches(0.65)
    start_y = Inches(1.65)

    for idx, (name, cat) in enumerate(tools):
        col = idx % cols
        row = idx // cols
        tx = start_x + col * (tw + gx)
        ty = start_y + row * (th + gy)
        c = add_rect(slide, tx, ty, tw, th, fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_textbox(slide, name, tx + Inches(0.15), ty + Inches(0.12), tw - Inches(0.2), Inches(0.42),
                    font_size=14, bold=True, color=NAVY)
        add_textbox(slide, cat, tx + Inches(0.15), ty + Inches(0.55), tw - Inches(0.2), Inches(0.35),
                    font_size=11, color=MED)

    add_textbox(slide, "These are excellent tools — but researchers are forced to jump between all of them.",
                Inches(0.7), H - Inches(1.4), Inches(11.5), Inches(0.5),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide


def slide06_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "The Problem")
    slide_number_tag(slide, 6)

    heading(slide, "The Problem With Existing Tools",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    stats = [
        ("1", "Researcher"),
        ("10", "Tools"),
        ("20+", "Browser Tabs"),
    ]
    sw = Inches(3.6)
    sx = Inches(0.65)
    for i, (num, label) in enumerate(stats):
        cx = sx + i * (sw + Inches(0.2))
        cy = Inches(1.7)
        c = add_rect(slide, cx, cy, sw, Inches(2.4), fill=NAVY)
        add_textbox(slide, num, cx, cy + Inches(0.2), sw, Inches(1.5),
                    font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, cx, cy + Inches(1.7), sw, Inches(0.5),
                    font_size=16, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Pain points row
    pains = [
        ("Broken Workflow", "Context switching kills deep thinking"),
        ("Multiple Subscriptions", "Cost adds up — and nothing integrates"),
        ("No Single Source of Truth", "Work scattered across tools and tabs"),
    ]
    pw = Inches(3.6)
    for i, (title, body) in enumerate(pains):
        cx = sx + i * (pw + Inches(0.2))
        cy = Inches(4.3)
        c = add_rect(slide, cx, cy, pw, Inches(1.4), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_textbox(slide, title, cx + Inches(0.15), cy + Inches(0.1), pw - Inches(0.2), Inches(0.4),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, body, cx + Inches(0.15), cy + Inches(0.52), pw - Inches(0.2), Inches(0.7),
                    font_size=11, color=MED, wrap=True)

    add_textbox(slide, "Does this look familiar?",
                Inches(0.7), H - Inches(1.38), Inches(8), Inches(0.45),
                font_size=16, bold=True, italic=True, color=NAVY)

    return slide


def slide07_born(prs):
    """Quote + Vision — section transition feel."""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    add_footer_bar(slide)
    slide_number_tag(slide, 7)

    # Subtle gold rectangle decoration top-right
    add_rect(slide, W - Inches(1.5), 0, Inches(1.5), Inches(1.5), fill=GOLD)

    # Big quote
    add_textbox(slide,
                '"Researchers don\'t need\nanother AI chatbot.\nThey need one intelligent\nresearch workspace."',
                Inches(1.0), Inches(0.7), Inches(9.5), Inches(3.2),
                font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT, wrap=True)

    gold_accent_line(slide, Inches(1.0), Inches(4.0), Inches(2.6))

    # Vision labels
    add_textbox(slide, "Dynamo AI Vision",
                Inches(1.0), Inches(4.15), Inches(5), Inches(0.45),
                font_size=13, bold=True, color=GOLD)

    vision = ["Search", "Analyze", "Write", "Cite", "Present"]
    vx = Inches(1.0)
    for i, v in enumerate(vision):
        bx = vx + i * Inches(2.2)
        b = add_rect(slide, bx, Inches(4.7), Inches(1.9), Inches(0.6), fill=RGBColor(0x25, 0x40, 0x63))
        b.line.color.rgb = GOLD
        b.line.width = Pt(0.6)
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = v
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

    add_textbox(slide, "All in one place.",
                Inches(1.0), Inches(5.55), Inches(5), Inches(0.45),
                font_size=16, bold=True, color=GOLD)

    add_textbox(slide, "Let me show you how modern researchers work.",
                Inches(1.0), H - Inches(1.38), Inches(10), Inches(0.45),
                font_size=13, italic=True, color=RGBColor(0x94, 0xA3, 0xB8))

    return slide


def slide08_workflow(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "The Workflow")
    slide_number_tag(slide, 8)

    heading(slide, "The AI-Powered Research Workflow",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    steps = [
        "Research Question",
        "Literature Review",
        "Gap Identification",
        "Writing & Drafting",
        "Citation Validation",
        "Plagiarism Check",
        "Presentation",
        "Publication",
    ]
    sw = Inches(1.38)
    sh = Inches(0.72)
    gx = Inches(0.16)
    start_x = Inches(0.6)
    sy = Inches(2.0)

    for i, step in enumerate(steps):
        bx = start_x + i * (sw + gx)
        b = add_rect(slide, bx, sy, sw, sh, fill=NAVY if i != 6 else GOLD)
        tf = b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step
        run.font.name = "Arial"
        run.font.size = Pt(10.5)
        run.font.bold = True
        run.font.color.rgb = WHITE if i != 6 else NAVY

        if i < len(steps) - 1:
            ax = bx + sw + Inches(0.02)
            ay = sy + sh / 2 - Inches(0.03)
            add_textbox(slide, ">", ax, ay, gx + Inches(0.04), Inches(0.25),
                        font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Description blocks below
    descs = [
        ("Define your question", Inches(0.6), Inches(3.1), Inches(2.8)),
        ("Find existing knowledge", Inches(3.55), Inches(3.1), Inches(2.8)),
        ("Spot what's missing", Inches(6.5), Inches(3.1), Inches(2.8)),
        ("Draft with AI assistance", Inches(9.45), Inches(3.1), Inches(3.1)),
    ]
    for text, dx, dy, dw in descs:
        c = add_rect(slide, dx, dy, dw, Inches(0.85), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_textbox(slide, text, dx + Inches(0.12), dy + Inches(0.12), dw - Inches(0.2), Inches(0.6),
                    font_size=12, color=MED, wrap=True, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                "Dynamo AI supports every stage of this workflow — in one unified workspace.",
                Inches(0.7), H - Inches(1.42), Inches(11.5), Inches(0.5),
                font_size=14, bold=True, color=NAVY)

    return slide


def slide09_write_faster(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Writing")
    slide_number_tag(slide, 9)

    heading(slide, "From Blank Page to First Draft Faster",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    # Two columns
    col_w = Inches(5.6)
    # Left — traditional
    add_rect(slide, Inches(0.6), Inches(1.65), col_w, Inches(0.42), fill=RED)
    add_textbox(slide, "  Traditional Approach",
                Inches(0.6), Inches(1.65), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)

    trad_steps = [
        "Research phase",
        "Create outline",
        "First draft",
        "Rewrite once",
        "Rewrite again",
        "Rewrite (weeks of effort)...",
    ]
    for i, s in enumerate(trad_steps):
        add_textbox(slide, "  " + s,
                    Inches(0.6), Inches(2.2) + i * Inches(0.52), Inches(5.5), Inches(0.45),
                    font_size=13, color=MED, italic=(i > 2))

    # Right — AI assisted
    add_rect(slide, Inches(7.0), Inches(1.65), col_w, Inches(0.42), fill=GREEN)
    add_textbox(slide, "  AI-Assisted Approach",
                Inches(7.0), Inches(1.65), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)

    ai_steps = [
        "1.  Define objective",
        "2.  Generate structure",
        "3.  Draft section-by-section",
        "4.  Refine with your expertise",
        "5.  Finalize",
    ]
    for i, s in enumerate(ai_steps):
        add_textbox(slide, s,
                    Inches(7.0), Inches(2.2) + i * Inches(0.52), Inches(5.5), Inches(0.45),
                    font_size=13, color=DARK)

    # Example prompt box
    add_rect(slide, Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.85), fill=LIGHT)
    add_textbox(slide,
                'Example Prompt:  "Draft the Introduction for a paper on AI adoption in higher education using a formal academic tone."',
                Inches(0.75), Inches(5.25), Inches(11.7), Inches(0.7),
                font_size=12, italic=True, color=NAVY, wrap=True)

    return slide


def slide10_framework(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Framework")
    slide_number_tag(slide, 10)

    heading(slide, "The 5-Step Research Writing Framework",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.2))

    steps = [
        ("Define Objectives", "Clarify your research question and scope before writing a single word."),
        ("Generate Structure", "AI builds a logical outline — introduction, methods, discussion, conclusion."),
        ("Draft Strategically", "Write section-by-section with AI assistance, maintaining your academic voice."),
        ("Verify Evidence", "Cross-check claims against real papers and sources. Validate every fact."),
        ("Refine and Humanize", "Add your expertise, insights, and perspective. Make it genuinely yours."),
    ]

    bw = Inches(2.28)
    bh = Inches(2.3)
    gap = Inches(0.14)
    start_x = Inches(0.5)
    by = Inches(1.75)

    for i, (title, desc) in enumerate(steps):
        bx = start_x + i * (bw + gap)
        # numbered circle
        number_circle(slide, str(i + 1),
                      bx + bw / 2, by + Inches(0.38), radius=Inches(0.33))
        c = add_rect(slide, bx, by, bw, bh, fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        # Gold top strip
        add_rect(slide, bx, by, bw, Inches(0.055), fill=GOLD)
        # Number (redrawn on top of card)
        number_circle(slide, str(i + 1),
                      bx + bw / 2, by + Inches(0.46), radius=Inches(0.33))
        add_textbox(slide, title, bx + Inches(0.1), by + Inches(0.95), bw - Inches(0.2), Inches(0.52),
                    font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, wrap=True)
        add_textbox(slide, desc, bx + Inches(0.1), by + Inches(1.5), bw - Inches(0.2), Inches(0.7),
                    font_size=10.5, color=MED, align=PP_ALIGN.CENTER, wrap=True)

    add_textbox(slide,
                "This framework applies to journal articles, conference papers, theses, and research reports.",
                Inches(0.7), H - Inches(1.4), Inches(11.5), Inches(0.45),
                font_size=13, italic=True, color=MED)

    return slide


def slide11_litreview(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Literature Review")
    slide_number_tag(slide, 11)

    heading(slide, "Literature Reviews Without the Chaos",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    # Two columns
    col_w = Inches(5.6)

    # Traditional
    add_rect(slide, Inches(0.6), Inches(1.7), col_w, Inches(0.42), fill=RED)
    add_textbox(slide, "  Traditional Way",
                Inches(0.6), Inches(1.7), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)

    trad = [
        "Download 50+ PDFs manually",
        "Highlight and annotate by hand",
        "Lose track of key findings",
        "Start over when gaps appear",
        "Hours lost — weeks of effort",
    ]
    for i, t in enumerate(trad):
        add_textbox(slide, "    " + t,
                    Inches(0.6), Inches(2.28) + i * Inches(0.54), Inches(5.5), Inches(0.48),
                    font_size=13, color=MED)

    # Smarter way
    add_rect(slide, Inches(7.0), Inches(1.7), col_w, Inches(0.42), fill=GREEN)
    add_textbox(slide, "  The Smarter Way (Dynamo AI)",
                Inches(7.0), Inches(1.7), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)

    smart = [
        "Upload papers — Dynamo reads them all",
        "Extract key insights automatically",
        "Compare findings across papers",
        "Organize by theme with one click",
        "Build synthesis in minutes",
    ]
    for i, s in enumerate(smart):
        add_textbox(slide, "  " + s,
                    Inches(7.0), Inches(2.28) + i * Inches(0.54), Inches(5.5), Inches(0.48),
                    font_size=13, color=DARK)

    # Example
    add_rect(slide, Inches(0.6), Inches(5.05), Inches(12.0), Inches(0.7), fill=LIGHT)
    add_textbox(slide,
                'Example:  "Summarize the impact of Artificial Intelligence on Student Learning Outcomes" — themes identified in seconds.',
                Inches(0.75), Inches(5.1), Inches(11.7), Inches(0.58),
                font_size=12, italic=True, color=NAVY, wrap=True)

    return slide


def slide12_litreview_example(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Literature Review")
    slide_number_tag(slide, 12)

    heading(slide, "See the Bigger Picture Faster",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(2.8))

    add_textbox(slide, "Topic: Impact of AI on Higher Education",
                Inches(0.7), Inches(1.55), Inches(8), Inches(0.45),
                font_size=14, color=MED, italic=True)

    # Table
    headers = ["Theme", "Key Findings", "Implication"]
    rows_data = [
        ("Student Engagement", "Increased participation in AI-assisted courses", "Interactive tools improve attendance"),
        ("Learning Outcomes", "Mixed evidence — strong in STEM, weaker in humanities", "Context matters significantly"),
        ("Personalization", "Positive impact on self-paced learning", "AI tutors adapt to student pace"),
        ("Ethics & Privacy", "Privacy concerns around data collection", "Policy gaps need urgent attention"),
    ]

    tx = Inches(0.6)
    ty = Inches(2.15)
    tw_cols = [Inches(2.6), Inches(4.5), Inches(5.0)]

    # Header row
    for ci, (hdr, cw) in enumerate(zip(headers, tw_cols)):
        cx = tx + sum(tw_cols[:ci])
        b = add_rect(slide, cx, ty, cw, Inches(0.5), fill=NAVY)
        add_textbox(slide, hdr, cx + Inches(0.1), ty + Inches(0.07), cw - Inches(0.15), Inches(0.38),
                    font_size=13, bold=True, color=WHITE)

    for ri, row in enumerate(rows_data):
        ry = ty + Inches(0.5) + ri * Inches(0.72)
        fill = WHITE if ri % 2 == 0 else LIGHT
        for ci, (cell, cw) in enumerate(zip(row, tw_cols)):
            cx = tx + sum(tw_cols[:ci])
            b = add_rect(slide, cx, ry, cw, Inches(0.65), fill=fill)
            b.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
            b.line.width = Pt(0.4)
            if ci == 0:
                add_rect(slide, cx, ry, Inches(0.07), Inches(0.65), fill=GOLD)
            add_textbox(slide, cell, cx + Inches(0.15), ry + Inches(0.1), cw - Inches(0.2), Inches(0.5),
                        font_size=11, color=DARK, wrap=True)

    add_textbox(slide,
                "Instead of reading dozens of papers separately, researchers immediately identify patterns and contradictions.",
                Inches(0.7), H - Inches(1.4), Inches(11.5), Inches(0.5),
                font_size=13, italic=True, color=NAVY, wrap=True)

    return slide


def slide13_gaps(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Gap Discovery")
    slide_number_tag(slide, 13)

    heading(slide, "Finding Research Gaps Faster",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(2.8))

    # Left — framework steps
    add_textbox(slide, "Gap Discovery Framework",
                Inches(0.6), Inches(1.6), Inches(4.5), Inches(0.42),
                font_size=13, bold=True, color=NAVY)

    fsteps = [
        "Known Research",
        "Emerging Trends",
        "Contradictions",
        "Underrepresented Populations",
        "Methodological Limitations",
        "Research Gap Identified",
    ]
    step_colors = [NAVY, NAVY, NAVY, NAVY, NAVY, GOLD]
    for i, (s, sc) in enumerate(zip(fsteps, step_colors)):
        by = Inches(2.1) + i * Inches(0.66)
        b = add_rect(slide, Inches(0.6), by, Inches(3.8), Inches(0.54), fill=sc)
        add_textbox(slide, s, Inches(0.75), by + Inches(0.08), Inches(3.5), Inches(0.4),
                    font_size=12, bold=(i == 5), color=WHITE if sc != GOLD else NAVY)
        if i < len(fsteps) - 1:
            add_textbox(slide, "v", Inches(1.6), by + Inches(0.54), Inches(0.5), Inches(0.2),
                        font_size=10, color=GOLD, align=PP_ALIGN.CENTER)

    # Right — example
    add_textbox(slide, "Example: AI in Higher Education",
                Inches(5.2), Inches(1.6), Inches(7.5), Inches(0.42),
                font_size=13, bold=True, color=NAVY)

    # Two sub-boxes
    # Existing research
    eb = add_rect(slide, Inches(5.2), Inches(2.1), Inches(3.5), Inches(3.2), fill=LIGHT)
    eb.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    eb.line.width = Pt(0.5)
    add_textbox(slide, "Existing Research",
                Inches(5.35), Inches(2.15), Inches(3.2), Inches(0.38),
                font_size=12, bold=True, color=NAVY)
    existing = ["Undergraduate students", "Student engagement", "Learning outcomes", "Developed countries"]
    for i, e in enumerate(existing):
        add_textbox(slide, "+ " + e,
                    Inches(5.35), Inches(2.6) + i * Inches(0.52), Inches(3.2), Inches(0.45),
                    font_size=11, color=GREEN)

    # Missing areas
    mb = add_rect(slide, Inches(9.0), Inches(2.1), Inches(3.7), Inches(3.2), fill=LIGHT)
    mb.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    mb.line.width = Pt(0.5)
    add_textbox(slide, "Missing Areas",
                Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.38),
                font_size=12, bold=True, color=RED)
    missing = ["Rural institutions", "Faculty adoption barriers", "Long-term impacts", "Developing countries"]
    for i, m in enumerate(missing):
        add_textbox(slide, "? " + m,
                    Inches(9.15), Inches(2.6) + i * Inches(0.52), Inches(3.4), Inches(0.45),
                    font_size=11, color=RED)

    # Gap statement
    gap_b = add_rect(slide, Inches(5.2), Inches(5.45), Inches(7.5), Inches(0.75), fill=NAVY)
    add_textbox(slide,
                '"How do faculty in developing countries adopt AI-assisted teaching practices?"',
                Inches(5.35), Inches(5.5), Inches(7.2), Inches(0.65),
                font_size=12, bold=True, italic=True, color=GOLD, wrap=True)

    return slide


def slide14_gap_formula(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Gap Discovery")
    slide_number_tag(slide, 14)

    heading(slide, "A Simple Formula for Discovering Novel Ideas",
            Inches(0.7), Inches(0.45), Inches(11), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.2))

    # Formula elements
    elements = [
        ("Known Findings", "What the existing research has established"),
        ("Conflicting Results", "Where studies contradict each other"),
        ("Understudied Populations", "Groups not yet represented in research"),
        ("Emerging Trends", "New phenomena that need investigation"),
    ]
    ew = Inches(2.75)
    eh = Inches(2.0)
    ex = Inches(0.55)

    for i, (title, desc) in enumerate(elements):
        bx = ex + i * (ew + Inches(0.18))
        c = add_rect(slide, bx, Inches(1.7), ew, eh, fill=NAVY)
        add_textbox(slide, title, bx + Inches(0.12), Inches(1.82), ew - Inches(0.2), Inches(0.52),
                    font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER, wrap=True)
        add_textbox(slide, desc, bx + Inches(0.1), Inches(2.45), ew - Inches(0.18), Inches(0.75),
                    font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER, wrap=True)

        if i < len(elements) - 1:
            ox = bx + ew + Inches(0.02)
            oy = Inches(1.7) + eh / 2 - Inches(0.15)
            add_textbox(slide, "+", ox, oy, Inches(0.18), Inches(0.35),
                        font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Equals result
    add_textbox(slide, "=", Inches(5.8), Inches(3.85), Inches(0.6), Inches(0.55),
                font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    result = add_rect(slide, Inches(3.0), Inches(4.55), Inches(6.85), Inches(0.88), fill=GOLD)
    add_textbox(slide, "RESEARCH OPPORTUNITY",
                Inches(3.0), Inches(4.62), Inches(6.85), Inches(0.7),
                font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Dynamo AI's Research Gap feature applies this formula automatically — surfacing opportunities you might have missed.",
                Inches(0.7), H - Inches(1.38), Inches(11.5), Inches(0.48),
                font_size=13, italic=True, color=MED, wrap=True)

    return slide


def slide15_ai_workflow(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Product Demo")
    slide_number_tag(slide, 15)

    heading(slide, "AI Research Workflow in Practice",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_accent_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))

    steps = [
        ("Upload Papers", "PDF, DOCX, TXT — Dynamo reads everything"),
        ("Extract Themes", "Automatic topic clustering and synthesis"),
        ("Draft Content", "AI generates structured academic text"),
        ("Validate Sources", "Cross-check against real publications"),
    ]

    sw = Inches(2.85)
    sh = Inches(2.6)
    gap = Inches(0.22)
    sx = Inches(0.6)

    for i, (title, desc) in enumerate(steps):
        bx = sx + i * (sw + gap)
        c = add_rect(slide, bx, Inches(1.75), sw, sh, fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        # Step number
        nb = add_rect(slide, bx + Inches(0.15), Inches(1.9), Inches(0.45), Inches(0.45), fill=NAVY)
        add_textbox(slide, str(i + 1), bx + Inches(0.15), Inches(1.9), Inches(0.45), Inches(0.45),
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, bx + Inches(0.15), Inches(2.48), sw - Inches(0.25), Inches(0.52),
                    font_size=14, bold=True, color=NAVY, wrap=True)
        add_textbox(slide, desc, bx + Inches(0.15), Inches(3.05), sw - Inches(0.25), Inches(0.8),
                    font_size=12, color=MED, wrap=True)

        if i < len(steps) - 1:
            ax = bx + sw + Inches(0.04)
            ay = Inches(1.75) + sh / 2 - Inches(0.12)
            add_textbox(slide, ">", ax, ay, gap - Inches(0.04), Inches(0.35),
                        font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Result banner
    rb = add_rect(slide, Inches(0.6), Inches(4.6), Inches(12.0), Inches(0.72), fill=NAVY)
    add_textbox(slide, "Result: A fully structured, cited, verified research draft — in a fraction of the time",
                Inches(0.75), Inches(4.67), Inches(11.7), Inches(0.6),
                font_size=14, bold=True, color=GOLD, wrap=True, align=PP_ALIGN.CENTER)

    return slide


def slide16_assistant(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Features")
    slide_number_tag(slide, 16)

    # Left navy panel
    panel_w = Inches(4.8)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)

    add_textbox(slide, "FEATURE", Inches(0.5), Inches(0.8), Inches(4.0), Inches(0.35),
                font_size=9, bold=True, color=GOLD)
    heading(slide, "AI Research\nAssistant",
            Inches(0.5), Inches(1.2), panel_w - Inches(0.7), Inches(1.8),
            size=40, color=WHITE)
    gold_accent_line(slide, Inches(0.5), Inches(3.1), Inches(2.5))
    add_textbox(slide, "Chat directly with your research papers. Compare authors. Explore contradictions.",
                Inches(0.5), Inches(3.2), panel_w - Inches(0.7), Inches(1.0),
                font_size=14, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

    # Right side — use cases
    rx = panel_w + Inches(0.5)
    rw = W - panel_w - Inches(0.7)

    add_textbox(slide, "What you can do:",
                rx, Inches(0.8), rw, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    use_cases = [
        ("Chat with PDFs", "Upload any paper and ask questions in plain language"),
        ("Compare Authors", "See where researchers agree, disagree, and contradict"),
        ("Extract Key Arguments", "Pull the core claims from any document instantly"),
        ("Build Literature Matrix", "Organize findings across multiple papers visually"),
    ]
    for i, (uc_title, uc_body) in enumerate(use_cases):
        uy = Inches(1.35) + i * Inches(1.2)
        c = add_rect(slide, rx, uy, rw, Inches(1.0), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, rx, uy, Inches(0.07), Inches(1.0), fill=GOLD)
        add_textbox(slide, uc_title, rx + Inches(0.2), uy + Inches(0.1), rw - Inches(0.3), Inches(0.38),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, uc_body, rx + Inches(0.2), uy + Inches(0.5), rw - Inches(0.3), Inches(0.4),
                    font_size=12, color=MED, wrap=True)

    return slide


def slide17_detector(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Features")
    slide_number_tag(slide, 17)

    # Left navy panel
    panel_w = Inches(4.8)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)

    add_textbox(slide, "FEATURE", Inches(0.5), Inches(0.8), Inches(4.0), Inches(0.35),
                font_size=9, bold=True, color=GOLD)
    heading(slide, "AI Content\nDetector",
            Inches(0.5), Inches(1.2), panel_w - Inches(0.7), Inches(1.8),
            size=40, color=WHITE)
    gold_accent_line(slide, Inches(0.5), Inches(3.1), Inches(2.5))
    add_textbox(slide, "Review your writing for AI-generated signals before journal submission.",
                Inches(0.5), Inches(3.2), panel_w - Inches(0.7), Inches(1.0),
                font_size=14, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

    # Right side
    rx = panel_w + Inches(0.5)
    rw = W - panel_w - Inches(0.7)

    add_textbox(slide, "Built for academic integrity:",
                rx, Inches(0.8), rw, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    features = [
        ("AI Score Meter", "Visual 0-100 score showing AI-generation probability"),
        ("Writing Pattern Analysis", "Detects unnaturally uniform sentence structure and vocabulary"),
        ("Originality Review", "Plagiarism check against web sources and academic papers"),
        ("Actionable Guidance", "Specific signals highlighted so you can refine your writing"),
    ]
    for i, (f_title, f_body) in enumerate(features):
        fy = Inches(1.35) + i * Inches(1.2)
        c = add_rect(slide, rx, fy, rw, Inches(1.0), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, rx, fy, Inches(0.07), Inches(1.0), fill=GOLD)
        add_textbox(slide, f_title, rx + Inches(0.2), fy + Inches(0.1), rw - Inches(0.3), Inches(0.38),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, f_body, rx + Inches(0.2), fy + Inches(0.5), rw - Inches(0.3), Inches(0.4),
                    font_size=12, color=MED, wrap=True)

    return slide


def slide18_citation(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Features")
    slide_number_tag(slide, 18)

    # Left navy panel
    panel_w = Inches(4.8)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)

    add_textbox(slide, "FEATURE", Inches(0.5), Inches(0.8), Inches(4.0), Inches(0.35),
                font_size=9, bold=True, color=GOLD)
    heading(slide, "Citation\nChecker",
            Inches(0.5), Inches(1.2), panel_w - Inches(0.7), Inches(1.8),
            size=40, color=WHITE)
    gold_accent_line(slide, Inches(0.5), Inches(3.1), Inches(2.5))
    add_textbox(slide, "Fix missing, broken, and unverified references before they cost you a rejection.",
                Inches(0.5), Inches(3.2), panel_w - Inches(0.7), Inches(1.0),
                font_size=14, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

    rx = panel_w + Inches(0.5)
    rw = W - panel_w - Inches(0.7)

    add_textbox(slide, "What it checks:",
                rx, Inches(0.8), rw, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    features = [
        ("Missing Citations", "Identifies claims made without supporting references"),
        ("Broken References", "Flags citations that can't be verified in academic databases"),
        ("Format Validation", "Checks APA, MLA, Chicago, Harvard compliance automatically"),
        ("DOI Verification", "Confirms each paper actually exists and is retrievable"),
    ]
    for i, (f_title, f_body) in enumerate(features):
        fy = Inches(1.35) + i * Inches(1.2)
        c = add_rect(slide, rx, fy, rw, Inches(1.0), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, rx, fy, Inches(0.07), Inches(1.0), fill=GOLD)
        add_textbox(slide, f_title, rx + Inches(0.2), fy + Inches(0.1), rw - Inches(0.3), Inches(0.38),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, f_body, rx + Inches(0.2), fy + Inches(0.5), rw - Inches(0.3), Inches(0.4),
                    font_size=12, color=MED, wrap=True)

    return slide


def slide19_watcher(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Features")
    slide_number_tag(slide, 19)

    # Left navy panel
    panel_w = Inches(4.8)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)

    add_textbox(slide, "FEATURE", Inches(0.5), Inches(0.8), Inches(4.0), Inches(0.35),
                font_size=9, bold=True, color=GOLD)
    heading(slide, "Research\nWatcher",
            Inches(0.5), Inches(1.2), panel_w - Inches(0.7), Inches(1.8),
            size=40, color=WHITE)
    gold_accent_line(slide, Inches(0.5), Inches(3.1), Inches(2.5))
    add_textbox(slide, "Monitor emerging publications in your field — without manual searching.",
                Inches(0.5), Inches(3.2), panel_w - Inches(0.7), Inches(1.0),
                font_size=14, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

    rx = panel_w + Inches(0.5)
    rw = W - panel_w - Inches(0.7)

    add_textbox(slide, "Never miss a relevant paper again:",
                rx, Inches(0.8), rw, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    features = [
        ("Topic Alerts", "Set up keywords — get notified when new papers are published"),
        ("Author Tracking", "Follow specific researchers and see their latest work"),
        ("Journal Monitoring", "Track new issues from your target journals automatically"),
        ("Weekly Digest", "Curated summary of what's new in your research area"),
    ]
    for i, (f_title, f_body) in enumerate(features):
        fy = Inches(1.35) + i * Inches(1.2)
        c = add_rect(slide, rx, fy, rw, Inches(1.0), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, rx, fy, Inches(0.07), Inches(1.0), fill=GOLD)
        add_textbox(slide, f_title, rx + Inches(0.2), fy + Inches(0.1), rw - Inches(0.3), Inches(0.38),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, f_body, rx + Inches(0.2), fy + Inches(0.5), rw - Inches(0.3), Inches(0.4),
                    font_size=12, color=MED, wrap=True)

    return slide


def slide20_demo(prs):
    """Live demo placeholder — dark/transition slide."""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    add_footer_bar(slide)
    slide_number_tag(slide, 20)

    add_rect(slide, W - Inches(1.5), 0, Inches(1.5), Inches(1.5), fill=GOLD)

    add_textbox(slide, "LIVE DEMO", Inches(1.0), Inches(1.6), Inches(11), Inches(0.5),
                font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    heading(slide, "See Dynamo AI in Action",
            Inches(1.0), Inches(2.1), Inches(11), Inches(1.3),
            size=48, color=WHITE, align=PP_ALIGN.CENTER)

    gold_accent_line(slide, Inches(4.5), Inches(3.55), Inches(4.2))

    demo_steps = [
        "Upload a research paper",
        "Extract themes and summary",
        "Find research gaps",
        "Draft an Introduction section",
    ]
    for i, ds in enumerate(demo_steps):
        dx = Inches(1.0) + i * Inches(3.1)
        db = add_rect(slide, dx, Inches(3.85), Inches(2.85), Inches(0.65),
                      fill=RGBColor(0x25, 0x40, 0x63))
        db.line.color.rgb = GOLD
        db.line.width = Pt(0.6)
        add_textbox(slide, str(i + 1) + ".  " + ds, dx + Inches(0.1), Inches(3.88), Inches(2.7), Inches(0.52),
                    font_size=12, color=WHITE, wrap=True, align=PP_ALIGN.LEFT)

    add_textbox(slide, "app.dynamoai.in",
                Inches(1.0), H - Inches(1.6), Inches(11), Inches(0.5),
                font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    return slide


def slide21_validation(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, "Validation")
    slide_number_tag(slide, 21)

    add_rect(slide, 0, 0, Inches(0.1), H, fill=GOLD)

    heading(slide, "We Are Building Dynamo AI\nWith Researchers Like You",
            Inches(0.55), Inches(0.45), Inches(9), Inches(1.6),
            size=36, wrap=True)
    gold_accent_line(slide, Inches(0.55), Inches(2.15), Inches(3.2))

    add_textbox(slide,
                "Your feedback is not just welcome — it shapes every feature we build.",
                Inches(0.55), Inches(2.3), Inches(9), Inches(0.5),
                font_size=16, color=MED, wrap=True)

    points = [
        ("Your Voice Matters", "Every piece of feedback directly influences our product roadmap."),
        ("Real Research. Real Problems.", "We focus on what actually slows researchers down — not assumptions."),
        ("Built for Indian Academia", "Designed with the constraints, workflows, and goals of local researchers."),
        ("Continuous Improvement", "Weekly updates based on what our validation community tells us."),
    ]

    pw = Inches(5.6)
    for i, (title, body) in enumerate(points):
        px = Inches(0.55) if i % 2 == 0 else Inches(6.8)
        py = Inches(3.0) + (i // 2) * Inches(1.4)
        c = add_rect(slide, px, py, pw, Inches(1.2), fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        c.line.width = Pt(0.5)
        add_rect(slide, px, py, Inches(0.08), Inches(1.2), fill=NAVY)
        add_textbox(slide, title, px + Inches(0.22), py + Inches(0.1), pw - Inches(0.3), Inches(0.4),
                    font_size=14, bold=True, color=NAVY)
        add_textbox(slide, body, px + Inches(0.22), py + Inches(0.55), pw - Inches(0.3), Inches(0.55),
                    font_size=12, color=MED, wrap=True)

    return slide


def slide22_bonus(prs):
    """Special offer slide — gold background."""
    slide = blank_slide(prs)
    fill_bg(slide, GOLD)
    add_footer_bar(slide)
    slide_number_tag(slide, 22)

    add_rect(slide, 0, 0, Inches(0.12), H, fill=NAVY)

    add_textbox(slide, "EXCLUSIVE FOR WEBINAR ATTENDEES",
                Inches(0.5), Inches(0.65), Inches(12), Inches(0.42),
                font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    heading(slide, "14 Days of Dynamo AI Pro",
            Inches(0.5), Inches(1.1), Inches(12), Inches(1.3),
            size=52, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Completely Free — No Strings Attached",
                Inches(0.5), Inches(2.45), Inches(12), Inches(0.5),
                font_size=20, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.0), Inches(3.05), Inches(5.3), Inches(0.05), fill=NAVY)

    # What's included
    included = [
        "300 AI Research Messages / Day",
        "100 AI-Generated Images / Month",
        "Full Literature Review & Gap Discovery",
        "AI Detector + Plagiarism Checker",
        "Citation Checker + Research Watcher",
        "Priority Support",
    ]
    inc_w = Inches(3.8)
    for i, item in enumerate(included):
        col = i % 2
        row = i // 2
        ix = Inches(2.1) + col * (inc_w + Inches(0.5))
        iy = Inches(3.3) + row * Inches(0.58)
        add_rect(slide, ix, iy + Inches(0.15), Inches(0.18), Inches(0.18), fill=NAVY)
        add_textbox(slide, "      " + item, ix, iy, inc_w, Inches(0.5),
                    font_size=14, bold=False, color=NAVY)

    add_textbox(slide, "How to claim: Fill out the feedback form at the end of this session",
                Inches(0.5), H - Inches(1.42), Inches(12), Inches(0.45),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    return slide


def slide23_thankyou(prs):
    """Closing slide — mirrors title slide treatment."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)

    panel_w = Inches(5.6)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)
    add_rect(slide, 0, 0, Inches(0.08), H, fill=GOLD)

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.5), height=Inches(0.85))

    add_textbox(slide, "THANK YOU",
                Inches(0.5), Inches(1.7), panel_w - Inches(0.7), Inches(0.6),
                font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    heading(slide, "From Research\nOverwhelm to\nResearch Confidence",
            Inches(0.5), Inches(2.35), panel_w - Inches(0.65), Inches(2.6),
            size=32, color=WHITE)

    add_rect(slide, Inches(0.5), Inches(5.05), Inches(2.6), Inches(0.05), fill=GOLD)

    add_textbox(slide, "Anish Krisna\nFounder, Dynamo AI",
                Inches(0.5), Inches(5.2), panel_w - Inches(0.7), Inches(0.7),
                font_size=13, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

    # Right side — next steps
    rx = panel_w + Inches(0.55)
    rw = W - panel_w - Inches(0.75)

    heading(slide, "What happens next?",
            rx, Inches(0.8), rw, Inches(0.5),
            size=20, color=NAVY, align=PP_ALIGN.LEFT)
    gold_accent_line(slide, rx, Inches(1.38), Inches(2.8))

    steps = [
        ("1.", "Fill out the feedback form in the chat"),
        ("2.", "Receive your Pro invite code privately"),
        ("3.", "Visit app.dynamoai.in and sign up"),
        ("4.", "Enter your code — start your 14-day Pro trial"),
        ("5.", "Share your experience — help us build better"),
    ]
    for i, (num, txt) in enumerate(steps):
        sy = Inches(1.55) + i * Inches(0.75)
        add_textbox(slide, num, rx, sy, Inches(0.4), Inches(0.55),
                    font_size=18, bold=True, color=GOLD)
        add_textbox(slide, txt, rx + Inches(0.4), sy + Inches(0.05), rw - Inches(0.45), Inches(0.5),
                    font_size=14, color=DARK, wrap=True)

    add_rect(slide, rx, H - Inches(1.55), rw, Inches(0.62), fill=LIGHT)
    add_textbox(slide, "app.dynamoai.in",
                rx + Inches(0.15), H - Inches(1.5), rw - Inches(0.2), Inches(0.52),
                font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#   MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide01_welcome(prs)
    slide02_why_here(prs)
    slide03_reality(prs)
    slide04_story(prs)
    slide05_toolkit(prs)
    slide06_problem(prs)
    slide07_born(prs)
    slide08_workflow(prs)
    slide09_write_faster(prs)
    slide10_framework(prs)
    slide11_litreview(prs)
    slide12_litreview_example(prs)
    slide13_gaps(prs)
    slide14_gap_formula(prs)
    slide15_ai_workflow(prs)
    slide16_assistant(prs)
    slide17_detector(prs)
    slide18_citation(prs)
    slide19_watcher(prs)
    slide20_demo(prs)
    slide21_validation(prs)
    slide22_bonus(prs)
    slide23_thankyou(prs)

    out = "Dynamo_AI_Webinar_2026.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
