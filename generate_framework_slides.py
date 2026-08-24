"""
Dynamo AI — 5-Step Framework Deep-Dive Slides
One detailed slide per step, matching the Webinar 2026 design exactly.
Output: Dynamo_AI_Framework_5Steps.pptx  (Google Slides compatible)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── DESIGN TOKENS (identical to generate_webinar_pptx.py) ─────────────────────
NAVY  = RGBColor(0x1B, 0x2E, 0x4B)
GOLD  = RGBColor(0xEA, 0xB3, 0x08)
DARK  = RGBColor(0x1F, 0x29, 0x37)
MED   = RGBColor(0x4B, 0x55, 0x63)
LIGHT = RGBColor(0xF0, 0xF4, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG    = RGBColor(0xF8, 0xF9, 0xFB)

LOGO_PATH = "frontend/assets/dynamo-logo.png"
W = Inches(13.33)
H = Inches(7.5)

# Left panel ends here — titles must not cross this boundary
LEFT_PANEL_END = Inches(6.3)
RIGHT_X        = Inches(6.85)
RIGHT_W        = Inches(6.1)


# ── HELPERS ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_pt:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_name="Arial", font_size=24, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_footer_bar(slide):
    add_rect(slide, 0, H - Inches(0.6), W, Inches(0.05), fill=GOLD)
    add_textbox(slide, "Dynamo AI   |   Webinar 2026",
                Inches(0.35), H - Inches(0.56), Inches(12.5), Inches(0.45),
                font_size=9, color=MED, align=PP_ALIGN.LEFT)


def slide_number_tag(slide, num):
    add_textbox(slide, str(num),
                W - Inches(0.55), H - Inches(0.52), Inches(0.4), Inches(0.35),
                font_size=9, color=MED, align=PP_ALIGN.RIGHT)


def section_label(slide, text):
    lw, lh = Inches(2.2), Inches(0.3)
    lx, ly = W - lw - Inches(0.35), Inches(0.22)
    r = add_rect(slide, lx, ly, lw, lh, fill=NAVY)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Arial"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = WHITE


def number_circle(slide, num_text, cx, cy, radius=Inches(0.42)):
    x, y = cx - radius, cy - radius
    d = radius * 2
    shape = slide.shapes.add_shape(9, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.auto_size = None
    shape.text_frame.margin_top    = Inches(0.0)
    shape.text_frame.margin_bottom = Inches(0.0)
    run = p.add_run()
    run.text = num_text
    run.font.name  = "Arial"
    run.font.size  = Pt(18)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    return shape


def gold_accent_line(slide, x, y, w):
    add_rect(slide, x, y, w, Inches(0.045), fill=GOLD)


# ── SLIDE CONTENT ──────────────────────────────────────────────────────────────
# bullets = list of (short title, one crisp example line)
# Max 4 bullets — kept short so nothing overflows

STEPS = [
    {
        "num":    "1",
        "title":  "Define Objectives",
        "tag":    "Step 1 of 5",
        "desc":   (
            "Clarify your research question and scope before writing "
            "a single word. A sharp objective saves weeks of rework."
        ),
        "bullets": [
            ("One-sentence research question",
             'e.g. "How does AI adoption affect student outcomes in Indian higher education?"'),
            ("Define the scope",
             "Field, time period, population — be specific, not broad."),
            ("Identify your gap",
             "What does existing literature NOT answer? That is your contribution."),
            ("Set your success criteria",
             "Target journal, word count, deadline. Know what done looks like."),
        ],
        "outcome": "Key Outcome: A clear, testable research objective written in one sentence.",
        "tip":     "Narrow scope = deeper contribution. Broad topics make weak papers.",
    },
    {
        "num":    "2",
        "title":  "Generate Structure",
        "tag":    "Step 2 of 5",
        "desc":   (
            "Use AI to build a complete outline before writing anything. "
            "A solid structure prevents writer's block and keeps every "
            "section aligned with your objective."
        ),
        "bullets": [
            ("Generate a full outline with AI",
             'Prompt: "Outline a journal article on [topic] for [target journal]."'),
            ("Standard sections to include",
             "Introduction · Literature Review · Methodology · Results · Discussion · Conclusion"),
            ("Customise for your paper type",
             "Conference papers are 8 pages. Theses have separate chapters. Adjust accordingly."),
            ("Assign target word counts",
             "e.g. Introduction 600w · Methods 800w · Discussion 1,200w"),
        ],
        "outcome": "Key Outcome: A section-by-section blueprint ready for drafting in Step 3.",
        "tip":     "Edit the outline until every section feels essential — AI is the scaffolding, you are the architect.",
    },
    {
        "num":    "3",
        "title":  "Draft Strategically",
        "tag":    "Step 3 of 5",
        "desc":   (
            "Write one section at a time using AI assistance. "
            "The goal here is a complete rough draft — not perfection. "
            "Speed and coverage now, polish in Step 5."
        ),
        "bullets": [
            ("One section at a time",
             'Prompt: "Draft the Introduction (600 words) for a paper on [topic], formal academic tone."'),
            ("Insert your own data",
             "AI cannot know your results. Paste in your numbers, tables, and observations directly."),
            ("Mark citations as placeholders",
             "Write [CITE] wherever a reference is needed. Verify all sources in Step 4."),
            ("Don't edit while drafting",
             "Finish the full draft first — editing mid-draft kills momentum."),
        ],
        "outcome": "Key Outcome: A complete rough draft covering every section from introduction to conclusion.",
        "tip":     "A bad first draft you can edit is infinitely better than a blank page.",
    },
    {
        "num":    "4",
        "title":  "Verify Evidence",
        "tag":    "Step 4 of 5",
        "desc":   (
            "Every claim needs a real, verifiable source. This is the "
            "step most researchers skip — and the one that decides "
            "whether a paper passes peer review."
        ),
        "bullets": [
            ("Resolve every [CITE] placeholder",
             "Search Dynamo AI Citation Checker or Google Scholar for each claim."),
            ("Verify AI-suggested references",
             "AI hallucinates citations. Check every DOI before including it."),
            ("Confirm all statistics",
             'e.g. "80% of researchers…" — find the original source or remove the stat.'),
            ("Run a plagiarism check",
             "Use Dynamo AI's detector. Target below 15% similarity before submission."),
        ],
        "outcome": "Key Outcome: Every claim in the paper has a verified, real, citeable source.",
        "tip":     "One fabricated citation can get a paper retracted. Verification is non-negotiable.",
    },
    {
        "num":    "5",
        "title":  "Refine and Humanize",
        "tag":    "Step 5 of 5",
        "desc":   (
            "Transform a well-structured, evidenced draft into a paper "
            "that is genuinely yours. Add your voice, your insight, "
            "your expertise. This is what separates accepted from rejected."
        ),
        "bullets": [
            ("Add your unique perspective",
             "Where does your fieldwork, data, or expertise say something AI cannot?"),
            ("Rephrase in your voice",
             "Read each paragraph aloud. Rewrite anything that doesn't sound like you."),
            ("Check argument flow",
             "Read the full paper in one sitting. Does every section build on the last?"),
            ("Final proofread",
             "Tense, abbreviations, figure captions, reference formatting — check each once."),
        ],
        "outcome": "Key Outcome: A submission-ready paper that reflects your voice and original thinking.",
        "tip":     "The final polish separates good papers from great ones. Reserve time for this step.",
    },
]


# ── SLIDE BUILDER ──────────────────────────────────────────────────────────────

def build_step_slide(prs, step, slide_num):
    slide = blank_slide(prs)
    fill_bg(slide, BG)
    add_footer_bar(slide)
    section_label(slide, step["tag"])
    slide_number_tag(slide, slide_num)

    # ── TOP BREADCRUMB ─────────────────────────────────────────────────────────
    add_textbox(slide, f"Step {step['num']}  —  The 5-Step Research Writing Framework",
                Inches(0.7), Inches(0.26), Inches(9.0), Inches(0.35),
                font_size=10, color=MED)

    # ── NUMBER CIRCLE ──────────────────────────────────────────────────────────
    circle_cx = Inches(1.22)
    circle_cy = Inches(1.42)
    number_circle(slide, step["num"], circle_cx, circle_cy, radius=Inches(0.44))

    # ── TITLE — constrained to left panel (max end x = LEFT_PANEL_END) ─────────
    title_x = Inches(2.1)
    title_w = LEFT_PANEL_END - title_x          # ~Inches(4.2) — never overflows
    add_textbox(slide, step["title"],
                title_x, Inches(0.98), title_w, Inches(1.0),
                font_size=32, bold=True, color=NAVY,
                align=PP_ALIGN.LEFT, wrap=True)

    # ── GOLD RULE ──────────────────────────────────────────────────────────────
    gold_accent_line(slide, Inches(0.7), Inches(2.02), Inches(3.2))

    # ── DESCRIPTION ────────────────────────────────────────────────────────────
    add_textbox(slide, step["desc"],
                Inches(0.7), Inches(2.15), Inches(5.7), Inches(1.45),
                font_size=13, color=DARK, wrap=True)

    # ── RIGHT PANEL HEADER ─────────────────────────────────────────────────────
    add_textbox(slide, "How to execute this step:",
                RIGHT_X, Inches(0.56), RIGHT_W, Inches(0.38),
                font_size=10.5, bold=True, color=NAVY)

    gold_accent_line(slide, RIGHT_X, Inches(0.97), Inches(2.6))

    # ── BULLET CARDS ──────────────────────────────────────────────────────────
    for i, (point_title, point_body) in enumerate(step["bullets"]):
        by = Inches(1.12) + i * Inches(1.02)

        card = add_rect(slide, RIGHT_X, by, RIGHT_W, Inches(0.88), fill=WHITE)
        card.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        card.line.width = Pt(0.5)
        add_rect(slide, RIGHT_X, by, Inches(0.055), Inches(0.88), fill=GOLD)

        add_textbox(slide, point_title,
                    RIGHT_X + Inches(0.14), by + Inches(0.08),
                    RIGHT_W - Inches(0.2), Inches(0.32),
                    font_size=12, bold=True, color=NAVY, wrap=True)

        add_textbox(slide, point_body,
                    RIGHT_X + Inches(0.14), by + Inches(0.42),
                    RIGHT_W - Inches(0.2), Inches(0.38),
                    font_size=10.5, italic=True, color=MED, wrap=True)

    # ── KEY OUTCOME BAR ────────────────────────────────────────────────────────
    callout_y = Inches(5.38)
    callout = add_rect(slide, Inches(0.7), callout_y, Inches(12.2), Inches(0.64), fill=NAVY)
    callout.line.fill.background()
    add_textbox(slide, step["outcome"],
                Inches(0.95), callout_y + Inches(0.07),
                Inches(11.7), Inches(0.5),
                font_size=12.5, bold=True, color=WHITE, wrap=True)

    # ── TIP LINE ──────────────────────────────────────────────────────────────
    add_textbox(slide, "Tip:  " + step["tip"],
                Inches(0.7), Inches(6.18), Inches(12.0), Inches(0.38),
                font_size=11, italic=True, color=MED, wrap=True)

    return slide


def build():
    prs = new_prs()
    for i, step in enumerate(STEPS):
        build_step_slide(prs, step, slide_num=i + 1)
    out = "Dynamo_AI_Framework_5Steps.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
