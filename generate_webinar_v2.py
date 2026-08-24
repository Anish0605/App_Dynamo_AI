"""
Dynamo AI Webinar 2026 — UPDATED DECK (v2)
27 slides — all 6 review fixes applied + Google Slides safe
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── DESIGN TOKENS ──────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1B, 0x2E, 0x4B)
GOLD  = RGBColor(0xEA, 0xB3, 0x08)
DARK  = RGBColor(0x1F, 0x29, 0x37)
MED   = RGBColor(0x4B, 0x55, 0x63)
LIGHT = RGBColor(0xF0, 0xF4, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG    = RGBColor(0xF8, 0xF9, 0xFB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED   = RGBColor(0xDC, 0x26, 0x26)
SLATE = RGBColor(0xCB, 0xD5, 0xE1)

LOGO_PATH = "frontend/assets/dynamo-logo.png"
W = Inches(13.33)
H = Inches(7.5)


# ── HELPERS ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_pt:
        shape.line.color.rgb = line; shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_name="Arial", font_size=24, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font_name; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color
    return txBox

def add_footer_bar(slide, text="Dynamo AI   |   Webinar 2026"):
    add_rect(slide, 0, H - Inches(0.6), W, Inches(0.05), fill=GOLD)
    add_textbox(slide, text, Inches(0.35), H - Inches(0.56), Inches(12.5), Inches(0.45),
                font_size=9, color=MED, align=PP_ALIGN.LEFT)

def slide_num(slide, num):
    add_textbox(slide, str(num), W - Inches(0.55), H - Inches(0.52), Inches(0.4), Inches(0.35),
                font_size=9, color=MED, align=PP_ALIGN.RIGHT)

def section_label(slide, text):
    lw, lh = Inches(2.2), Inches(0.3)
    lx, ly = W - lw - Inches(0.35), Inches(0.22)
    r = add_rect(slide, lx, ly, lw, lh, fill=NAVY)
    tf = r.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text.upper()
    run.font.name = "Arial"; run.font.size = Pt(8)
    run.font.bold = True; run.font.color.rgb = WHITE

def heading(slide, text, x, y, w, h, size=36, color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or NAVY
    return add_textbox(slide, text, x, y, w, h,
                       font_size=size, bold=True, color=color, align=align, wrap=wrap)

def gold_line(slide, x, y, w):
    add_rect(slide, x, y, w, Inches(0.045), fill=GOLD)

def num_circle(slide, num_text, cx, cy, radius=Inches(0.42)):
    x, y = cx - radius, cy - radius; d = radius * 2
    shape = slide.shapes.add_shape(9, x, y, d, d)
    shape.fill.solid(); shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    tf.auto_size = None
    shape.text_frame.margin_top    = Inches(0.0)
    shape.text_frame.margin_bottom = Inches(0.0)
    run = p.add_run(); run.text = num_text
    run.font.name = "Arial"; run.font.size = Pt(18)
    run.font.bold = True; run.font.color.rgb = WHITE
    return shape

def card(slide, x, y, w, h, fill=None):
    fill = fill or LIGHT
    r = add_rect(slide, x, y, w, h, fill=fill)
    r.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); r.line.width = Pt(0.5)
    return r

def left_strip(slide, x, y, h, color=None):
    add_rect(slide, x, y, Inches(0.07), h, fill=color or GOLD)

def navy_panel_slide(prs, num, section, feature_title, subtitle, right_label, items):
    """Reusable navy-left-panel feature slide (slides 16, 18)."""
    slide = blank_slide(prs)
    fill_bg(slide, BG); add_footer_bar(slide)
    section_label(slide, section); slide_num(slide, num)
    panel_w = Inches(4.8)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)
    add_textbox(slide, "FEATURE", Inches(0.5), Inches(0.8), Inches(4.0), Inches(0.35),
                font_size=9, bold=True, color=GOLD)
    heading(slide, feature_title, Inches(0.5), Inches(1.2), panel_w - Inches(0.7),
            Inches(1.8), size=40, color=WHITE)
    gold_line(slide, Inches(0.5), Inches(3.1), Inches(2.5))
    add_textbox(slide, subtitle, Inches(0.5), Inches(3.2), panel_w - Inches(0.7), Inches(1.0),
                font_size=14, color=SLATE, wrap=True)
    rx = panel_w + Inches(0.5); rw = W - panel_w - Inches(0.7)
    add_textbox(slide, right_label, rx, Inches(0.8), rw, Inches(0.4),
                font_size=14, bold=True, color=NAVY)
    for i, (t, b) in enumerate(items):
        fy = Inches(1.35) + i * Inches(1.2)
        c = card(slide, rx, fy, rw, Inches(1.0))
        left_strip(slide, rx, fy, Inches(1.0))
        add_textbox(slide, t, rx + Inches(0.2), fy + Inches(0.1), rw - Inches(0.3), Inches(0.38),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, b, rx + Inches(0.2), fy + Inches(0.5), rw - Inches(0.3), Inches(0.4),
                    font_size=12, color=MED, wrap=True)
    return slide


# ══ SLIDE BUILDERS ═════════════════════════════════════════════════════════════

def slide01_welcome(prs):
    slide = blank_slide(prs); fill_bg(slide, WHITE)
    panel_w = Inches(5.6)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)
    add_rect(slide, 0, 0, Inches(0.08), H, fill=GOLD)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.5), height=Inches(0.85))
    add_textbox(slide, "WEBINAR  2026", Inches(0.5), Inches(1.6), panel_w - Inches(0.8),
                Inches(0.35), font_size=9, bold=True, color=GOLD)
    add_textbox(slide, "Research Smarter,\nPublish Faster",
                Inches(0.5), Inches(2.0), panel_w - Inches(0.7), Inches(2.4),
                font_size=40, bold=True, color=WHITE, wrap=True)
    add_rect(slide, Inches(0.5), Inches(4.55), Inches(2.6), Inches(0.05), fill=GOLD)
    add_textbox(slide,
                "The AI-Powered Future of Academic Research\nA practical workshop for researchers, PhD scholars & faculty",
                Inches(0.5), Inches(4.7), panel_w - Inches(0.7), Inches(1.0),
                font_size=13, color=SLATE, wrap=True)
    add_textbox(slide, "Anish Krisna\nFounder, Dynamo AI",
                Inches(0.5), H - Inches(1.5), panel_w - Inches(0.7), Inches(0.9),
                font_size=12, color=RGBColor(0x94, 0xA3, 0xB8), wrap=True)
    rx = panel_w + Inches(0.5); rw = W - panel_w - Inches(0.7)
    add_textbox(slide,
                '"Researchers don\'t need another AI chatbot.\nThey need one intelligent\nresearch workspace."',
                rx, Inches(1.5), rw, Inches(2.4),
                font_size=22, italic=True, color=NAVY, wrap=True)
    gold_line(slide, rx, Inches(1.45), Inches(3.2))
    for i, feat in enumerate(["AI-Powered Research", "Cite & Verify Sources", "Gap Discovery"]):
        fy = Inches(4.1) + i * Inches(0.75)
        c = add_rect(slide, rx, fy, Inches(3.2), Inches(0.55), fill=LIGHT)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); c.line.width = Pt(0.5)
        tf = c.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = "  " + feat
        run.font.name = "Arial"; run.font.size = Pt(13); run.font.color.rgb = NAVY
    return slide


def slide02_why_here(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Opening"); slide_num(slide, 2)
    heading(slide, "Why Are We Here Today?", Inches(0.7), Inches(0.5), Inches(9), Inches(0.85), size=34)
    gold_line(slide, Inches(0.7), Inches(1.45), Inches(3.0))
    add_textbox(slide, "How many of you have experienced at least one of these?",
                Inches(0.7), Inches(1.6), Inches(8), Inches(0.5), font_size=16, color=MED)
    checks = [
        "Spending weeks on literature review",
        "Not knowing if your research idea is truly novel",
        "Struggling with citations and references",
        "Fear of accidental plagiarism before submission",
        "Drowning in information overload",
        "Pressure to publish — limited time",
    ]
    for i, item in enumerate(checks[:3]):
        y = Inches(2.3) + i * Inches(0.78)
        c = add_rect(slide, Inches(0.7), y, Inches(5.8), Inches(0.62), fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); c.line.width = Pt(0.5)
        add_rect(slide, Inches(0.7), y, Inches(0.07), Inches(0.62), fill=GOLD)
        add_textbox(slide, "    " + item, Inches(0.7), y, Inches(5.8), Inches(0.62),
                    font_size=14, color=DARK)
    for i, item in enumerate(checks[3:]):
        y = Inches(2.3) + i * Inches(0.78)
        c = add_rect(slide, Inches(7.0), y, Inches(5.8), Inches(0.62), fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); c.line.width = Pt(0.5)
        add_rect(slide, Inches(7.0), y, Inches(0.07), Inches(0.62), fill=GOLD)
        add_textbox(slide, "    " + item, Inches(7.0), y, Inches(5.8), Inches(0.62),
                    font_size=14, color=DARK)
    add_textbox(slide, "You are not alone.", Inches(0.7), H - Inches(1.4), Inches(8), Inches(0.5),
                font_size=18, bold=True, color=NAVY)
    return slide


def slide03_reality(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "The Problem"); slide_num(slide, 3)
    heading(slide, "The Reality of Research Today",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    add_textbox(slide, "80%", Inches(0.5), Inches(1.6), Inches(4.0), Inches(2.2),
                font_size=110, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "of research time is spent\nsearching, organizing & rewriting",
                Inches(0.5), Inches(3.7), Inches(4.2), Inches(0.9),
                font_size=15, color=MED, align=PP_ALIGN.CENTER, wrap=True)
    for label, lx, ly, lw, lh in [
        ("SEARCHING",   Inches(6.0), Inches(1.6), Inches(2.8), Inches(1.0)),
        ("READING",     Inches(9.1), Inches(1.6), Inches(2.0), Inches(1.0)),
        ("FORMATTING",  Inches(6.0), Inches(2.8), Inches(2.0), Inches(1.0)),
        ("CITATIONS",   Inches(8.2), Inches(2.8), Inches(2.9), Inches(1.0)),
    ]:
        b = add_rect(slide, lx, ly, lw, lh, fill=NAVY)
        p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = label
        run.font.name = "Arial"; run.font.size = Pt(14)
        run.font.bold = True; run.font.color.rgb = WHITE
    b = add_rect(slide, Inches(6.0), Inches(4.0), Inches(5.1), Inches(0.75), fill=GOLD)
    p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "ACTUAL DISCOVERY  —  only 20%"
    run.font.name = "Arial"; run.font.size = Pt(14)
    run.font.bold = True; run.font.color.rgb = NAVY
    add_textbox(slide, "Researchers deserve to spend more time on discovery, less on management.",
                Inches(0.7), H - Inches(1.45), Inches(11.5), Inches(0.5),
                font_size=14, bold=True, italic=True, color=NAVY)
    return slide


# FIX 2a — Slide 4 split: personal story only
def slide04a_story(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "My Story"); slide_num(slide, 4)
    add_rect(slide, 0, 0, Inches(0.1), H, fill=GOLD)
    heading(slide, "I Didn't Set Out to Build\nAnother AI Tool",
            Inches(0.55), Inches(0.4), Inches(9), Inches(1.5), size=34, wrap=True)
    gold_line(slide, Inches(0.55), Inches(2.0), Inches(3.2))
    add_textbox(slide, "One researcher's frustration became a platform.",
                Inches(0.55), Inches(2.15), Inches(9), Inches(0.42),
                font_size=15, italic=True, color=MED)
    story = [
        ("1 Conference Paper",
         "It took me 6 months — juggling 12 different tools, chasing citations, and rewriting sections from scratch."),
        ("After Building Dynamo AI",
         "The same quality paper done in 1 week. Submitted. Accepted. Published."),
        ("The Turning Point",
         "I realized the bottleneck wasn't intelligence — it was workflow. The research process itself was broken."),
    ]
    for i, (title, body) in enumerate(story):
        cy = Inches(2.75) + i * Inches(1.3)
        c = card(slide, Inches(0.55), cy, Inches(11.8), Inches(1.1))
        add_rect(slide, Inches(0.55), cy, Inches(0.1), Inches(1.1), fill=NAVY)
        add_textbox(slide, title, Inches(0.85), cy + Inches(0.1), Inches(11.1), Inches(0.35),
                    font_size=14, bold=True, color=NAVY)
        add_textbox(slide, body, Inches(0.85), cy + Inches(0.48), Inches(11.1), Inches(0.52),
                    font_size=13, color=MED, wrap=True)
    return slide


# FIX 2b — Slide 4b: the pattern I kept seeing
def slide04b_pattern(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "My Story"); slide_num(slide, 5)
    add_rect(slide, 0, 0, Inches(0.1), H, fill=GOLD)
    heading(slide, "The Same Pattern — Everywhere",
            Inches(0.55), Inches(0.4), Inches(9), Inches(0.85), size=34)
    gold_line(slide, Inches(0.55), Inches(1.38), Inches(3.2))
    add_textbox(slide, "Every researcher I spoke to had the same story:",
                Inches(0.55), Inches(1.55), Inches(10), Inches(0.42),
                font_size=15, color=MED)
    patterns = [
        ("Hours lost finding papers", "Brilliant researchers spending evenings manually searching databases."),
        ("Scattered notes and citations", "Findings spread across PDFs, Word docs, and browser bookmarks."),
        ("Uncertainty about novelty", "No clear way to know if the research idea was truly original."),
        ("Pressure with limited time", "Publish-or-perish culture, but the tools hadn't kept up."),
        ("Tool overload", "Jumping between ChatGPT, Zotero, Turnitin, and 8 other apps per paper."),
    ]
    for i, (title, body) in enumerate(patterns):
        col = i % 2; row = i // 2
        px = Inches(0.55) + col * Inches(6.4)
        py = Inches(2.1) + row * Inches(1.3)
        if i == 4:
            px = Inches(3.5)
        c = card(slide, px, py, Inches(5.9), Inches(1.1))
        add_rect(slide, px, py, Inches(0.07), Inches(1.1), fill=GOLD)
        add_textbox(slide, title, px + Inches(0.2), py + Inches(0.1), Inches(5.5), Inches(0.35),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, body, px + Inches(0.2), py + Inches(0.48), Inches(5.5), Inches(0.52),
                    font_size=12, color=MED, wrap=True)
    add_textbox(slide, "I built Dynamo AI so researchers spend more time discovering, less time managing.",
                Inches(0.55), H - Inches(1.35), Inches(12.0), Inches(0.5),
                font_size=14, bold=True, italic=True, color=NAVY, wrap=True)
    return slide


# FIX 1 — Slide 5 (now 6): Toolkit redesigned into 3 categories
def slide05_toolkit_v2(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "The Landscape"); slide_num(slide, 6)
    heading(slide, "The Current Research Toolkit",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    add_textbox(slide, "Researchers already use excellent tools — just not together.",
                Inches(0.7), Inches(1.55), Inches(10), Inches(0.42),
                font_size=14, color=MED)
    categories = [
        ("AI Writing & Analysis",
         NAVY,
         ["ChatGPT", "Claude", "Gemini"],
         "Draft, summarise, and analyse research content"),
        ("Citation & Integrity",
         RGBColor(0x1D, 0x40, 0xAF),
         ["Zotero", "Turnitin", "Scite"],
         "Manage references and check for plagiarism"),
        ("Paper Discovery",
         RGBColor(0x07, 0x5E, 0x54),
         ["Perplexity", "Consensus", "ResearchRabbit"],
         "Find papers, track authors, and map literature"),
    ]
    cw = Inches(3.85); gap = Inches(0.22); cx_start = Inches(0.65)
    for i, (cat_title, cat_color, tools, cat_desc) in enumerate(categories):
        cx = cx_start + i * (cw + gap)
        # Category header
        b = add_rect(slide, cx, Inches(2.1), cw, Inches(0.52), fill=cat_color)
        p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = "  " + cat_title
        run.font.name = "Arial"; run.font.size = Pt(13)
        run.font.bold = True; run.font.color.rgb = WHITE
        # Tool pills
        for j, tool in enumerate(tools):
            ty = Inches(2.75) + j * Inches(0.62)
            tc = card(slide, cx, ty, cw, Inches(0.52), fill=WHITE)
            add_textbox(slide, tool, cx + Inches(0.18), ty + Inches(0.08), cw - Inches(0.3), Inches(0.38),
                        font_size=13, bold=True, color=NAVY)
        # Description
        add_textbox(slide, cat_desc, cx, Inches(4.7), cw, Inches(0.55),
                    font_size=11, italic=True, color=MED, wrap=True)
    # Problem callout
    pb = add_rect(slide, Inches(0.65), Inches(5.45), Inches(12.0), Inches(0.68), fill=NAVY)
    add_textbox(slide,
                "The problem: researchers switch between all 3 categories — losing context and time with every switch.",
                Inches(0.85), Inches(5.52), Inches(11.7), Inches(0.55),
                font_size=13, bold=True, color=WHITE, wrap=True)
    return slide


def slide06_problem(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "The Problem"); slide_num(slide, 7)
    heading(slide, "The Problem With Existing Tools",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    sw = Inches(3.6); sx = Inches(0.65)
    for i, (num, label) in enumerate([("1", "Researcher"), ("10", "Tools"), ("20+", "Browser Tabs")]):
        cx = sx + i * (sw + Inches(0.2))
        c = add_rect(slide, cx, Inches(1.7), sw, Inches(2.4), fill=NAVY)
        add_textbox(slide, num, cx, Inches(1.9), sw, Inches(1.5),
                    font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, cx, Inches(3.4), sw, Inches(0.5),
                    font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    pains = [
        ("Broken Workflow",         "Context switching kills deep thinking"),
        ("Multiple Subscriptions",  "Cost adds up — and nothing integrates"),
        ("No Single Source of Truth","Work scattered across tools and tabs"),
    ]
    for i, (title, body) in enumerate(pains):
        cx = sx + i * (sw + Inches(0.2))
        cy = Inches(4.3)
        c = card(slide, cx, cy, sw, Inches(1.4))
        add_textbox(slide, title, cx + Inches(0.15), cy + Inches(0.1), sw - Inches(0.2), Inches(0.4),
                    font_size=13, bold=True, color=NAVY)
        add_textbox(slide, body, cx + Inches(0.15), cy + Inches(0.52), sw - Inches(0.2), Inches(0.7),
                    font_size=11, color=MED, wrap=True)
    add_textbox(slide, "Does this look familiar?", Inches(0.7), H - Inches(1.38), Inches(8), Inches(0.45),
                font_size=16, bold=True, italic=True, color=NAVY)
    return slide


def slide07_born(prs):
    slide = blank_slide(prs); fill_bg(slide, NAVY)
    add_footer_bar(slide); slide_num(slide, 8)
    add_rect(slide, W - Inches(1.5), 0, Inches(1.5), Inches(1.5), fill=GOLD)
    add_textbox(slide,
                '"Researchers don\'t need\nanother AI chatbot.\nThey need one intelligent\nresearch workspace."',
                Inches(1.0), Inches(0.7), Inches(9.5), Inches(3.2),
                font_size=38, bold=True, italic=True, color=WHITE, wrap=True)
    gold_line(slide, Inches(1.0), Inches(4.0), Inches(2.6))
    add_textbox(slide, "Dynamo AI Vision", Inches(1.0), Inches(4.15), Inches(5), Inches(0.45),
                font_size=13, bold=True, color=GOLD)
    for i, v in enumerate(["Search", "Analyze", "Write", "Cite", "Present"]):
        bx = Inches(1.0) + i * Inches(2.2)
        b = add_rect(slide, bx, Inches(4.7), Inches(1.9), Inches(0.6),
                     fill=RGBColor(0x25, 0x40, 0x63))
        b.line.color.rgb = GOLD; b.line.width = Pt(0.6)
        p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = v
        run.font.name = "Arial"; run.font.size = Pt(14)
        run.font.bold = True; run.font.color.rgb = WHITE
    add_textbox(slide, "All in one place.", Inches(1.0), Inches(5.55), Inches(5), Inches(0.45),
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide, "Let me show you how modern researchers work.",
                Inches(1.0), H - Inches(1.38), Inches(10), Inches(0.45),
                font_size=13, italic=True, color=RGBColor(0x94, 0xA3, 0xB8))
    return slide


# FIX 3a — Slide 8: trimmed to 3 bullets per column
def slide08_workflow_v2(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "The Workflow"); slide_num(slide, 9)
    heading(slide, "The AI-Powered Research Workflow",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    col_w = Inches(5.6)
    # Traditional — 3 bullets
    add_rect(slide, Inches(0.6), Inches(1.65), col_w, Inches(0.42), fill=RED)
    add_textbox(slide, "  Traditional Approach",
                Inches(0.6), Inches(1.65), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)
    for i, s in enumerate([
        "Weeks searching databases manually",
        "Rewrite draft 4–5 times from scratch",
        "Submit and hope for the best",
    ]):
        add_textbox(slide, "  " + s, Inches(0.6), Inches(2.2) + i * Inches(0.6),
                    Inches(5.5), Inches(0.5), font_size=13, color=MED)
    # AI-Assisted — 3 bullets
    add_rect(slide, Inches(7.0), Inches(1.65), col_w, Inches(0.42), fill=GREEN)
    add_textbox(slide, "  AI-Assisted (Dynamo AI)",
                Inches(7.0), Inches(1.65), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)
    for i, s in enumerate([
        "Research question to outline in minutes",
        "Draft section-by-section with AI support",
        "Verify citations before submission",
    ]):
        add_textbox(slide, "  " + s, Inches(7.0), Inches(2.2) + i * Inches(0.6),
                    Inches(5.5), Inches(0.5), font_size=13, color=DARK)
    # Time saved banner
    tb = add_rect(slide, Inches(0.6), Inches(4.25), Inches(12.0), Inches(0.72), fill=NAVY)
    add_textbox(slide, "Time saved: weeks of manual effort reduced to focused, structured sessions.",
                Inches(0.8), Inches(4.32), Inches(11.6), Inches(0.56),
                font_size=14, bold=True, color=GOLD, wrap=True)
    # Example prompt
    add_rect(slide, Inches(0.6), Inches(5.15), Inches(12.0), Inches(0.78), fill=LIGHT)
    add_textbox(slide,
                'Example: "Draft the Introduction for a paper on AI adoption in higher education using a formal academic tone."',
                Inches(0.75), Inches(5.22), Inches(11.7), Inches(0.65),
                font_size=12, italic=True, color=NAVY, wrap=True)
    return slide


# FIX 3b — Slide 9: trimmed to 3 bullets per column
def slide09_write_faster_v2(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Writing"); slide_num(slide, 10)
    heading(slide, "Write Your Paper Faster Without Cutting Corners",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=32)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    col_w = Inches(5.6)
    add_rect(slide, Inches(0.6), Inches(1.65), col_w, Inches(0.42), fill=RED)
    add_textbox(slide, "  Old Way", Inches(0.6), Inches(1.65), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)
    for i, s in enumerate([
        "Stare at blank page for hours",
        "Rewrite entire sections when structure is wrong",
        "Citations added last — often wrong or missing",
    ]):
        add_textbox(slide, "  " + s, Inches(0.6), Inches(2.2) + i * Inches(0.62),
                    Inches(5.5), Inches(0.5), font_size=13, color=MED, italic=(i > 0))
    add_rect(slide, Inches(7.0), Inches(1.65), col_w, Inches(0.42), fill=GREEN)
    add_textbox(slide, "  With Dynamo AI", Inches(7.0), Inches(1.65), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)
    for i, s in enumerate([
        "1. Define objective",
        "2. Generate structure",
        "3. Draft section-by-section",
    ]):
        add_textbox(slide, "  " + s, Inches(7.0), Inches(2.2) + i * Inches(0.62),
                    Inches(5.5), Inches(0.5), font_size=13, color=DARK)
    add_rect(slide, Inches(0.6), Inches(4.3), Inches(12.0), Inches(0.78), fill=LIGHT)
    add_textbox(slide,
                'Example Prompt: "Draft the Introduction for a paper on AI adoption in higher education using a formal academic tone."',
                Inches(0.75), Inches(4.37), Inches(11.7), Inches(0.65),
                font_size=12, italic=True, color=NAVY, wrap=True)
    return slide


def slide10_framework(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Framework"); slide_num(slide, 11)
    heading(slide, "The 5-Step Research Writing Framework",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.2))
    steps = [
        ("Define Objectives",   "Clarify your research question and scope before writing a single word."),
        ("Generate Structure",  "AI builds a logical outline — introduction, methods, discussion, conclusion."),
        ("Draft Strategically", "Write section-by-section with AI assistance, maintaining your academic voice."),
        ("Verify Evidence",     "Cross-check claims against real papers and sources. Validate every fact."),
        ("Refine and Humanize", "Add your expertise, insights, and perspective. Make it genuinely yours."),
    ]
    bw = Inches(2.28); bh = Inches(2.3); gap = Inches(0.14)
    start_x = Inches(0.5); by = Inches(1.75)
    for i, (title, desc) in enumerate(steps):
        bx = start_x + i * (bw + gap)
        c = add_rect(slide, bx, by, bw, bh, fill=WHITE)
        c.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); c.line.width = Pt(0.5)
        add_rect(slide, bx, by, bw, Inches(0.055), fill=GOLD)
        num_circle(slide, str(i + 1), bx + bw / 2, by + Inches(0.46), radius=Inches(0.33))
        add_textbox(slide, title, bx + Inches(0.1), by + Inches(0.95), bw - Inches(0.2), Inches(0.52),
                    font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, wrap=True)
        add_textbox(slide, desc, bx + Inches(0.1), by + Inches(1.5), bw - Inches(0.2), Inches(0.7),
                    font_size=10.5, color=MED, align=PP_ALIGN.CENTER, wrap=True)
    add_textbox(slide,
                "This framework applies to journal articles, conference papers, theses, and research reports.",
                Inches(0.7), H - Inches(1.4), Inches(11.5), Inches(0.45),
                font_size=13, italic=True, color=MED)
    return slide


# POLL SLIDE BUILDER
def poll_slide(prs, num, question, options, tip="Use Mentimeter or Slido — show live results on screen"):
    slide = blank_slide(prs); fill_bg(slide, NAVY)
    add_footer_bar(slide, "Dynamo AI   |   Webinar 2026   |   Audience Poll")
    slide_num(slide, num)
    # Gold corner
    add_rect(slide, 0, 0, Inches(0.12), H, fill=GOLD)
    # Label
    lb = add_rect(slide, Inches(0.4), Inches(0.3), Inches(2.6), Inches(0.38), fill=GOLD)
    p = lb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "AUDIENCE POLL"
    run.font.name = "Arial"; run.font.size = Pt(10)
    run.font.bold = True; run.font.color.rgb = NAVY
    # Question
    add_textbox(slide, question, Inches(0.55), Inches(0.88), Inches(12.3), Inches(1.2),
                font_size=26, bold=True, color=WHITE, wrap=True)
    gold_line(slide, Inches(0.55), Inches(2.2), Inches(4.0))
    # Options — 2 columns of up to 3
    left_opts  = options[:3]
    right_opts = options[3:]
    for i, opt in enumerate(left_opts):
        oy = Inches(2.45) + i * Inches(0.88)
        ob = add_rect(slide, Inches(0.55), oy, Inches(5.9), Inches(0.72),
                      fill=RGBColor(0x25, 0x40, 0x63))
        ob.line.color.rgb = GOLD; ob.line.width = Pt(0.6)
        add_textbox(slide, opt, Inches(0.75), oy + Inches(0.13), Inches(5.6), Inches(0.48),
                    font_size=13, color=WHITE, wrap=True)
    for i, opt in enumerate(right_opts):
        oy = Inches(2.45) + i * Inches(0.88)
        ob = add_rect(slide, Inches(6.9), oy, Inches(5.9), Inches(0.72),
                      fill=RGBColor(0x25, 0x40, 0x63))
        ob.line.color.rgb = GOLD; ob.line.width = Pt(0.6)
        add_textbox(slide, opt, Inches(7.1), oy + Inches(0.13), Inches(5.6), Inches(0.48),
                    font_size=13, color=WHITE, wrap=True)
    # Tip
    add_textbox(slide, tip, Inches(0.55), H - Inches(1.35), Inches(12.0), Inches(0.42),
                font_size=10.5, italic=True, color=RGBColor(0x94, 0xA3, 0xB8), wrap=True)
    return slide


def slide11_litreview(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Literature Review"); slide_num(slide, 13)
    heading(slide, "Literature Reviews Without the Chaos",
            Inches(0.7), Inches(0.5), Inches(10), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    col_w = Inches(5.6)
    add_rect(slide, Inches(0.6), Inches(1.7), col_w, Inches(0.42), fill=RED)
    add_textbox(slide, "  Traditional Way", Inches(0.6), Inches(1.7), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)
    for i, t in enumerate([
        "Download 50+ PDFs manually",
        "Highlight and annotate by hand",
        "Hours lost — weeks of effort",
    ]):
        add_textbox(slide, "    " + t, Inches(0.6), Inches(2.28) + i * Inches(0.54),
                    Inches(5.5), Inches(0.48), font_size=13, color=MED)
    add_rect(slide, Inches(7.0), Inches(1.7), col_w, Inches(0.42), fill=GREEN)
    add_textbox(slide, "  The Smarter Way (Dynamo AI)", Inches(7.0), Inches(1.7), col_w, Inches(0.42),
                font_size=13, bold=True, color=WHITE)
    for i, s in enumerate([
        "Upload papers — Dynamo reads them all",
        "Extract key insights automatically",
        "Build synthesis in minutes",
    ]):
        add_textbox(slide, "  " + s, Inches(7.0), Inches(2.28) + i * Inches(0.54),
                    Inches(5.5), Inches(0.48), font_size=13, color=DARK)
    add_rect(slide, Inches(0.6), Inches(4.1), Inches(12.0), Inches(0.7), fill=LIGHT)
    add_textbox(slide,
                'Example: "Summarize the impact of AI on Student Learning Outcomes" — themes identified in seconds.',
                Inches(0.75), Inches(4.17), Inches(11.7), Inches(0.58),
                font_size=12, italic=True, color=NAVY, wrap=True)
    return slide


def slide12_litreview_example(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Literature Review"); slide_num(slide, 14)
    heading(slide, "See the Bigger Picture Faster",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(2.8))
    add_textbox(slide, "Topic: Impact of AI on Higher Education",
                Inches(0.7), Inches(1.55), Inches(8), Inches(0.45),
                font_size=14, color=MED, italic=True)
    headers = ["Theme", "Key Findings", "Implication"]
    rows_data = [
        ("Student Engagement",   "Increased participation in AI-assisted courses",              "Interactive tools improve attendance"),
        ("Learning Outcomes",    "Mixed evidence — strong in STEM, weaker in humanities",       "Context matters significantly"),
        ("Personalization",      "Positive impact on self-paced learning",                      "AI tutors adapt to student pace"),
        ("Ethics & Privacy",     "Privacy concerns around data collection",                     "Policy gaps need urgent attention"),
    ]
    tx, ty = Inches(0.6), Inches(2.15)
    tw_cols = [Inches(2.6), Inches(4.5), Inches(5.0)]
    for ci, (hdr, cw) in enumerate(zip(headers, tw_cols)):
        cx = tx + sum(tw_cols[:ci])
        b = add_rect(slide, cx, ty, cw, Inches(0.5), fill=NAVY)
        add_textbox(slide, hdr, cx + Inches(0.1), ty + Inches(0.07), cw - Inches(0.15), Inches(0.38),
                    font_size=13, bold=True, color=WHITE)
    for ri, row in enumerate(rows_data):
        ry = ty + Inches(0.5) + ri * Inches(0.72)
        fill = WHITE if ri % 2 == 0 else LIGHT
        for ci, (cell, cw) in enumerate(zip(row, tw_cols)):
            cx = tx + sum(tw_cols[:ci])
            b = add_rect(slide, cx, ry, cw, Inches(0.65), fill=fill)
            b.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); b.line.width = Pt(0.4)
            if ci == 0:
                add_rect(slide, cx, ry, Inches(0.07), Inches(0.65), fill=GOLD)
            add_textbox(slide, cell, cx + Inches(0.15), ry + Inches(0.1), cw - Inches(0.2), Inches(0.5),
                        font_size=11, color=DARK, wrap=True)
    add_textbox(slide,
                "Instead of reading dozens of papers separately, researchers immediately identify patterns and contradictions.",
                Inches(0.7), H - Inches(1.4), Inches(11.5), Inches(0.5),
                font_size=13, italic=True, color=NAVY, wrap=True)
    return slide


def slide13_gaps(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Gap Discovery"); slide_num(slide, 15)
    heading(slide, "Finding Research Gaps Faster",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(2.8))
    add_textbox(slide, "Gap Discovery Framework", Inches(0.6), Inches(1.6), Inches(4.5), Inches(0.42),
                font_size=13, bold=True, color=NAVY)
    fsteps = ["Known Research", "Emerging Trends", "Contradictions",
              "Understudied Groups", "Methodological Gaps", "Research Gap Found"]
    for i, (s, sc) in enumerate(zip(fsteps, [NAVY]*5 + [GOLD])):
        by = Inches(2.1) + i * Inches(0.66)
        b = add_rect(slide, Inches(0.6), by, Inches(3.8), Inches(0.54), fill=sc)
        add_textbox(slide, s, Inches(0.75), by + Inches(0.08), Inches(3.5), Inches(0.4),
                    font_size=12, bold=(i == 5), color=WHITE if sc != GOLD else NAVY)
        if i < 5:
            add_textbox(slide, "v", Inches(1.6), by + Inches(0.54), Inches(0.5), Inches(0.2),
                        font_size=10, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Example: AI in Higher Education",
                Inches(5.2), Inches(1.6), Inches(7.5), Inches(0.42),
                font_size=13, bold=True, color=NAVY)
    eb = card(slide, Inches(5.2), Inches(2.1), Inches(3.5), Inches(3.2))
    add_textbox(slide, "Existing Research", Inches(5.35), Inches(2.15), Inches(3.2), Inches(0.38),
                font_size=12, bold=True, color=NAVY)
    for i, e in enumerate(["Undergraduate students", "Student engagement", "Learning outcomes", "Developed countries"]):
        add_textbox(slide, "+ " + e, Inches(5.35), Inches(2.6) + i * Inches(0.52), Inches(3.2), Inches(0.45),
                    font_size=11, color=GREEN)
    mb = card(slide, Inches(9.0), Inches(2.1), Inches(3.7), Inches(3.2))
    add_textbox(slide, "Missing Areas", Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.38),
                font_size=12, bold=True, color=RED)
    for i, m in enumerate(["Rural institutions", "Faculty adoption", "Long-term impacts", "Developing countries"]):
        add_textbox(slide, "? " + m, Inches(9.15), Inches(2.6) + i * Inches(0.52), Inches(3.4), Inches(0.45),
                    font_size=11, color=RED)
    gb = add_rect(slide, Inches(5.2), Inches(5.45), Inches(7.5), Inches(0.75), fill=NAVY)
    add_textbox(slide, '"How do faculty in developing countries adopt AI-assisted teaching?"',
                Inches(5.35), Inches(5.5), Inches(7.2), Inches(0.65),
                font_size=12, bold=True, italic=True, color=GOLD, wrap=True)
    return slide


def slide14_gap_formula(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Gap Discovery"); slide_num(slide, 16)
    heading(slide, "A Simple Formula for Discovering Novel Ideas",
            Inches(0.7), Inches(0.45), Inches(11), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.2))
    elements = [
        ("Known Findings",           "What existing research has established"),
        ("Conflicting Results",      "Where studies contradict each other"),
        ("Understudied Groups",      "Populations not yet represented"),
        ("Emerging Trends",          "New phenomena needing investigation"),
    ]
    ew = Inches(2.75); ex = Inches(0.55)
    for i, (title, desc) in enumerate(elements):
        bx = ex + i * (ew + Inches(0.18))
        c = add_rect(slide, bx, Inches(1.7), ew, Inches(2.0), fill=NAVY)
        add_textbox(slide, title, bx + Inches(0.12), Inches(1.82), ew - Inches(0.2), Inches(0.52),
                    font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER, wrap=True)
        add_textbox(slide, desc, bx + Inches(0.1), Inches(2.45), ew - Inches(0.18), Inches(0.75),
                    font_size=11, color=SLATE, align=PP_ALIGN.CENTER, wrap=True)
        if i < 3:
            add_textbox(slide, "+", bx + ew + Inches(0.02), Inches(2.55), Inches(0.18), Inches(0.35),
                        font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "=", Inches(5.8), Inches(3.85), Inches(0.6), Inches(0.55),
                font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.0), Inches(4.55), Inches(6.85), Inches(0.88), fill=GOLD)
    add_textbox(slide, "RESEARCH OPPORTUNITY", Inches(3.0), Inches(4.62), Inches(6.85), Inches(0.7),
                font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Dynamo AI's Research Gap feature applies this formula automatically — surfacing opportunities you might have missed.",
                Inches(0.7), H - Inches(1.38), Inches(11.5), Inches(0.48),
                font_size=13, italic=True, color=MED, wrap=True)
    return slide


def slide15_ai_workflow(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Product Demo"); slide_num(slide, 17)
    heading(slide, "AI Research Workflow in Practice",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(3.0))
    steps = [
        ("Upload Papers",     "PDF, DOCX, TXT — Dynamo reads everything"),
        ("Extract Themes",    "Automatic topic clustering and synthesis"),
        ("Draft Content",     "AI generates structured academic text"),
        ("Validate Sources",  "Cross-check against real publications"),
    ]
    sw = Inches(2.85); sh = Inches(2.6); gap = Inches(0.22); sx = Inches(0.6)
    for i, (title, desc) in enumerate(steps):
        bx = sx + i * (sw + gap)
        c = card(slide, bx, Inches(1.75), sw, sh, fill=WHITE)
        nb = add_rect(slide, bx + Inches(0.15), Inches(1.9), Inches(0.45), Inches(0.45), fill=NAVY)
        add_textbox(slide, str(i + 1), bx + Inches(0.15), Inches(1.9), Inches(0.45), Inches(0.45),
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, bx + Inches(0.15), Inches(2.48), sw - Inches(0.25), Inches(0.52),
                    font_size=14, bold=True, color=NAVY, wrap=True)
        add_textbox(slide, desc, bx + Inches(0.15), Inches(3.05), sw - Inches(0.25), Inches(0.8),
                    font_size=12, color=MED, wrap=True)
        if i < 3:
            add_textbox(slide, ">", bx + sw + Inches(0.04), Inches(2.93), gap - Inches(0.04), Inches(0.35),
                        font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rb = add_rect(slide, Inches(0.6), Inches(4.6), Inches(12.0), Inches(0.72), fill=NAVY)
    add_textbox(slide, "Result: A fully structured, cited, verified research draft — in a fraction of the time",
                Inches(0.75), Inches(4.67), Inches(11.7), Inches(0.6),
                font_size=14, bold=True, color=GOLD, wrap=True, align=PP_ALIGN.CENTER)
    return slide


def slide16_assistant(prs):
    return navy_panel_slide(prs, 18, "Features", "AI Research\nAssistant",
        "Chat directly with your research papers. Compare authors. Explore contradictions.",
        "What you can do:",
        [
            ("Chat with PDFs",           "Upload any paper and ask questions in plain language"),
            ("Compare Authors",          "See where researchers agree, disagree, and contradict"),
            ("Extract Key Arguments",    "Pull the core claims from any document instantly"),
            ("Build Literature Matrix",  "Organize findings across multiple papers visually"),
        ])


# FIX 4a — Slide 17: AI Detector redesigned (stat-centred layout)
def slide17_detector_v2(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Features"); slide_num(slide, 19)
    heading(slide, "AI Content Detector",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(2.8))
    add_textbox(slide, "Know exactly how AI-generated your writing looks before you submit.",
                Inches(0.7), Inches(1.58), Inches(9), Inches(0.45), font_size=14, color=MED)
    # Large score visual
    score_box = add_rect(slide, Inches(0.65), Inches(2.15), Inches(4.2), Inches(3.2), fill=NAVY)
    add_textbox(slide, "AI Score", Inches(0.65), Inches(2.35), Inches(4.2), Inches(0.45),
                font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "27", Inches(0.65), Inches(2.75), Inches(4.2), Inches(1.6),
                font_size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "LOW AI PROBABILITY — SAFE TO SUBMIT",
                Inches(0.65), Inches(4.35), Inches(4.2), Inches(0.52),
                font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # Feature cards right
    features = [
        ("Score Meter 0–100",        "Instant visual probability of AI-generation in your text."),
        ("Writing Pattern Analysis", "Flags unnaturally uniform sentences — a common AI signal."),
        ("Plagiarism Check",         "Cross-checks web sources and academic paper databases."),
        ("Actionable Highlights",    "Exact sentences flagged so you know what to rewrite."),
    ]
    rx = Inches(5.3); rw = Inches(7.6)
    add_textbox(slide, "Built for academic integrity:", rx, Inches(2.05), rw, Inches(0.4),
                font_size=13, bold=True, color=NAVY)
    for i, (t, b) in enumerate(features):
        fy = Inches(2.55) + i * Inches(0.75)
        c = card(slide, rx, fy, rw, Inches(0.65), fill=WHITE)
        left_strip(slide, rx, fy, Inches(0.65))
        add_textbox(slide, t, rx + Inches(0.18), fy + Inches(0.06), rw - Inches(0.28), Inches(0.28),
                    font_size=12, bold=True, color=NAVY)
        add_textbox(slide, b, rx + Inches(0.18), fy + Inches(0.34), rw - Inches(0.28), Inches(0.26),
                    font_size=10.5, color=MED, wrap=True)
    return slide


def slide18_citation(prs):
    return navy_panel_slide(prs, 20, "Features", "Citation\nChecker",
        "Fix missing, broken, and unverified references before they cost you a rejection.",
        "What it checks:",
        [
            ("Missing Citations",    "Identifies claims made without supporting references"),
            ("Broken References",    "Flags citations that can't be verified in academic databases"),
            ("Format Validation",    "Checks APA, MLA, Chicago, Harvard compliance automatically"),
            ("DOI Verification",     "Confirms each paper actually exists and is retrievable"),
        ])


# FIX 4b — Slide 19: Research Watcher redesigned (alert-flow layout)
def slide19_watcher_v2(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Features"); slide_num(slide, 21)
    heading(slide, "Research Watcher",
            Inches(0.7), Inches(0.5), Inches(9), Inches(0.8), size=34)
    gold_line(slide, Inches(0.7), Inches(1.4), Inches(2.6))
    add_textbox(slide, "Never miss a new paper in your field — without manual searching.",
                Inches(0.7), Inches(1.58), Inches(9), Inches(0.45), font_size=14, color=MED)
    # Horizontal flow: Set Topic → Monitor → Alert → Read
    flow_steps = [
        ("1. Set Topics",    "Add your research keywords or author names"),
        ("2. Monitor",       "Dynamo watches 100M+ papers automatically"),
        ("3. Get Alerted",   "Email or in-app notification when new work appears"),
        ("4. Stay Current",  "Weekly digest: curated summary of what is new"),
    ]
    fw = Inches(2.85); fh = Inches(1.9); gap = Inches(0.2); fx_start = Inches(0.6)
    for i, (title, desc) in enumerate(flow_steps):
        fx = fx_start + i * (fw + gap)
        c = add_rect(slide, fx, Inches(2.15), fw, fh, fill=NAVY)
        add_rect(slide, fx, Inches(2.15), fw, Inches(0.06), fill=GOLD)
        add_textbox(slide, title, fx + Inches(0.15), Inches(2.32), fw - Inches(0.25), Inches(0.48),
                    font_size=14, bold=True, color=GOLD, wrap=True)
        add_textbox(slide, desc, fx + Inches(0.15), Inches(2.9), fw - Inches(0.25), Inches(0.95),
                    font_size=12, color=SLATE, wrap=True)
        if i < 3:
            add_textbox(slide, ">", fx + fw + Inches(0.03), Inches(2.9), Inches(0.18), Inches(0.35),
                        font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Example alert card
    alert_y = Inches(4.3)
    ab = card(slide, Inches(0.65), alert_y, Inches(12.0), Inches(0.88), fill=LIGHT)
    add_rect(slide, Inches(0.65), alert_y, Inches(0.07), Inches(0.88), fill=GOLD)
    add_textbox(slide, "Example Alert:",
                Inches(0.9), alert_y + Inches(0.1), Inches(2.0), Inches(0.32),
                font_size=11, bold=True, color=NAVY)
    add_textbox(slide,
                '"3 new papers published this week matching: AI-assisted teaching in developing countries"',
                Inches(3.0), alert_y + Inches(0.1), Inches(9.4), Inches(0.65),
                font_size=12, italic=True, color=NAVY, wrap=True)
    return slide


# FIX 5 — Slide 20: Demo enhanced with clear 3-step script
def slide20_demo_v2(prs):
    slide = blank_slide(prs); fill_bg(slide, NAVY)
    add_footer_bar(slide, "Dynamo AI   |   Webinar 2026   |   Live Demo")
    slide_num(slide, 22)
    add_rect(slide, W - Inches(1.5), 0, Inches(1.5), Inches(1.5), fill=GOLD)
    add_textbox(slide, "LIVE DEMO", Inches(1.0), Inches(0.65), Inches(11), Inches(0.45),
                font_size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    heading(slide, "See Dynamo AI in Action",
            Inches(1.0), Inches(1.1), Inches(11), Inches(1.1),
            size=44, color=WHITE, align=PP_ALIGN.CENTER)
    gold_line(slide, Inches(4.0), Inches(2.35), Inches(5.2))
    add_textbox(slide, "3-step live walkthrough:", Inches(1.0), Inches(2.55), Inches(11), Inches(0.42),
                font_size=13, color=GOLD, align=PP_ALIGN.CENTER)
    demo_steps = [
        ("Step 1", "Upload a paper + ask a question", "Show: Chat with PDF"),
        ("Step 2", "Run the Research Gap finder",     "Show: Gap Discovery"),
        ("Step 3", "Draft an Introduction section",   "Show: AI Writing"),
    ]
    for i, (num, action, show) in enumerate(demo_steps):
        dx = Inches(1.2) + i * Inches(3.75)
        db = add_rect(slide, dx, Inches(3.1), Inches(3.4), Inches(1.6),
                      fill=RGBColor(0x25, 0x40, 0x63))
        db.line.color.rgb = GOLD; db.line.width = Pt(0.8)
        add_textbox(slide, num, dx + Inches(0.15), Inches(3.18), Inches(3.1), Inches(0.35),
                    font_size=11, bold=True, color=GOLD)
        add_textbox(slide, action, dx + Inches(0.15), Inches(3.52), Inches(3.1), Inches(0.5),
                    font_size=13, bold=True, color=WHITE, wrap=True)
        add_textbox(slide, show, dx + Inches(0.15), Inches(4.05), Inches(3.1), Inches(0.45),
                    font_size=11, italic=True, color=RGBColor(0x94, 0xA3, 0xB8), wrap=True)
    add_textbox(slide, "app.dynamoai.in",
                Inches(1.0), H - Inches(1.55), Inches(11), Inches(0.45),
                font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    return slide


def slide21_validation(prs):
    slide = blank_slide(prs); fill_bg(slide, BG)
    add_footer_bar(slide); section_label(slide, "Validation"); slide_num(slide, 23)
    add_rect(slide, 0, 0, Inches(0.1), H, fill=GOLD)
    heading(slide, "We Are Building Dynamo AI\nWith Researchers Like You",
            Inches(0.55), Inches(0.45), Inches(9), Inches(1.6), size=36, wrap=True)
    gold_line(slide, Inches(0.55), Inches(2.15), Inches(3.2))
    add_textbox(slide, "Your feedback shapes every feature we build.",
                Inches(0.55), Inches(2.3), Inches(9), Inches(0.5), font_size=16, color=MED)
    points = [
        ("Your Voice Matters",       "Every piece of feedback directly influences our product roadmap."),
        ("Real Research. Real Problems.", "We focus on what actually slows researchers down."),
        ("Built for Indian Academia","Designed with local researcher workflows and goals."),
        ("Continuous Improvement",   "Weekly updates based on what our validation community tells us."),
    ]
    pw = Inches(5.6)
    for i, (title, body) in enumerate(points):
        px = Inches(0.55) if i % 2 == 0 else Inches(6.8)
        py = Inches(3.0) + (i // 2) * Inches(1.4)
        c = card(slide, px, py, pw, Inches(1.2), fill=WHITE)
        add_rect(slide, px, py, Inches(0.08), Inches(1.2), fill=NAVY)
        add_textbox(slide, title, px + Inches(0.22), py + Inches(0.1), pw - Inches(0.3), Inches(0.4),
                    font_size=14, bold=True, color=NAVY)
        add_textbox(slide, body, px + Inches(0.22), py + Inches(0.55), pw - Inches(0.3), Inches(0.55),
                    font_size=12, color=MED, wrap=True)
    return slide


def slide22_bonus(prs):
    slide = blank_slide(prs); fill_bg(slide, GOLD)
    add_footer_bar(slide); slide_num(slide, 25)
    add_rect(slide, 0, 0, Inches(0.12), H, fill=NAVY)
    add_textbox(slide, "EXCLUSIVE FOR WEBINAR ATTENDEES", Inches(0.5), Inches(0.65), Inches(12), Inches(0.42),
                font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    heading(slide, "14 Days of Dynamo AI Pro",
            Inches(0.5), Inches(1.1), Inches(12), Inches(1.3),
            size=52, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Completely Free — No Strings Attached", Inches(0.5), Inches(2.45), Inches(12), Inches(0.5),
                font_size=20, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.0), Inches(3.05), Inches(5.3), Inches(0.05), fill=NAVY)
    included = [
        "300 AI Research Messages / Day",
        "Deep Research Agent",
        "Full Literature Review & Gap Discovery",
        "AI Detector + Plagiarism Checker",
        "Citation Checker + Research Watcher",
        "Priority Support",
    ]
    inc_w = Inches(3.8)
    for i, item in enumerate(included):
        col = i % 2; row = i // 2
        ix = Inches(2.1) + col * (inc_w + Inches(0.5))
        iy = Inches(3.3) + row * Inches(0.58)
        add_rect(slide, ix, iy + Inches(0.15), Inches(0.18), Inches(0.18), fill=NAVY)
        add_textbox(slide, "      " + item, ix, iy, inc_w, Inches(0.5),
                    font_size=14, color=NAVY)
    add_textbox(slide, "How to claim: Fill out the feedback form at the end of this session",
                Inches(0.5), H - Inches(1.42), Inches(12), Inches(0.45),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return slide


def slide23_thankyou(prs):
    slide = blank_slide(prs); fill_bg(slide, WHITE)
    panel_w = Inches(5.6)
    add_rect(slide, 0, 0, panel_w, H, fill=NAVY)
    add_rect(slide, 0, 0, Inches(0.08), H, fill=GOLD)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.5), height=Inches(0.85))
    add_textbox(slide, "THANK YOU", Inches(0.5), Inches(1.7), panel_w - Inches(0.7), Inches(0.6),
                font_size=11, bold=True, color=GOLD)
    heading(slide, "From Research\nOverwhelm to\nResearch Confidence",
            Inches(0.5), Inches(2.35), panel_w - Inches(0.65), Inches(2.6),
            size=32, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(5.05), Inches(2.6), Inches(0.05), fill=GOLD)
    add_textbox(slide, "Anish Krisna\nFounder, Dynamo AI",
                Inches(0.5), Inches(5.2), panel_w - Inches(0.7), Inches(0.7),
                font_size=13, color=SLATE, wrap=True)
    rx = panel_w + Inches(0.55); rw = W - panel_w - Inches(0.75)
    heading(slide, "What happens next?", rx, Inches(0.8), rw, Inches(0.5), size=20, color=NAVY)
    gold_line(slide, rx, Inches(1.38), Inches(2.8))
    for i, (num, txt) in enumerate([
        ("1.", "Fill out the feedback form in the chat"),
        ("2.", "Receive your Pro invite code privately"),
        ("3.", "Visit app.dynamoai.in and sign up"),
        ("4.", "Enter your code — start your 14-day Pro trial"),
        ("5.", "Share your experience — help us build better"),
    ]):
        sy = Inches(1.55) + i * Inches(0.75)
        add_textbox(slide, num, rx, sy, Inches(0.4), Inches(0.55),
                    font_size=18, bold=True, color=GOLD)
        add_textbox(slide, txt, rx + Inches(0.4), sy + Inches(0.05), rw - Inches(0.45), Inches(0.5),
                    font_size=14, color=DARK, wrap=True)
    add_rect(slide, rx, H - Inches(1.55), rw, Inches(0.62), fill=LIGHT)
    add_textbox(slide, "app.dynamoai.in", rx + Inches(0.15), H - Inches(1.5), rw - Inches(0.2), Inches(0.52),
                font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return slide


# ══ BUILD ═══════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide01_welcome(prs)                 # 1
    slide02_why_here(prs)                # 2
    slide03_reality(prs)                 # 3
    slide04a_story(prs)                  # 4  ← SPLIT: personal journey
    slide04b_pattern(prs)                # 5  ← SPLIT: the pattern
    slide05_toolkit_v2(prs)              # 6  ← REDESIGNED: 3 categories
    slide06_problem(prs)                 # 7
    slide07_born(prs)                    # 8
    slide08_workflow_v2(prs)             # 9  ← TRIMMED: 3 bullets
    slide09_write_faster_v2(prs)         # 10 ← TRIMMED: 3 bullets
    slide10_framework(prs)               # 11

    poll_slide(prs, 12,                  # 12 ← POLL 1
        "Which step of the 5-Step Framework do you struggle with most?",
        ["A.  Defining a clear research objective",
         "B.  Structuring the paper logically",
         "C.  Actually drafting — getting words on the page",
         "D.  Verifying sources and citations",
         "E.  Making it sound like me, not AI"])

    slide11_litreview(prs)               # 13
    slide12_litreview_example(prs)       # 14
    slide13_gaps(prs)                    # 15
    slide14_gap_formula(prs)             # 16
    slide15_ai_workflow(prs)             # 17
    slide16_assistant(prs)               # 18
    slide17_detector_v2(prs)             # 19 ← REDESIGNED: stat layout
    slide18_citation(prs)                # 20

    poll_slide(prs, 21,                  # 21 ← POLL 2
        "Have you ever worried about submitting writing that contained AI-generated text?",
        ["A.  Yes — and it caused a problem",
         "B.  Yes — I caught it before submitting",
         "C.  Not yet, but I worry about it",
         "D.  I don't use AI in my writing"])

    slide19_watcher_v2(prs)              # 22 ← REDESIGNED: alert-flow
    slide20_demo_v2(prs)                 # 23 ← ENHANCED: 3-step script

    poll_slide(prs, 24,                  # 24 ← POLL 3
        "What would make you most confident to try Dynamo AI for your research?",
        ["A.  A live demo on a real topic (just happened!)",
         "B.  Reviews from researchers like me",
         "C.  Free trial — no credit card required",
         "D.  Recommendation from my supervisor"])

    slide21_validation(prs)              # 25
    slide22_bonus(prs)                   # 26
    slide23_thankyou(prs)                # 27

    out = "Dynamo_AI_Webinar_2026_v2.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
