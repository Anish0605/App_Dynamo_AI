"""
Dynamo AI Pitch Deck — PPTX Generator (v2, May 2026)
17 slides, 16:9 widescreen, brand colors: #0a0e1a / #FFC107 / #ffffff

Changes from v1 (based on incubator/ChatGPT feedback):
- Updated cover positioning: "The Research Operating System for Academia"
- Slide 3 (Cost): removed unverifiable hour figures → qualitative pains
- Slides 6 + 7: real product screenshots (was: text-only cards)
- NEW Slide 8: Demo Video (YouTube placeholder)
- NEW Slide 9: Why Dynamo Wins
- Slide 10 (Competition): softened "only" claim
- Slide 11 (Why Now): qualitative claims, no unsourced stats
- Slide 12 (Market): TAM / SAM / SOM framing
- Slide 14 (Validation): unchanged per founder's request
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Brand tokens ──────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x0a, 0x0e, 0x1a)
BLACK2 = RGBColor(0x14, 0x19, 0x2b)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xAA, 0xAA, 0xBB)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
RED    = RGBColor(0xEF, 0x44, 0x44)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)
CARD   = RGBColor(0x12, 0x17, 0x28)
BORDER = RGBColor(0x25, 0x2A, 0x40)

W  = Inches(13.33)
H  = Inches(7.5)
TOTAL_SLIDES = 17

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

LOGO_PATH = "attached_assets/Dynamo_AI_New_Logo_1779360398074.png"

# Product screenshots
SHOT_DASHBOARD_PRO  = "attached_assets/Screenshot_2026-05-21_at_9.07.27_PM_1779378444860.png"
SHOT_DASHBOARD_OUT  = "attached_assets/Screenshot_2026-05-21_at_9.06.48_PM_1779378444856.png"
SHOT_SIDEBAR_TOOLS  = "attached_assets/Screenshot_2026-05-21_at_9.07.44_PM_1779378444862.png"
SHOT_SIDEBAR_CHATS  = "attached_assets/Screenshot_2026-05-21_at_9.07.59_PM_1779378444862.png"
SHOT_SIDEBAR_FOLDER = "attached_assets/Screenshot_2026-05-21_at_9.08.07_PM_1779378444863.png"
SHOT_PROFILE        = "attached_assets/Screenshot_2026-05-21_at_9.08.16_PM_1779378444863.png"
SHOT_AI_MEMORY      = "attached_assets/Screenshot_2026-05-21_at_9.08.27_PM_1779378444864.png"
SHOT_DOC_LIBRARY    = "attached_assets/Screenshot_2026-05-21_at_9.08.38_PM_1779378444864.png"
SHOT_PLUS_MENU      = "attached_assets/Screenshot_2026-05-21_at_9.08.49_PM_1779378444865.png"
SHOT_TOOLS_MENU     = "attached_assets/Screenshot_2026-05-21_at_9.08.58_PM_1779378444865.png"

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(BLANK)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return sl

def txb(sl, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return tx

def rect(sl, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = sl.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = line_width
    else:
        shape.line.fill.background()
    return shape

def header(sl, page_num, section_label=""):
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(0.4), Inches(0.2),
                              width=Inches(0.5), height=Inches(0.5))
    txb(sl, "Dynamo AI", Inches(1.0), Inches(0.25), Inches(2), Inches(0.4),
        size=16, bold=True, color=WHITE)
    if section_label:
        txb(sl, section_label.upper(), Inches(4), Inches(0.28), Inches(5.33), Inches(0.35),
            size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txb(sl, f"{page_num:02d} / {TOTAL_SLIDES}", Inches(12.3), Inches(0.28), Inches(0.9), Inches(0.35),
        size=9, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

def footer(sl):
    txb(sl, "DYNAMO AI · CONFIDENTIAL",
        Inches(0.4), Inches(7.1), Inches(4), Inches(0.28),
        size=8, color=MUTED)
    txb(sl, "app.dynamoai.in",
        Inches(9.5), Inches(7.1), Inches(3.4), Inches(0.28),
        size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def section_tag(sl, label):
    rect(sl, Inches(0.4), Inches(0.95), Inches(len(label)*0.11 + 0.4), Inches(0.3),
         fill_color=RGBColor(0x26, 0x1E, 0x00), line_color=RGBColor(0x80, 0x60, 0x03))
    txb(sl, f"● {label.upper()}", Inches(0.45), Inches(0.96), Inches(len(label)*0.11 + 0.3),
        Inches(0.28), size=8, bold=True, color=YELLOW)

def picture_card(sl, img_path, l, t, w, h, caption=None, sub=None):
    """Render a screenshot inside a polished frame with caption."""
    # Outer card
    rect(sl, l, t, w, h, fill_color=WHITE, line_color=YELLOW)
    # Browser chrome (thin bar)
    chrome_h = Inches(0.22)
    rect(sl, l, t, w, chrome_h, fill_color=RGBColor(0xEE, 0xEE, 0xF2),
         line_color=RGBColor(0xDD, 0xDD, 0xE5))
    txb(sl, "● ● ●   app.dynamoai.in", l + Inches(0.1), t + Inches(0.02),
        w - Inches(0.2), Inches(0.18), size=7, color=RGBColor(0x66, 0x66, 0x77))
    # Image area
    img_top = t + chrome_h
    img_h   = h - chrome_h - (Inches(0.5) if caption else Inches(0))
    if os.path.exists(img_path):
        # Add picture sized to fit width; centered vertically inside img area
        sl.shapes.add_picture(img_path, l, img_top, width=w, height=img_h)
    # Caption strip
    if caption:
        cap_top = t + h - Inches(0.5)
        rect(sl, l, cap_top, w, Inches(0.5),
             fill_color=BLACK, line_color=YELLOW)
        txb(sl, caption, l + Inches(0.12), cap_top + Inches(0.04),
            w - Inches(0.24), Inches(0.25),
            size=11, bold=True, color=YELLOW)
        if sub:
            txb(sl, sub, l + Inches(0.12), cap_top + Inches(0.26),
                w - Inches(0.24), Inches(0.22),
                size=8, color=MUTED)

def check(ok):
    return "✓" if ok == "y" else ("✗" if ok == "n" else "~")

def check_color(ok):
    return GREEN if ok == "y" else (RED if ok == "n" else YELLOW)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — Cover  (new positioning)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
if os.path.exists(LOGO_PATH):
    sl.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.25),
                          width=Inches(0.7), height=Inches(0.7))
txb(sl, "Dynamo AI", Inches(1.3), Inches(0.3), Inches(3), Inches(0.55),
    size=20, bold=True, color=WHITE)
txb(sl, "🇮🇳  MADE IN INDIA", Inches(10.8), Inches(0.35), Inches(2.1), Inches(0.35),
    size=9, bold=True, color=YELLOW)
rect(sl, Inches(0.5), Inches(1.55), Inches(2.8), Inches(0.33),
     fill_color=RGBColor(0x1A, 0x14, 0x00), line_color=RGBColor(0x80, 0x60, 0x03))
txb(sl, "● INCUBATOR PITCH · 2026", Inches(0.55), Inches(1.57), Inches(2.7), Inches(0.3),
    size=8, bold=True, color=YELLOW)
txb(sl, "Dynamo", Inches(0.5), Inches(2.0), Inches(5), Inches(1.3),
    size=90, bold=True, color=WHITE)
txb(sl, "AI", Inches(4.9), Inches(2.0), Inches(2.2), Inches(1.3),
    size=90, bold=True, color=YELLOW)
# Stronger positioning
txb(sl, "The Research Operating System for Academia.",
    Inches(0.5), Inches(3.5), Inches(12), Inches(0.55),
    size=24, bold=True, color=WHITE)
txb(sl, "RESEARCH  ·  WRITE  ·  CITE  ·  PRESENT  —  ALL IN ONE PLACE.",
    Inches(0.5), Inches(4.15), Inches(12), Inches(0.4),
    size=13, bold=True, color=YELLOW)
txb(sl, "Built for researchers, scholars and educators.",
    Inches(0.5), Inches(4.7), Inches(10), Inches(0.4),
    size=13, color=MUTED)
txb(sl, "FOUNDER: ANISH KRISNA S  ·  MS DATA SCIENCE",
    Inches(0.5), Inches(7.1), Inches(5), Inches(0.28), size=8, color=MUTED)
txb(sl, "APP.DYNAMOAI.IN", Inches(10.5), Inches(7.1), Inches(2.4), Inches(0.28),
    size=8, color=MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 2, "THE PROBLEM")
section_tag(sl, "THE PROBLEM")
txb(sl, "Academic work is broken.",
    Inches(0.4), Inches(1.35), Inches(9), Inches(0.75),
    size=38, bold=True, color=WHITE)
txb(sl, "Today's researchers juggle a dozen disconnected tools just to complete a single piece of academic work.",
    Inches(0.4), Inches(2.1), Inches(9), Inches(0.45), size=14, color=MUTED)

tools = [
    ("💬", "ChatGPT", "General AI chat"),
    ("🔎", "Google Scholar", "Paper search"),
    ("📚", "Research DBs", "JSTOR · Scopus"),
    ("📋", "Citation Tools", "Zotero · Mendeley"),
    ("📝", "Word", "Writing & drafts"),
    ("📊", "PowerPoint", "Presentations"),
]
for i, (icon, name, desc) in enumerate(tools):
    col = i % 3; row = i // 3
    lft = Inches(0.4 + col * 2.45)
    top = Inches(2.65 + row * 0.85)
    rect(sl, lft, top, Inches(2.3), Inches(0.75),
         fill_color=CARD, line_color=BORDER)
    txb(sl, f"{icon} {name}", lft + Inches(0.1), top + Inches(0.04), Inches(2.1), Inches(0.38),
        size=12, bold=True, color=WHITE)
    txb(sl, desc, lft + Inches(0.1), top + Inches(0.4), Inches(2.1), Inches(0.28),
        size=9, color=MUTED)

pains = [
    "Context switching across 6+ tools",
    "Lost productivity & wasted time",
    "Repeated effort & duplicate work",
    "Fragmented, disconnected workflow",
]
txb(sl, "THE RESULT", Inches(7.9), Inches(2.55), Inches(5), Inches(0.3),
    size=9, bold=True, color=YELLOW)
rect(sl, Inches(7.9), Inches(2.85), Inches(5.0), Inches(2.25),
     fill_color=RGBColor(0x14, 0x10, 0x00), line_color=RGBColor(0x60, 0x45, 0x00))
for i, pain in enumerate(pains):
    txb(sl, f"✗  {pain}", Inches(8.05), Inches(2.95 + i * 0.48), Inches(4.7), Inches(0.4),
        size=13, color=WHITE)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — Cost of Fragmentation  (SOFTENED — no unverifiable hour stats)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 3, "THE COST")
section_tag(sl, "THE COST OF FRAGMENTATION")
txb(sl, "The hidden tax on every researcher.",
    Inches(0.4), Inches(1.35), Inches(12), Inches(0.75), size=34, bold=True, color=WHITE)
txb(sl, "Fragmented tools don't just slow researchers down — they steal focus from the actual work that matters.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.45), size=14, color=MUTED)

pains = [
    ("🔄", "Context Switching",        "Constantly jumping between 6+ tools breaks flow and deep thinking."),
    ("📋", "Manual Formatting",        "Citation rules, paper formatting, and references rewritten by hand."),
    ("🗂️",  "Fragmented Notes",         "Findings scattered across tabs, docs, and offline notebooks."),
    ("⏱️",  "Lost Time, Lost Focus",    "Time spent stitching tools together — instead of producing insight."),
    ("🔁", "Duplicated Effort",        "Same searches, same summaries, same formatting — across every project."),
]
card_w = Inches(2.4)
for i, (icon, title, desc) in enumerate(pains):
    lft = Inches(0.4 + i * 2.58)
    top = Inches(2.6)
    rect(sl, lft, top, card_w, Inches(3.6), fill_color=CARD, line_color=BORDER)
    txb(sl, icon,  lft + Inches(0.15), top + Inches(0.2),  Inches(0.5), Inches(0.5), size=26)
    txb(sl, title, lft + Inches(0.15), top + Inches(0.85), card_w - Inches(0.3), Inches(0.5),
        size=13, bold=True, color=WHITE)
    txb(sl, desc,  lft + Inches(0.15), top + Inches(1.45),  card_w - Inches(0.3), Inches(2.0),
        size=10, color=MUTED)

rect(sl, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.75),
     fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
txb(sl, "⚡  The real cost isn't just hours — it's the insight, ideas and research output lost to busywork.",
    Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.55), size=12, color=WHITE)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 04 — Introducing Dynamo AI
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 4, "THE SOLUTION")
section_tag(sl, "INTRODUCING")
txb(sl, "One platform. One continuous academic workflow.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.75), size=36, bold=True, color=WHITE)
txb(sl, "Every step of your research — from discovery to final presentation — happens in one place.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

features = [
    ("🔬", "Discover relevant papers and evidence instantly",
     "Search across web + Semantic Scholar in seconds — no tab switching"),
    ("📖", "Identify insights, patterns and research gaps",
     "Analyse and summarise papers automatically, surface what matters"),
    ("✍️",  "Draft academic content with structured assistance",
     "Write with DeepThink AI — built for academic rigour and depth"),
    ("📋", "Generate publication-ready citations automatically",
     "APA, MLA, IEEE, Chicago, Harvard, Vancouver — formatted on demand"),
    ("📊", "Convert research into presentations, mind maps and flowcharts",
     "Go from notes to polished deck or visual in minutes"),
    ("🧠", "Preserve research context across documents and sessions",
     "AI Memory + Document Library keep your work accessible — always"),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 3; row = i // 3
    lft = Inches(0.4 + col * 4.3)
    top = Inches(2.65 + row * 1.65)
    rect(sl, lft, top, Inches(4.1), Inches(1.5),
         fill_color=CARD, line_color=RGBColor(0x3A, 0x30, 0x05))
    txb(sl, f"{icon}  {title}", lft + Inches(0.15), top + Inches(0.12),
        Inches(3.8), Inches(0.45), size=14, bold=True, color=WHITE)
    txb(sl, desc, lft + Inches(0.15), top + Inches(0.65), Inches(3.8), Inches(0.7),
        size=11, color=MUTED)

txb(sl, "20+ AI capabilities  ·  3 AI model tiers  ·  6 citation formats  ·  Built for researchers",
    Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.35),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 05 — Workflow
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 5, "THE WORKFLOW")
section_tag(sl, "THE DYNAMO AI WORKFLOW")
txb(sl, "One continuous workflow — start to finish.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.75), size=36, bold=True, color=WHITE)
txb(sl, "No more tab-switching. From research question to final deck — every step in one platform.",
    Inches(0.4), Inches(2.1), Inches(9.5), Inches(0.4), size=14, color=MUTED)

steps = [
    ("01", "🎯", "Define Topic"),
    ("02", "🔍", "Search"),
    ("03", "🧠", "Analyse"),
    ("04", "✍️",  "Write"),
    ("05", "📋", "Cite"),
    ("06", "📊", "Present"),
    ("07", "📂", "Store"),
]
step_w = Inches(1.6)
for i, (n, icon, title) in enumerate(steps):
    lft = Inches(0.4 + i * 1.83)
    top = Inches(2.7)
    is_first_last = i == 0 or i == len(steps)-1
    bg_c = RGBColor(0x20, 0x19, 0x00) if is_first_last else CARD
    bd_c = YELLOW if is_first_last else BORDER
    rect(sl, lft, top, step_w, Inches(2.6), fill_color=bg_c, line_color=bd_c)
    rect(sl, lft + Inches(0.45), top - Inches(0.18), Inches(0.7), Inches(0.28),
         fill_color=YELLOW, line_color=None)
    txb(sl, n, lft + Inches(0.45), top - Inches(0.18), Inches(0.7), Inches(0.28),
        size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(sl, icon,  lft + Inches(0.1), top + Inches(0.2), step_w - Inches(0.2), Inches(0.55),
        size=28, align=PP_ALIGN.CENTER)
    txb(sl, title, lft + Inches(0.1), top + Inches(0.9), step_w - Inches(0.2), Inches(0.5),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txb(sl, "›", lft + step_w + Inches(0.05), top + Inches(0.9), Inches(0.25), Inches(0.5),
            size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

rect(sl, Inches(3.0), Inches(5.7), Inches(7.3), Inches(0.45),
     fill_color=RGBColor(0x1A, 0x14, 0x00), line_color=RGBColor(0x80, 0x60, 0x03))
txb(sl, "⚡  Each step removes friction. The full workflow lives inside a single workspace.",
    Inches(3.0), Inches(5.73), Inches(7.3), Inches(0.38),
    size=11, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — Product Showcase Part 1 — The Workspace  (REAL SCREENSHOTS)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 6, "PRODUCT · PART 1")
section_tag(sl, "PRODUCT · THE WORKSPACE")
txb(sl, "A workspace built for thinking.",
    Inches(0.4), Inches(1.35), Inches(10), Inches(0.7), size=32, bold=True, color=WHITE)
txb(sl, "Clean, chat-first interface. Modes for every research task. Built and live today at app.dynamoai.in",
    Inches(0.4), Inches(2.05), Inches(11), Inches(0.4), size=13, color=MUTED)

# Main dashboard (large) — left ~7" wide
picture_card(sl, SHOT_DASHBOARD_PRO,
             Inches(0.4), Inches(2.6), Inches(7.5), Inches(4.0),
             caption="Main Workspace",
             sub="Chat-first interface · Recents · Quick Tools · Persistent memory across sessions")

# Plus / Mode menu — top right
picture_card(sl, SHOT_PLUS_MENU,
             Inches(8.1), Inches(2.6), Inches(4.8), Inches(1.9),
             caption="Mode Selector",
             sub="Fast · Research (Plus) · DeepThink (Pro)")

# Tools menu — bottom right
picture_card(sl, SHOT_TOOLS_MENU,
             Inches(8.1), Inches(4.7), Inches(4.8), Inches(1.9),
             caption="Tools Menu",
             sub="Study guides · Radio mode · Image, Slides, Mindmaps")

# Live badge
rect(sl, Inches(11.0), Inches(0.95), Inches(1.9), Inches(0.3),
     fill_color=RGBColor(0x05, 0x28, 0x12), line_color=GREEN)
txb(sl, "● LIVE IN PRODUCTION", Inches(11.05), Inches(0.97),
    Inches(1.85), Inches(0.28), size=8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txb(sl, "Try it live · app.dynamoai.in",
    Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.3),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 07 — Product Showcase Part 2 — Memory & Library  (REAL SCREENSHOTS)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 7, "PRODUCT · PART 2")
section_tag(sl, "PRODUCT · PERSISTENT INTELLIGENCE")
txb(sl, "It remembers. So you don't have to.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.7), size=32, bold=True, color=WHITE)
txb(sl, "Your research context persists across sessions, documents and projects.",
    Inches(0.4), Inches(2.05), Inches(12), Inches(0.45), size=13, color=MUTED)

# Profile / Quota
picture_card(sl, SHOT_PROFILE,
             Inches(0.4), Inches(2.6), Inches(3.9), Inches(4.0),
             caption="User Profile & Quotas",
             sub="Plan tier · Chat / Image / Video quotas · One-click upgrade")

# AI Memory
picture_card(sl, SHOT_AI_MEMORY,
             Inches(4.55), Inches(2.6), Inches(4.2), Inches(4.0),
             caption="AI Memory — Research Context Retention",
             sub="Dynamo remembers your topics, goals & interests across every session")

# Document Library
picture_card(sl, SHOT_DOC_LIBRARY,
             Inches(9.0), Inches(2.6), Inches(3.9), Inches(4.0),
             caption="Document Library",
             sub="Access all research materials in one place — summaries injected into every chat")

txb(sl, "Your research context remains available across sessions, projects and uploaded documents — automatically.",
    Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.3),
    size=10, color=YELLOW, align=PP_ALIGN.CENTER, bold=True)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — Researcher Journey (Before / After)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 8, "THE JOURNEY")
section_tag(sl, "THE RESEARCHER JOURNEY")
txb(sl, "The same workflow. A completely different experience.",
    Inches(0.4), Inches(1.35), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
txb(sl, "Dynamo AI doesn't add another tool. It replaces the entire fragmented workflow with one.",
    Inches(0.4), Inches(2.05), Inches(12), Inches(0.38), size=13, color=MUTED)

# ── LEFT COLUMN: Without Dynamo ───────────────────────────────────────────────
lc_l, lc_t, lc_w = Inches(0.4), Inches(2.55), Inches(5.8)
rect(sl, lc_l, lc_t, lc_w, Inches(4.55),
     fill_color=RGBColor(0x14, 0x08, 0x08), line_color=RED)
txb(sl, "✗  WITHOUT DYNAMO AI", lc_l + Inches(0.2), lc_t + Inches(0.15),
    lc_w - Inches(0.3), Inches(0.28), size=9, bold=True, color=RED)
txb(sl, "Multiple tools. Fragmented workflow.",
    lc_l + Inches(0.2), lc_t + Inches(0.45), lc_w - Inches(0.3), Inches(0.28),
    size=10, color=MUTED)

before_steps = [
    ("Search Papers",        "Google Scholar, databases"),
    ("Take Notes",           "Word, Notion, sticky notes"),
    ("Write Draft",          "MS Word / Docs"),
    ("Generate Citations",   "Zotero / Mendeley"),
    ("Create Slides",        "PowerPoint / Canva"),
    ("Store Documents",      "Drive / email to self"),
]
for i, (step, tool) in enumerate(before_steps):
    sy = lc_t + Inches(0.85 + i * 0.58)
    rect(sl, lc_l + Inches(0.2), sy, lc_w - Inches(0.4), Inches(0.48),
         fill_color=RGBColor(0x20, 0x10, 0x10), line_color=RGBColor(0x55, 0x20, 0x20))
    txb(sl, f"→  {step}", lc_l + Inches(0.32), sy + Inches(0.04),
        Inches(2.2), Inches(0.25), size=11, bold=True, color=WHITE)
    txb(sl, tool, lc_l + Inches(0.32), sy + Inches(0.26),
        Inches(3.0), Inches(0.2), size=8, color=RED)
    if i < len(before_steps) - 1:
        txb(sl, "↕  switch app", lc_l + Inches(0.32),
            sy + Inches(0.5), Inches(2), Inches(0.18),
            size=7, color=RGBColor(0x77, 0x33, 0x33))

# ── BIG ARROW centre ──────────────────────────────────────────────────────────
txb(sl, "⟹", Inches(6.35), Inches(4.5), Inches(0.8), Inches(0.6),
    size=36, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ── RIGHT COLUMN: With Dynamo ─────────────────────────────────────────────────
rc_l = Inches(7.3)
rect(sl, rc_l, lc_t, lc_w, Inches(4.55),
     fill_color=RGBColor(0x06, 0x18, 0x0c), line_color=GREEN)
txb(sl, "✓  WITH DYNAMO AI", rc_l + Inches(0.2), lc_t + Inches(0.15),
    lc_w - Inches(0.3), Inches(0.28), size=9, bold=True, color=GREEN)
txb(sl, "Single platform. Single workflow.",
    rc_l + Inches(0.2), lc_t + Inches(0.45), lc_w - Inches(0.3), Inches(0.28),
    size=10, color=MUTED)

after_steps = [
    ("Discover evidence-backed papers",  "Evidence-backed research assistance"),
    ("Analyse & take smart notes",       "Insights, gaps and summaries — instant"),
    ("Draft academic content",           "Structured assistance with DeepThink AI"),
    ("Generate citations automatically", "6 formats — APA, MLA, IEEE and more"),
    ("Convert research to a deck",       "Slides, mindmaps, flowcharts — in seconds"),
    ("Store in Document Library",        "Access all research materials in one place"),
]
for i, (step, tool) in enumerate(after_steps):
    sy = lc_t + Inches(0.85 + i * 0.58)
    rect(sl, rc_l + Inches(0.2), sy, lc_w - Inches(0.4), Inches(0.48),
         fill_color=RGBColor(0x08, 0x22, 0x10), line_color=RGBColor(0x18, 0x55, 0x28))
    txb(sl, f"✓  {step}", rc_l + Inches(0.32), sy + Inches(0.04),
        Inches(2.8), Inches(0.25), size=11, bold=True, color=WHITE)
    txb(sl, tool, rc_l + Inches(0.32), sy + Inches(0.26),
        Inches(3.2), Inches(0.2), size=8, color=GREEN)

# ── Bottom callout ────────────────────────────────────────────────────────────
rect(sl, Inches(0.4), Inches(7.18), Inches(12.5), Inches(0.2),
     fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
txb(sl, "Dynamo AI is not another AI chatbot. It is the academic workflow platform where researchers do all their work.",
    Inches(0.6), Inches(7.2), Inches(12.1), Inches(0.18),
    size=9, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — Why Dynamo Wins  (NEW)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 9, "WHY RESEARCHERS CHOOSE DYNAMO AI")
section_tag(sl, "WHY RESEARCHERS CHOOSE DYNAMO AI")
txb(sl, "Why researchers choose Dynamo AI.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.75), size=36, bold=True, color=WHITE)
txb(sl, "Dynamo AI is not another AI chatbot. It is the academic workflow platform where researchers do all their work.",
    Inches(0.4), Inches(2.1), Inches(12), Inches(0.4), size=14, color=MUTED)

reasons = [
    ("🎓", "Built for Academic Workflows",
     "Not a generic AI chat. Every feature — modes, citation engine, study guides — is shaped around how researchers actually work."),
    ("🔗", "Research → Write → Cite → Present",
     "The full academic pipeline lives inside a single workspace. Switching tools is the problem we exist to remove."),
    ("🧠", "Academic Memory Layer",
     "Research context remains available across sessions, projects and uploaded documents — automatically, without any manual effort."),
    ("📋", "Generate publication-ready citations automatically",
     "6 academic citation formats, structured assistance, and evidence-backed research in one place."),
    ("👨‍🔬", "Built by Someone Who Lived It",
     "Founder is a data scientist who built Dynamo to solve a workflow problem he experienced first-hand."),
]
for i, (icon, title, desc) in enumerate(reasons):
    # 2 + 3 layout — top row 2 wide cards, bottom row 3
    if i < 2:
        col = i; row = 0
        lft = Inches(0.4 + col * 6.45)
        cw  = Inches(6.25)
    else:
        col = i - 2; row = 1
        lft = Inches(0.4 + col * 4.3)
        cw  = Inches(4.1)
    top = Inches(2.6 + row * 2.05)
    ch  = Inches(1.85)
    rect(sl, lft, top, cw, ch,
         fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
    txb(sl, icon, lft + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.55), size=26)
    txb(sl, title, lft + Inches(0.95), top + Inches(0.18), cw - Inches(1.1), Inches(0.5),
        size=14, bold=True, color=YELLOW)
    txb(sl, desc, lft + Inches(0.2), top + Inches(0.8), cw - Inches(0.35), Inches(1.0),
        size=10, color=WHITE)

txb(sl, "We don't compete on AI capability. We compete on workflow ownership.",
    Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.3),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Competitive Landscape  (SOFTENED)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 10, "COMPETITION")
section_tag(sl, "COMPETITIVE LANDSCAPE")
txb(sl, "One workflow. Everyone else is a slice.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.75), size=34, bold=True, color=WHITE)
txb(sl, "Other tools each solve a single piece of the research workflow. Dynamo AI brings them together in a single workspace.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=13, color=MUTED)

cols_hdr = ["Researcher Outcome", "ChatGPT", "Consensus", "Jenni", "⚡ Dynamo AI"]
col_ws   = [4.2, 1.8, 1.8, 1.8, 1.95]
rows_data = [
    ["Find evidence from research papers",      "~", "y", "n", "y"],
    ["Write academic research papers",          "y", "n", "y", "y"],
    ["Manage research knowledge & memory",      "n", "n", "n", "y"],
    ["Generate formatted citations",            "~", "y", "y", "y"],
    ["Create presentations & visuals",          "n", "n", "n", "y"],
    ["Complete research workflow end-to-end",   "n", "n", "n", "y"],
]
row_h = Inches(0.46)
lft0 = Inches(0.4)
top0 = Inches(2.65)
x = lft0
for j, (col_name, cw) in enumerate(zip(cols_hdr, col_ws)):
    bg = RGBColor(0x20, 0x19, 0x00) if j == 5 else BLACK2
    rect(sl, x, top0, Inches(cw), row_h, fill_color=bg, line_color=BORDER)
    c = YELLOW if j == 5 else MUTED
    txb(sl, col_name, x + Inches(0.08), top0 + Inches(0.1), Inches(cw - 0.1), Inches(0.28),
        size=9, bold=True, color=c, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    x += Inches(cw)
for ri, row in enumerate(rows_data):
    top = top0 + row_h * (ri + 1)
    x = lft0
    for j, (val, cw) in enumerate(zip(row, col_ws)):
        bg = RGBColor(0x20, 0x19, 0x00) if j == 5 else (
             CARD if ri % 2 == 0 else RGBColor(0x0E, 0x13, 0x22))
        rect(sl, x, top, Inches(cw), row_h, fill_color=bg, line_color=BORDER)
        if j == 0:
            txb(sl, val, x + Inches(0.1), top + Inches(0.1), Inches(cw - 0.15), Inches(0.28),
                size=12, bold=True, color=WHITE)
        else:
            sym_color = check_color(val)
            txb(sl, check(val), x, top + Inches(0.07), Inches(cw), Inches(0.32),
                size=14, bold=True, color=sym_color, align=PP_ALIGN.CENTER)
        x += Inches(cw)

rect(sl, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.65),
     fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
txb(sl, "💡  Dynamo AI is the only platform where a researcher can complete their entire workflow — without switching tools.",
    Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.45), size=11, color=WHITE)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Why Now  (SOFTENED — qualitative claims, no unsourced stats)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 11, "WHY NOW")
section_tag(sl, "WHY NOW")
txb(sl, "Three forces. One window.",
    Inches(0.4), Inches(1.35), Inches(9), Inches(0.75), size=38, bold=True, color=WHITE)
txb(sl, "Academia is at an inflection point. Researchers need integrated AI tools — and they need them now.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

trends = [
    ("📈", "01", "Research Output is Exploding",
     "Exponential growth",
     "Global research volume is growing exponentially. Manual literature reviews don't scale to the pace of modern science."),
    ("🤖", "02", "AI Adoption is Mainstream",
     "Already in workflows",
     "AI tools are now part of daily academic life. But most are generic — none are designed around how researchers actually work."),
    ("⚡", "03", "Workflow Tools are Demanded",
     "Integration is the gap",
     "Researchers are vocally tired of stitching 6+ tools together. The next winner won't add another chatbot — it'll consolidate the workflow."),
]
for i, (icon, num, title, stat, desc) in enumerate(trends):
    lft = Inches(0.4 + i * 4.3)
    top = Inches(2.6)
    rect(sl, lft, top, Inches(4.1), Inches(3.8),
         fill_color=RGBColor(0x18, 0x14, 0x00), line_color=RGBColor(0x60, 0x45, 0x00))
    txb(sl, icon, lft + Inches(0.2), top + Inches(0.2), Inches(0.6), Inches(0.55), size=30)
    txb(sl, num,  lft + Inches(3.5), top + Inches(0.2), Inches(0.5), Inches(0.4),
        size=22, color=RGBColor(0x60, 0x45, 0x03))
    txb(sl, title, lft + Inches(0.2), top + Inches(0.95), Inches(3.8), Inches(0.9),
        size=18, bold=True, color=WHITE)
    txb(sl, stat, lft + Inches(0.2), top + Inches(1.95), Inches(3.8), Inches(0.45),
        size=20, bold=True, color=YELLOW)
    txb(sl, desc, lft + Inches(0.2), top + Inches(2.5), Inches(3.7), Inches(1.2),
        size=10, color=MUTED)

txb(sl, "The market is ready. Researchers want integrated, intelligent tools — built for them, in India.",
    Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.4),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Target Market (TAM / SAM / SOM)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 12, "MARKET")
section_tag(sl, "TARGET MARKET · TAM / SAM / SOM")
txb(sl, "Big market. Sharp entry point.",
    Inches(0.4), Inches(1.35), Inches(10), Inches(0.75), size=38, bold=True, color=WHITE)
txb(sl, "Start with the highest-pain users. Expand outward to institutions.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

tiers = [
    ("TAM", "Total Addressable Market",
     "Researchers, educators and institutions across India using digital academic tools.",
     "India-wide", "academic research & education sector", MUTED,
     RGBColor(0x12, 0x17, 0x28), RGBColor(0x30, 0x35, 0x50)),
    ("SAM", "Serviceable Available Market",
     "India's 200K+ PhD scholars + 1.5M professors & academic researchers.",
     "1.7M+", "active researchers & faculty (India)", WHITE,
     CARD, RGBColor(0x60, 0x45, 0x00)),
    ("SOM", "Serviceable Obtainable Market",
     "First wedge: 10,000 paying scholars & researchers in 24 months.",
     "10K", "users · ₹4–10 Cr ARR target by Year 2", YELLOW,
     RGBColor(0x20, 0x19, 0x00), YELLOW),
]
for i, (badge, name, desc, big, sub, big_c, bg_c, bd_c) in enumerate(tiers):
    lft = Inches(0.4 + i * 4.3)
    top = Inches(2.6)
    rect(sl, lft, top, Inches(4.1), Inches(4.1), fill_color=bg_c, line_color=bd_c)
    # Badge pill
    rect(sl, lft + Inches(0.2), top + Inches(0.2), Inches(0.8), Inches(0.35),
         fill_color=YELLOW if i == 2 else BLACK2, line_color=YELLOW if i == 2 else BORDER)
    txb(sl, badge, lft + Inches(0.2), top + Inches(0.22), Inches(0.8), Inches(0.32),
        size=11, bold=True, color=BLACK if i == 2 else YELLOW, align=PP_ALIGN.CENTER)
    txb(sl, name, lft + Inches(0.2), top + Inches(0.65), Inches(3.8), Inches(0.35),
        size=12, bold=True, color=WHITE)
    txb(sl, desc, lft + Inches(0.2), top + Inches(1.05), Inches(3.8), Inches(1.0),
        size=10, color=MUTED)
    txb(sl, big, lft + Inches(0.2), top + Inches(2.15), Inches(3.8), Inches(0.9),
        size=42, bold=True, color=big_c)
    txb(sl, sub, lft + Inches(0.2), top + Inches(3.15), Inches(3.8), Inches(0.7),
        size=10, color=MUTED)

txb(sl, "Wedge strategy: dominate PhD scholars first → expand to professors → institutional contracts.",
    Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.3),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Business Model
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 13, "BUSINESS MODEL")
section_tag(sl, "BUSINESS MODEL")
txb(sl, "Predictable revenue. Scalable margins.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.75), size=36, bold=True, color=WHITE)
txb(sl, "Subscription SaaS today. Institutional licensing tomorrow. INR-priced for the Indian market.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

plans = [
    ("Free Forever", "₹0", "/mo", "Acquisition layer", WHITE,
     CARD, RGBColor(0x30, 0x35, 0x50),
     ["10 AI chats / day", "Fast Mode (Gemini Flash)", "Web search + voice input",
      "AI Detector & Plagiarism", "Mindmaps & Flowcharts"]),
    ("Plus", "₹399", "/mo", "For scholars & researchers", YELLOW,
     RGBColor(0x20, 0x19, 0x00), YELLOW,
     ["✓ Unlimited AI chats", "✓ Research Mode (3 models)", "✓ 6-Format Citation Engine",
      "✓ Document Library", "✓ Find Research Gaps"]),
    ("Pro", "₹999", "/mo", "Power users & professors", WHITE,
     CARD, RGBColor(0x30, 0x35, 0x50),
     ["• Everything in Plus", "• DeepThink Mode", "• Deep Research Agent",
      "• Unlimited PDF uploads", "• Advanced academic workflows"]),
]
for i, (name, price, per, sub, txt_c, bg_c, bd_c, feats) in enumerate(plans):
    lft = Inches(0.4 + i * 4.3)
    top = Inches(2.55)
    rect(sl, lft, top, Inches(4.1), Inches(4.15), fill_color=bg_c, line_color=bd_c)
    if i == 1:
        rect(sl, lft + Inches(1.1), top - Inches(0.18), Inches(1.9), Inches(0.28),
             fill_color=YELLOW, line_color=None)
        txb(sl, "MOST POPULAR", lft + Inches(1.1), top - Inches(0.18), Inches(1.9), Inches(0.28),
            size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(sl, name,  lft + Inches(0.2), top + Inches(0.15), Inches(3.8), Inches(0.32),
        size=11, bold=True, color=txt_c)
    txb(sl, price, lft + Inches(0.2), top + Inches(0.55), Inches(3.8), Inches(0.8),
        size=48, bold=True, color=WHITE)
    txb(sl, per,   lft + Inches(1.85), top + Inches(0.85), Inches(1.0), Inches(0.3),
        size=12, color=MUTED)
    txb(sl, sub,   lft + Inches(0.2), top + Inches(1.45), Inches(3.8), Inches(0.28),
        size=10, color=MUTED)
    for j, feat in enumerate(feats):
        txb(sl, feat, lft + Inches(0.2), top + Inches(1.85 + j * 0.42),
            Inches(3.8), Inches(0.35), size=11, color=WHITE if i == 1 else RGBColor(0xCC,0xCC,0xDD))

for i, (icon, lbl, val) in enumerate([
    ("🏛️", "Institutional Licensing", "₹5L–50L / year per university"),
    ("🤝", "University Partnerships",  "Embedded into curriculum"),
]):
    lft = Inches(0.4 + i * 6.6)
    top = Inches(6.9)
    rect(sl, lft, top, Inches(6.2), Inches(0.45),
         fill_color=RGBColor(0x18, 0x14, 0x00), line_color=RGBColor(0x60, 0x45, 0x00))
    txb(sl, f"{icon}  FUTURE · YEAR 2  —  {lbl}: {val}", lft + Inches(0.15), top + Inches(0.08),
        Inches(6.0), Inches(0.3), size=10, color=YELLOW)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Early Validation  (UNCHANGED — founder will edit numbers themselves)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 14, "VALIDATION")
section_tag(sl, "EARLY VALIDATION")
txb(sl, "Real users. Real traction.",
    Inches(0.4), Inches(1.35), Inches(9), Inches(0.75), size=38, bold=True, color=WHITE)
txb(sl, "We're live. We have users. They're using it. Here's where we are and where we're heading.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

rect(sl, Inches(0.4), Inches(2.6), Inches(6.0), Inches(4.1),
     fill_color=CARD, line_color=RGBColor(0x30, 0x35, 0x50))
txb(sl, "● TODAY · LIVE", Inches(0.55), Inches(2.72), Inches(3), Inches(0.3),
    size=10, bold=True, color=GREEN)
metrics = [("[X]", "Beta users"), ("100%", "MVP shipped"), ("20+", "AI features live"), ("[X]", "Professor endorsements")]
for i, (stat, lbl) in enumerate(metrics):
    col = i % 2; row = i // 2
    ml = Inches(0.55 + col * 2.85)
    mt = Inches(3.15 + row * 0.9)
    rect(sl, ml, mt, Inches(2.65), Inches(0.78),
         fill_color=RGBColor(0x0E, 0x13, 0x22), line_color=BORDER)
    txb(sl, stat, ml + Inches(0.12), mt + Inches(0.04), Inches(2.4), Inches(0.42),
        size=28, bold=True, color=YELLOW)
    txb(sl, lbl,  ml + Inches(0.12), mt + Inches(0.48), Inches(2.4), Inches(0.26),
        size=10, color=MUTED)
rect(sl, Inches(0.55), Inches(5.2), Inches(5.65), Inches(0.9),
     fill_color=RGBColor(0x08, 0x0C, 0x18), line_color=YELLOW)
txb(sl, '"[ Add real user testimonial here ]"',
    Inches(0.7), Inches(5.3), Inches(5.4), Inches(0.5), size=11, color=WHITE)
txb(sl, "— [ Add attribution ]",
    Inches(0.7), Inches(5.82), Inches(5.4), Inches(0.25), size=9, color=MUTED)

rect(sl, Inches(6.8), Inches(2.6), Inches(6.1), Inches(4.1),
     fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
txb(sl, "● NEXT 90 DAYS", Inches(6.95), Inches(2.72), Inches(3), Inches(0.3),
    size=10, bold=True, color=YELLOW)
next_items = [
    ("50",   "Beta users onboarded",   "Scaling from launch to active researchers"),
    ("5+",   "Video testimonials",     "Professor + PhD scholar validation videos"),
    ("100",  "Paid users target",      "First ₹1L+ MRR from organic conversion"),
    ("3",    "University pilots",      "Initial institutional conversations"),
]
for i, (stat, lbl, desc) in enumerate(next_items):
    top = Inches(3.15 + i * 0.88)
    txb(sl, stat, Inches(6.95), top, Inches(0.9), Inches(0.4), size=22, bold=True, color=YELLOW)
    txb(sl, lbl,  Inches(7.95), top, Inches(4.8), Inches(0.3), size=12, bold=True, color=WHITE)
    txb(sl, desc, Inches(7.95), top + Inches(0.32), Inches(4.8), Inches(0.3), size=9, color=MUTED)
    if i < 3:
        rect(sl, Inches(6.95), top + Inches(0.82), Inches(5.8), Inches(0.02),
             fill_color=RGBColor(0x30, 0x28, 0x00))
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 15, "ROADMAP")
section_tag(sl, "THE ROADMAP")
txb(sl, "From MVP to market leader — in 18 months.",
    Inches(0.4), Inches(1.35), Inches(12), Inches(0.75), size=34, bold=True, color=WHITE)
txb(sl, "A focused, milestone-driven plan from where we are today to a category-defining position.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

phases = [
    ("Phase 1", "Q2–Q3 2026",       "Validate & Refine",  GREEN,
     ["100 paid users (₹1L+ MRR)", "5+ video testimonials", "Product-market fit signals", "Feature prioritisation"]),
    ("Phase 2", "Q4 2026 – Q1 2027","Scale & Expand",     YELLOW,
     ["Mobile apps (iOS + Android)", "First 3 university pilots", "1,000 paid users (₹10L+ MRR)", "Hindi language support"]),
    ("Phase 3", "Q2–Q4 2027",       "Institutional Play", BLUE,
     ["10,000+ active users", "10+ institutional contracts", "Team scale to 15", "National Academic Adoption"]),
]
for i, (phase, q, title, color, items) in enumerate(phases):
    lft = Inches(0.4 + i * 4.3)
    top = Inches(2.6)
    r = rect(sl, lft, top, Inches(4.1), Inches(4.1), fill_color=CARD, line_color=color)
    r.line.width = Pt(2)
    txb(sl, phase, lft + Inches(0.2), top + Inches(0.15), Inches(2), Inches(0.3),
        size=10, bold=True, color=color)
    txb(sl, q, lft + Inches(2.2), top + Inches(0.15), Inches(1.75), Inches(0.3),
        size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    txb(sl, title, lft + Inches(0.2), top + Inches(0.55), Inches(3.8), Inches(0.42),
        size=16, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txb(sl, f"›  {item}", lft + Inches(0.2), top + Inches(1.15 + j * 0.65),
            Inches(3.8), Inches(0.5), size=12, color=RGBColor(0xCC, 0xCC, 0xDD))
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Founder
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 16, "FOUNDER")
section_tag(sl, "WHO'S BUILDING THIS")
txb(sl, "Built by someone who lived the problem.",
    Inches(0.4), Inches(1.35), Inches(12), Inches(0.75), size=34, bold=True, color=WHITE)

rect(sl, Inches(0.4), Inches(2.55), Inches(3.8), Inches(4.3),
     fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
av = sl.shapes.add_shape(9, Inches(1.4), Inches(2.75), Inches(1.8), Inches(1.8))
av.fill.solid(); av.fill.fore_color.rgb = YELLOW
av.line.fill.background()
txb(sl, "AK", Inches(1.4), Inches(2.95), Inches(1.8), Inches(1.5),
    size=46, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txb(sl, "Anish Krisna S", Inches(0.55), Inches(4.65), Inches(3.5), Inches(0.45),
    size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Founder & CEO", Inches(0.55), Inches(5.15), Inches(3.5), Inches(0.3),
    size=10, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
rect(sl, Inches(0.55), Inches(5.55), Inches(3.5), Inches(0.02),
     fill_color=RGBColor(0x40, 0x30, 0x00))
txb(sl, "EDUCATION: MS in Data Science", Inches(0.55), Inches(5.65), Inches(3.5), Inches(0.28),
    size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "EXPERIENCE: 11 Years cross-functional", Inches(0.55), Inches(5.98), Inches(3.5), Inches(0.28),
    size=9, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

skills = [("📊", "Analytics", "Data-driven product decisions"),
          ("📣", "Marketing", "Brand, growth & demand gen"),
          ("🛠️", "Product",    "End-to-end product building")]
for i, (icon, skill, desc) in enumerate(skills):
    lft = Inches(4.6 + i * 2.95)
    rect(sl, lft, Inches(2.6), Inches(2.75), Inches(1.6),
         fill_color=CARD, line_color=BORDER)
    txb(sl, icon, lft + Inches(0.15), Inches(2.75), Inches(0.55), Inches(0.5), size=24)
    txb(sl, skill, lft + Inches(0.15), Inches(3.35), Inches(2.5), Inches(0.35),
        size=13, bold=True, color=WHITE)
    txb(sl, desc, lft + Inches(0.15), Inches(3.75), Inches(2.5), Inches(0.3),
        size=9, color=MUTED)

rect(sl, Inches(4.6), Inches(4.45), Inches(8.75), Inches(1.5),
     fill_color=RGBColor(0x18, 0x14, 0x00), line_color=YELLOW)
txb(sl, "MISSION", Inches(4.75), Inches(4.58), Inches(5), Inches(0.25),
    size=9, bold=True, color=YELLOW)
txb(sl, '"Build the operating system for academic work — so India\'s brightest minds spend their time on insight, not on busywork."',
    Inches(4.75), Inches(4.9), Inches(8.4), Inches(0.9), size=13, color=WHITE)

# Why I Built section
rect(sl, Inches(4.6), Inches(5.75), Inches(8.75), Inches(1.35),
     fill_color=CARD, line_color=BORDER)
txb(sl, "WHY I BUILT DYNAMO AI", Inches(4.75), Inches(5.88), Inches(5), Inches(0.25),
    size=9, bold=True, color=YELLOW)
txb(sl, "I experienced how fragmented academic workflows force researchers to spend more time managing tools than producing insight. Dynamo AI was created to unify research, writing, citations and knowledge management into a single academic workspace — so scholars can focus on what actually matters.",
    Inches(4.75), Inches(6.15), Inches(8.4), Inches(0.88), size=10, color=WHITE)
footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — The Ask
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, 17, "THE ASK")
section_tag(sl, "WHAT WE NEED FROM YOU")
txb(sl, "Let's build India's academic operating system — together.",
    Inches(0.4), Inches(1.35), Inches(12), Inches(0.75), size=32, bold=True, color=WHITE)
txb(sl, "We have the product, the vision, and early traction. We're looking for the right partners to scale.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

asks = [
    ("🎓", "Mentorship",            "Strategic guidance from founders who've scaled SaaS in India"),
    ("🚀", "Incubation",            "Workspace, infrastructure, and ecosystem support"),
    ("🏛️", "Academic Partnerships", "Warm intros to universities & research institutions"),
    ("🧪", "Pilot Institutions",    "3–5 partner institutions for institutional MVP"),
    ("🧭", "Strategic Guidance",    "Go-to-market, hiring, and fundraising playbooks"),
]
for i, (icon, title, desc) in enumerate(asks):
    lft = Inches(0.4 + i * 2.58)
    top = Inches(2.65)
    rect(sl, lft, top, Inches(2.4), Inches(2.8),
         fill_color=RGBColor(0x18, 0x14, 0x00), line_color=RGBColor(0x60, 0x45, 0x00))
    txb(sl, icon, lft + Inches(0.1), top + Inches(0.2), Inches(2.2), Inches(0.55),
        size=28, align=PP_ALIGN.CENTER)
    txb(sl, title, lft + Inches(0.1), top + Inches(0.9), Inches(2.2), Inches(0.38),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, desc, lft + Inches(0.1), top + Inches(1.35), Inches(2.2), Inches(1.1),
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.85),
     fill_color=RGBColor(0x20, 0x19, 0x00), line_color=YELLOW)
txb(sl, "LET'S TALK", Inches(0.6), Inches(5.92), Inches(2), Inches(0.28),
    size=9, bold=True, color=YELLOW)
txb(sl, "Try it live · See the product · Then decide.",
    Inches(0.6), Inches(6.22), Inches(5.5), Inches(0.35), size=16, bold=True, color=WHITE)
txb(sl, "Demo: app.dynamoai.in", Inches(8.5), Inches(5.95), Inches(3.9), Inches(0.28),
    size=11, color=WHITE)
txb(sl, "Email: anish@dynamoai.in", Inches(8.5), Inches(6.28), Inches(3.9), Inches(0.28),
    size=11, color=WHITE)

txb(sl, "⚡   THANK YOU   ⚡",
    Inches(0), Inches(7.0), Inches(13.33), Inches(0.38),
    size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "Dynamo_AI_Pitch_Deck.pptx"
prs.save(out)
print(f"✅ Saved: {out}  ({os.path.getsize(out)//1024} KB · {TOTAL_SLIDES} slides)")
