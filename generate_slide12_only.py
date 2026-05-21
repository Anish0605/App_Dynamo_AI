"""
Dynamo AI — Slide 12 (Market: TAM / SAM / SOM) — Standalone PPTX
Single slide, exactly matching the full deck's Slide 12.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BLACK  = RGBColor(0x0a, 0x0e, 0x1a)
BLACK2 = RGBColor(0x14, 0x19, 0x2b)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xAA, 0xAA, 0xBB)
CARD   = RGBColor(0x12, 0x17, 0x28)
BORDER = RGBColor(0x25, 0x2A, 0x40)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

LOGO_PATH = "attached_assets/Dynamo_AI_New_Logo_1779360398074.png"

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return sl

def txb(sl, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return tx

def rect(sl, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = sl.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = line_width
    else:
        shape.line.fill.background()
    return shape

# ── Build Slide 12 ────────────────────────────────────────────────────────────
sl = add_slide()

# Header
if os.path.exists(LOGO_PATH):
    sl.shapes.add_picture(LOGO_PATH, Inches(0.4), Inches(0.2),
                          width=Inches(0.5), height=Inches(0.5))
txb(sl, "Dynamo AI", Inches(1.0), Inches(0.25), Inches(2), Inches(0.4),
    size=16, bold=True, color=WHITE)
txb(sl, "MARKET", Inches(4), Inches(0.28), Inches(5.33), Inches(0.35),
    size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txb(sl, "12 / 17", Inches(12.3), Inches(0.28), Inches(0.9), Inches(0.35),
    size=9, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

# Section tag
rect(sl, Inches(0.4), Inches(0.95), Inches(3.8), Inches(0.3),
     fill_color=RGBColor(0x26, 0x1E, 0x00), line_color=RGBColor(0x80, 0x60, 0x03))
txb(sl, "● TARGET MARKET · TAM / SAM / SOM", Inches(0.45), Inches(0.96),
    Inches(3.7), Inches(0.28), size=8, bold=True, color=YELLOW)

# Headline
txb(sl, "Big market. Sharp entry point.",
    Inches(0.4), Inches(1.35), Inches(11), Inches(0.75), size=38, bold=True, color=WHITE)
txb(sl, "Start with the highest-pain users. Expand outward to institutions.",
    Inches(0.4), Inches(2.1), Inches(11), Inches(0.4), size=14, color=MUTED)

# ── Three tier cards ──────────────────────────────────────────────────────────
tiers = [
    ("TAM", "Total Addressable Market",
     "All researchers, educators & students using digital tools across India.",
     "₹2,400 Cr+", "estimated annual academic productivity market (India)",
     MUTED, CARD, RGBColor(0x30, 0x35, 0x50)),
    ("SAM", "Serviceable Available Market",
     "India's 200K+ PhD scholars + 1.5M professors & academic researchers.\n(Source: AISHE 2022-23, Ministry of Education, Govt. of India)",
     "1.7M+", "active researchers & faculty (India)",
     WHITE, CARD, RGBColor(0x60, 0x45, 0x00)),
    ("SOM", "Serviceable Obtainable Market",
     "First wedge: 10,000 paying scholars & researchers in 24 months.\n(0.6% penetration of SAM — conservative target)",
     "10K", "users · ₹4–10 Cr ARR target by Year 2",
     YELLOW, RGBColor(0x20, 0x19, 0x00), YELLOW),
]

for i, (badge, name, desc, big, sub, big_c, bg_c, bd_c) in enumerate(tiers):
    lft = Inches(0.4 + i * 4.3)
    top = Inches(2.6)
    rect(sl, lft, top, Inches(4.1), Inches(4.1), fill_color=bg_c, line_color=bd_c)

    # Badge pill
    pill_bg = YELLOW if i == 2 else BLACK2
    pill_bd = YELLOW if i == 2 else BORDER
    rect(sl, lft + Inches(0.2), top + Inches(0.2), Inches(0.8), Inches(0.35),
         fill_color=pill_bg, line_color=pill_bd)
    txb(sl, badge, lft + Inches(0.2), top + Inches(0.22), Inches(0.8), Inches(0.32),
        size=11, bold=True,
        color=BLACK if i == 2 else YELLOW,
        align=PP_ALIGN.CENTER)

    txb(sl, name, lft + Inches(0.2), top + Inches(0.65), Inches(3.8), Inches(0.35),
        size=12, bold=True, color=WHITE)
    txb(sl, desc, lft + Inches(0.2), top + Inches(1.05), Inches(3.8), Inches(1.1),
        size=10, color=MUTED)
    txb(sl, big, lft + Inches(0.2), top + Inches(2.25), Inches(3.8), Inches(0.9),
        size=42, bold=True, color=big_c)
    txb(sl, sub, lft + Inches(0.2), top + Inches(3.25), Inches(3.8), Inches(0.65),
        size=10, color=MUTED)

# Bottom note — added source callout
rect(sl, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.35),
     fill_color=RGBColor(0x12, 0x17, 0x28), line_color=BORDER)
txb(sl, "Wedge strategy: dominate PhD scholars first → expand to professors → institutional contracts.  |  SAM source: AISHE 2022-23, Ministry of Education, Govt. of India",
    Inches(0.55), Inches(6.93), Inches(12.2), Inches(0.28),
    size=8, color=MUTED, align=PP_ALIGN.CENTER)

# Footer
txb(sl, "DYNAMO AI · CONFIDENTIAL",
    Inches(0.4), Inches(7.22), Inches(4), Inches(0.22), size=7, color=MUTED)
txb(sl, "app.dynamoai.in",
    Inches(9.5), Inches(7.22), Inches(3.4), Inches(0.22),
    size=7, color=MUTED, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "Dynamo_AI_Slide12_Market.pptx"
prs.save(out)
print(f"✅ Saved: {out}  ({os.path.getsize(out)//1024} KB · 1 slide)")
