"""
pitch_export.py — Generates a 16:9 PPTX investor deck for Dynamo AI.
Uses python-pptx. No external images required (logo replaced by styled text).
"""

import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────────
W = Inches(13.333)  # 16:9 width  (standard widescreen)
H = Inches(7.5)     # 16:9 height

# Colours
BG     = RGBColor(0xFA, 0xFA, 0xF7)   # off-white slide background
YELLOW = RGBColor(0xF5, 0xC4, 0x00)   # brand yellow
DARK   = RGBColor(0x1A, 0x1A, 0x1A)   # near-black
MUTED  = RGBColor(0x6B, 0x6B, 0x6B)   # secondary text
GOLD   = RGBColor(0xB3, 0x8A, 0x00)   # dark-gold for yellow-on-white text
BGCARD = RGBColor(0xFF, 0xFF, 0xFF)   # card white
LINE   = RGBColor(0xE0, 0xE0, 0xDA)   # card border colour (not used as colour param)
WARMYELLOW = RGBColor(0xFF, 0xFB, 0xEA) # warm yellow bg for highlighted cards
INVBG  = RGBColor(0x1A, 0x1A, 0x1A)   # inverted card background

# Fonts
FONT = "Calibri"


# ─────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────

def _bg(slide, colour=BG):
    """Fill the slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _box(slide, left, top, width, height,
         text="", font_size=18, bold=False, colour=DARK,
         align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name=None):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name or FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return tf


def _rect(slide, left, top, width, height, fill_colour, line_colour=None):
    """Add a filled rectangle (optionally with a border)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _accent_bar(slide, top=Inches(1.1)):
    """Yellow accent bar just below the section label."""
    _rect(slide, Inches(0.6), top, Inches(0.35), Pt(4), YELLOW)


def _header(slide, section_label, slide_num, total=12):
    """Standard slide header: logo text | section | page num."""
    # Logo text left
    _box(slide, Inches(0.55), Inches(0.28), Inches(2.5), Inches(0.4),
         "Dynamo AI", font_size=13, bold=True, colour=DARK)
    # Section label centre
    _box(slide, Inches(4), Inches(0.28), Inches(5.3), Inches(0.4),
         section_label.upper(), font_size=9, colour=MUTED, align=PP_ALIGN.CENTER)
    # Page number right
    _box(slide, Inches(11.5), Inches(0.28), Inches(1.3), Inches(0.4),
         f"{slide_num:02d} / {total}", font_size=9, colour=MUTED, align=PP_ALIGN.RIGHT)
    # Thin separator line
    _rect(slide, Inches(0.55), Inches(0.75), Inches(12.23), Pt(1), LINE)


def _footer(slide):
    """Standard footer line."""
    _rect(slide, Inches(0.55), Inches(7.1), Inches(12.23), Pt(1), LINE)
    _box(slide, Inches(0.55), Inches(7.15), Inches(6), Inches(0.28),
         "DYNAMO AI · INVESTOR PRESENTATION · 2026", font_size=7.5, colour=MUTED)
    _box(slide, Inches(7), Inches(7.15), Inches(5.78), Inches(0.28),
         "app.dynamoai.in", font_size=7.5, colour=GOLD, align=PP_ALIGN.RIGHT)


def _section_intro(slide, label, headline, sub=None, top=Inches(0.9)):
    """Yellow bar + label + h2 headline + optional sub-line."""
    _accent_bar(slide, top + Inches(0.05))
    _box(slide, Inches(0.55), top, Inches(12), Inches(0.3),
         label.upper(), font_size=9, colour=MUTED)
    _box(slide, Inches(0.55), top + Inches(0.3), Inches(12), Inches(0.55),
         headline, font_size=26, bold=True, colour=DARK)
    if sub:
        _box(slide, Inches(0.55), top + Inches(0.88), Inches(12), Inches(0.35),
             sub, font_size=13, colour=MUTED)


def _card(slide, left, top, width, height, title, body,
          bg=BGCARD, title_colour=DARK, body_colour=MUTED,
          accent_left=None):
    """White card with title + body text."""
    _rect(slide, left, top, width, height, bg, LINE)
    if accent_left:
        _rect(slide, left, top, Pt(3), height, accent_left)
    x_inner = left + Inches(0.18)
    w_inner = width - Inches(0.36)
    _box(slide, x_inner, top + Inches(0.12), w_inner, Inches(0.28),
         title, font_size=11, bold=True, colour=title_colour)
    _box(slide, x_inner, top + Inches(0.42), w_inner, height - Inches(0.52),
         body, font_size=10, colour=body_colour, wrap=True)


def _bullet_list(slide, left, top, width, items, colour=DARK, size=11):
    """Vertical list of bullet items (•)."""
    y = top
    for item in items:
        _box(slide, left, y, width, Inches(0.28),
             f"  {item}", font_size=size, colour=colour)
        y += Inches(0.3)


# ─────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────

def _slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide, RGBColor(0xFA, 0xFA, 0xF7))

    # Right yellow accent block
    _rect(slide, Inches(10.5), Inches(0), Inches(2.83), Inches(7.5),
          RGBColor(0xFF, 0xFB, 0xD0))

    # Logo badge
    _rect(slide, Inches(0.75), Inches(1.4), Inches(0.8), Inches(0.8), YELLOW)
    _box(slide, Inches(0.75), Inches(1.4), Inches(0.8), Inches(0.8),
         "D", font_size=30, bold=True, colour=DARK, align=PP_ALIGN.CENTER)

    # Main headline
    _box(slide, Inches(0.75), Inches(2.4), Inches(9.5), Inches(0.75),
         "India's First", font_size=46, bold=True, colour=DARK)
    _box(slide, Inches(0.75), Inches(3.15), Inches(9.5), Inches(0.75),
         "Research OS", font_size=46, bold=True, colour=GOLD)

    # Tagline
    _box(slide, Inches(0.75), Inches(4.05), Inches(9.3), Inches(0.65),
         "One platform for research discovery, academic writing, citations,\n"
         "document analysis, and presentations.",
         font_size=14, colour=MUTED, wrap=True)

    # Three stats
    stats = [
        ("₹399/mo", "Starting price"),
        ("3 Modes",  "Fast · Research · DeepThink"),
        ("Made in India", "Built for Indian researchers"),
    ]
    x = Inches(0.75)
    _rect(slide, x, Inches(5.05), Inches(9.3), Pt(1), LINE)
    for i, (val, lbl) in enumerate(stats):
        sx = x + i * Inches(3.1)
        if i > 0:
            _rect(slide, sx - Inches(0.02), Inches(5.2), Pt(1), Inches(0.65), LINE)
        _box(slide, sx, Inches(5.2), Inches(2.9), Inches(0.38),
             val, font_size=19, bold=True, colour=DARK)
        _box(slide, sx, Inches(5.6), Inches(2.9), Inches(0.28),
             lbl, font_size=9, colour=MUTED)

    # Confidential footer
    _box(slide, Inches(0.75), Inches(6.95), Inches(6), Inches(0.28),
         "CONFIDENTIAL · NOT FOR DISTRIBUTION", font_size=7.5, colour=MUTED)
    _box(slide, Inches(7.5), Inches(6.95), Inches(5.28), Inches(0.28),
         "app.dynamoai.in", font_size=7.5, colour=GOLD, align=PP_ALIGN.RIGHT)


def _slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Problem Statement", 2)
    _footer(slide)
    _section_intro(slide, "The Problem",
                   "Research today is fragmented.",
                   "A researcher uses 5–10 different tools just to complete one project.",
                   top=Inches(0.85))

    cards = [
        ("Too Many Separate Tools",
         "Search engine, AI chatbot, citation tool, PDF reader, presentation software — all separate, all costly."),
        ("Generic AI is Not Enough",
         "Most AI tools are built for general use. None are designed specifically for academic workflows."),
        ("Citation Errors Waste Time",
         "Creating and verifying citations manually is slow and error-prone, especially across 8+ formats."),
        ("No End-to-End Workflow",
         "Existing solutions solve only parts of the research process. No platform covers the full journey."),
    ]
    cw, ch = Inches(5.95), Inches(1.35)
    positions = [
        (Inches(0.55), Inches(2.1)),
        (Inches(6.75), Inches(2.1)),
        (Inches(0.55), Inches(3.55)),
        (Inches(6.75), Inches(3.55)),
    ]
    for (t, b), (lx, ly) in zip(cards, positions):
        _card(slide, lx, ly, cw, ch, t, b, accent_left=YELLOW)

    # Wide bottom card
    _card(slide, Inches(0.55), Inches(5.0), Inches(12.23), Inches(1.1),
          "Overlooked Indian Researchers",
          "Global platforms are not designed for Indian academic standards, languages, or pricing realities. "
          "India has over 1.5 million PhD scholars — a hugely underserved market.",
          accent_left=YELLOW)


def _slide_03_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Solution", 3)
    _footer(slide)
    _section_intro(slide, "The Solution",
                   "One platform. The complete research workflow.",
                   "Dynamo AI replaces 5+ separate tools with a single unified workspace.",
                   top=Inches(0.85))

    # Workflow arrow strip
    _rect(slide, Inches(0.55), Inches(2.15), Inches(12.23), Inches(1.4),
          RGBColor(0xFF, 0xFF, 0xFF), LINE)
    steps = [
        ("🔍", "Research",  "Discover papers"),
        ("📄", "Analyse",   "Read documents"),
        ("✏",  "Write",     "Draft academic content"),
        ("📋", "Cite",      "Generate citations"),
        ("🎨", "Present",   "Create slides"),
    ]
    sw = Inches(2.0)
    for i, (icon, title, sub) in enumerate(steps):
        sx = Inches(0.9) + i * Inches(2.42)
        _box(slide, sx, Inches(2.2), Inches(1.8), Inches(0.42),
             title, font_size=13, bold=True, colour=DARK, align=PP_ALIGN.CENTER)
        _box(slide, sx, Inches(2.65), Inches(1.8), Inches(0.28),
             sub, font_size=9, colour=MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            _box(slide, sx + Inches(1.85), Inches(2.3), Inches(0.5), Inches(0.3),
                 "→", font_size=16, colour=MUTED, align=PP_ALIGN.CENTER)

    # Three feature pillars
    pillars = [
        ("Research-First AI",
         "Designed specifically for academic and research workflows, not a general chatbot."),
        ("Three AI Modes",
         "Fast, Research, and DeepThink modes for different levels of task complexity."),
        ("Privacy Protected",
         "User data, research work, and documents remain protected with strong security practices."),
    ]
    pw = Inches(3.9)
    for i, (t, b) in enumerate(pillars):
        px = Inches(0.55) + i * Inches(4.1)
        _card(slide, px, Inches(3.75), pw, Inches(1.3), t, b)

    # Bottom note
    _rect(slide, Inches(0.55), Inches(5.2), Inches(12.23), Inches(0.55),
          RGBColor(0xFF, 0xFB, 0xEA))
    _box(slide, Inches(0.75), Inches(5.27), Inches(12), Inches(0.35),
         "Replace multiple subscriptions with one platform. Starting at ₹399/month.",
         font_size=11, colour=GOLD)


def _slide_04_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Business Model", 4)
    _footer(slide)
    _section_intro(slide, "Business Model Type",
                   "SaaS Subscription with a Two-Sided Value Chain",
                   top=Inches(0.85))

    # Left column — platform diagram
    _box(slide, Inches(0.55), Inches(1.95), Inches(5.9), Inches(0.28),
         "PLATFORM STRUCTURE", font_size=9, bold=True, colour=MUTED)
    layers = [
        (RGBColor(0xFF, 0xFB, 0xEA), GOLD, "Users",
         "PhD Scholars · Professors · Research Assistants"),
        (RGBColor(0xF5, 0xF5, 0xF0), MUTED, "subscribe to access", ""),
        (RGBColor(0x1A, 0x1A, 0x1A), YELLOW, "Dynamo AI Platform",
         "AI Research Operating System"),
        (RGBColor(0xF5, 0xF5, 0xF0), MUTED, "powered by", ""),
        (RGBColor(0xFF, 0xFF, 0xFF), DARK, "AI Infrastructure",
         "Gemini · Groq · Tavily · Razorpay"),
    ]
    ly = Inches(2.3)
    for i, (bg, tc, title, sub) in enumerate(layers):
        lh = Inches(0.45) if sub else Inches(0.32)
        _rect(slide, Inches(0.55), ly, Inches(5.9), lh, bg, LINE)
        _box(slide, Inches(0.75), ly + Inches(0.04), Inches(5.5), Inches(0.22),
             title, font_size=11, bold=bool(sub), colour=tc, align=PP_ALIGN.CENTER)
        if sub:
            _box(slide, Inches(0.75), ly + Inches(0.24), Inches(5.5), Inches(0.18),
                 sub, font_size=8, colour=RGBColor(0x9A, 0x7E, 0x30) if bg == RGBColor(0xFF, 0xFB, 0xEA)
                 else RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.CENTER)
        ly += lh

    # Right column — pricing tiers
    _box(slide, Inches(6.75), Inches(1.95), Inches(5.9), Inches(0.28),
         "MONETIZATION MODEL", font_size=9, bold=True, colour=MUTED)
    tiers = [
        (RGBColor(0xFF, 0xFF, 0xFF), DARK,  DARK,  "Free",  "10 messages/day · User acquisition", "₹0"),
        (RGBColor(0xFF, 0xFD, 0xF5), GOLD,  GOLD,  "Plus",  "100 messages/day · Core research tools", "₹399/mo"),
        (RGBColor(0x1A, 0x1A, 0x1A), YELLOW, RGBColor(0xFF,0xFF,0xFF), "Pro",
         "300 messages/day · Full platform access", "₹999/mo"),
    ]
    ty = Inches(2.3)
    for bg, tc, vc, name, desc, price in tiers:
        _rect(slide, Inches(6.75), ty, Inches(5.9), Inches(0.7), bg, LINE)
        _box(slide, Inches(6.95), ty + Inches(0.06), Inches(3), Inches(0.28),
             name, font_size=13, bold=True, colour=tc)
        _box(slide, Inches(6.95), ty + Inches(0.36), Inches(3.5), Inches(0.22),
             desc, font_size=9, colour=RGBColor(0x9A, 0x7E, 0x30) if bg == RGBColor(0xFF,0xFD,0xF5)
             else (RGBColor(0x80,0x80,0x80) if bg == RGBColor(0xFF,0xFF,0xFF)
                   else RGBColor(0x99,0x99,0x99)))
        _box(slide, Inches(10.5), ty + Inches(0.12), Inches(2), Inches(0.4),
             price, font_size=17, bold=True, colour=vc, align=PP_ALIGN.RIGHT)
        ty += Inches(0.78)

    # Bottom note
    _rect(slide, Inches(6.75), ty + Inches(0.1), Inches(5.9), Inches(0.55),
          RGBColor(0xF5, 0xF5, 0xF0))
    _box(slide, Inches(6.95), ty + Inches(0.17), Inches(5.5), Inches(0.35),
         "Scalable SaaS revenue grows with every new subscriber. Low marginal cost per additional user.",
         font_size=10, colour=MUTED, wrap=True)


def _slide_05_customer_segment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Customer Segment", 5)
    _footer(slide)
    _section_intro(slide, "Customer Segment",
                   "Who We Serve",
                   "India has over 1.5 million active PhD scholars and thousands of research institutions.",
                   top=Inches(0.85))

    segments = [
        ("PhD Scholars", [
            "Literature reviews", "Thesis writing",
            "Citation management", "Research analysis"]),
        ("Professors", [
            "Research publications", "Teaching materials",
            "Academic writing", "Content creation"]),
        ("Research Assistants", [
            "Document analysis", "Summarisation",
            "Workflow efficiency", "Report writing"]),
        ("Institutions", [
            "Department productivity", "Standardised tools",
            "Cost consolidation", "Bulk licensing"]),
    ]
    cw = Inches(2.95)
    for i, (title, bullets) in enumerate(segments):
        cx = Inches(0.55) + i * Inches(3.13)
        _rect(slide, cx, Inches(2.15), cw, Inches(3.5), BGCARD, LINE)
        _rect(slide, cx, Inches(2.15), cw, Pt(3), YELLOW)
        _box(slide, cx + Inches(0.15), Inches(2.22), cw - Inches(0.3), Inches(0.32),
             title, font_size=12, bold=True, colour=DARK)
        _bullet_list(slide, cx + Inches(0.15), Inches(2.62), cw - Inches(0.3),
                     bullets, colour=MUTED, size=10)

    # Bottom note
    _rect(slide, Inches(0.55), Inches(5.82), Inches(12.23), Inches(0.55),
          RGBColor(0xFF, 0xFB, 0xEA))
    _box(slide, Inches(0.75), Inches(5.89), Inches(12), Inches(0.35),
         "Core need across all segments: save time, reduce costs, and improve research quality — all in one place.",
         font_size=11, colour=GOLD, bold=True, wrap=True)


def _slide_06_service_offering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Service Offering", 6)
    _footer(slide)
    _section_intro(slide, "Service Offering",
                   "What Dynamo AI Offers",
                   top=Inches(0.85))

    features = [
        ("Research Discovery",
         "Web search, Research Mode, and DeepThink AI to find and synthesise relevant academic sources."),
        ("Academic Writing",
         "Write a Paper tool generates full academic papers by topic, with structured sections and references."),
        ("Citation Checker",
         "Verifies in-text citations and reference lists across APA, IEEE, MLA, Harvard, Vancouver, Chicago."),
        ("AI Text Detector",
         "In-house AI detection and plagiarism checker — no external API dependency, powered by Gemini."),
        ("Document Intelligence",
         "Upload PDFs, DOCX, and TXT files. Dynamo AI summarises and stores them in a personal library."),
        ("Study Tools",
         "Study guides, quiz mode, flashcards, mindmaps, flowcharts, and AI-generated image creation."),
    ]
    cw, ch = Inches(3.9), Inches(1.4)
    for i, (t, b) in enumerate(features):
        row = i // 3
        col = i % 3
        cx = Inches(0.55) + col * Inches(4.1)
        cy = Inches(1.85) + row * Inches(1.55)
        _card(slide, cx, cy, cw, ch, t, b)

    # Tags row
    tags = ["Fast Mode", "Research Mode", "DeepThink Mode", "Voice + TTS", "PDF Export", "AI Memory", "Multi-language"]
    tx = Inches(0.55)
    for tag in tags:
        _rect(slide, tx, Inches(5.02), Inches(len(tag) * 0.1 + 0.6), Inches(0.32),
              RGBColor(0xF5, 0xF5, 0xF0), LINE)
        _box(slide, tx + Inches(0.1), Inches(5.05), Inches(len(tag) * 0.1 + 0.4), Inches(0.25),
             tag, font_size=9, colour=MUTED)
        tx += Inches(len(tag) * 0.1 + 0.75)


def _slide_07_revenue_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Revenue Model", 7)
    _footer(slide)
    _section_intro(slide, "Revenue Model",
                   "Subscription-Led, Scalable Revenue",
                   top=Inches(0.85))

    # Left: subscription tiers
    _box(slide, Inches(0.55), Inches(1.85), Inches(6.2), Inches(0.25),
         "SUBSCRIPTION TIERS", font_size=9, bold=True, colour=MUTED)
    tiers = [
        (BGCARD, DARK, MUTED, DARK, "Free",
         "10 messages/day · 0 images · Entry point", "₹0", "User acquisition"),
        (RGBColor(0xFF, 0xFD, 0xF5), GOLD, RGBColor(0x9A, 0x7E, 0x30), GOLD, "Plus",
         "100 messages/day · 25 images · Full tools", "₹399/mo", "Core revenue driver"),
        (INVBG, YELLOW, RGBColor(0x99, 0x99, 0x99), RGBColor(0xFF, 0xFF, 0xFF), "Pro",
         "300 messages/day · 100 images · All features", "₹999/mo", "High-value segment"),
    ]
    ty = Inches(2.15)
    for bg, tc, dc, vc, name, desc, price, note in tiers:
        _rect(slide, Inches(0.55), ty, Inches(6.2), Inches(0.82), bg, LINE)
        _box(slide, Inches(0.75), ty + Inches(0.07), Inches(3.5), Inches(0.28),
             name, font_size=14, bold=True, colour=tc)
        _box(slide, Inches(0.75), ty + Inches(0.42), Inches(3.5), Inches(0.26),
             desc, font_size=9, colour=dc)
        _box(slide, Inches(5), ty + Inches(0.1), Inches(1.55), Inches(0.4),
             price, font_size=17, bold=True, colour=vc, align=PP_ALIGN.RIGHT)
        _box(slide, Inches(5), ty + Inches(0.52), Inches(1.55), Inches(0.22),
             note, font_size=8, colour=dc, align=PP_ALIGN.RIGHT)
        ty += Inches(0.9)

    # Right: revenue streams
    _box(slide, Inches(6.95), Inches(1.85), Inches(5.85), Inches(0.25),
         "REVENUE STREAMS", font_size=9, bold=True, colour=MUTED)
    streams = [
        ("Monthly Subscriptions",
         "Recurring SaaS revenue from Plus and Pro users. Predictable, compounding growth."),
        ("Institutional Licences",
         "Bulk access for universities, research departments, and academic institutions."),
        ("Future: API Access",
         "Developer and enterprise API access planned as platform matures."),
    ]
    ry = Inches(2.15)
    for t, b in streams:
        _card(slide, Inches(6.95), ry, Inches(5.85), Inches(0.88), t, b)
        ry += Inches(0.95)

    _rect(slide, Inches(6.95), ry + Inches(0.05), Inches(5.85), Inches(0.55),
          RGBColor(0xF5, 0xF5, 0xF0))
    _box(slide, Inches(7.15), ry + Inches(0.12), Inches(5.5), Inches(0.35),
         "Low marginal cost per new subscriber — each additional user adds revenue with minimal incremental cost.",
         font_size=10, colour=MUTED, wrap=True)


def _slide_08_value_prop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Value Proposition", 8)
    _footer(slide)
    _section_intro(slide, "Value Proposition",
                   "Why Dynamo AI?",
                   top=Inches(0.85))

    props = [
        ("India's First Research OS",
         "No other platform combines research discovery, writing, citations, and document analysis "
         "into one unified workspace built for Indian researchers."),
        ("Research-First, Not Generic",
         "Built from the ground up for academic workflows. Every feature — from DeepThink mode to "
         "citation verification — is designed specifically for researchers."),
        ("Affordable for Indian Users",
         "Full research platform access from ₹399/month — a fraction of the cost of maintaining "
         "5+ separate subscriptions globally priced."),
        ("Multi-Language Support",
         "Supports Indian languages, making advanced research tools accessible to a wider audience "
         "beyond English-only platforms."),
    ]
    cw, ch = Inches(5.95), Inches(1.45)
    positions = [
        (Inches(0.55), Inches(1.9)),
        (Inches(6.75), Inches(1.9)),
        (Inches(0.55), Inches(3.45)),
        (Inches(6.75), Inches(3.45)),
    ]
    for (t, b), (px, py) in zip(props, positions):
        _card(slide, px, py, cw, ch, t, b, accent_left=YELLOW)

    # Dark summary bar
    _rect(slide, Inches(0.55), Inches(5.05), Inches(12.23), Inches(0.85), INVBG)
    _box(slide, Inches(0.75), Inches(5.15), Inches(4), Inches(0.3),
         "Replace multiple subscriptions with one platform", font_size=11,
         colour=RGBColor(0xCC, 0xCC, 0xCC))
    stats = [("5+", "Tools replaced"), ("₹399", "Starting price"), ("1", "Platform")]
    for i, (val, lbl) in enumerate(stats):
        sx = Inches(5.5) + i * Inches(2.3)
        _box(slide, sx, Inches(5.12), Inches(2), Inches(0.35),
             val, font_size=19, bold=True, colour=YELLOW, align=PP_ALIGN.CENTER)
        _box(slide, sx, Inches(5.5), Inches(2), Inches(0.25),
             lbl, font_size=8, colour=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.CENTER)


def _slide_09_growth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Growth Strategy", 9)
    _footer(slide)
    _section_intro(slide, "Growth Strategy",
                   "How We Plan to Grow",
                   top=Inches(0.85))

    pillars = [
        ("Geographic Expansion", [
            "1. Start: Indian universities and research institutions",
            "2. Expand to South Asia and Southeast Asian markets",
            "3. Long-term: global researcher communities",
        ]),
        ("Product Enhancement", [
            "+ Multi-agent AI research assistants",
            "+ Collaborative research workspaces",
            "+ Enhanced document intelligence",
            "+ Deeper regional language support",
        ]),
        ("Institutional Partnerships", [
            "+ University-level bulk licence agreements",
            "+ Research department integrations",
            "+ Corporate R&D teams",
            "+ Government research bodies",
        ]),
    ]
    pw = Inches(3.9)
    for i, (title, bullets) in enumerate(pillars):
        px = Inches(0.55) + i * Inches(4.1)
        _rect(slide, px, Inches(1.9), pw, Inches(3.45), BGCARD, LINE)
        _box(slide, px + Inches(0.18), Inches(1.97), pw - Inches(0.36), Inches(0.32),
             title, font_size=13, bold=True, colour=DARK)
        _bullet_list(slide, px + Inches(0.18), Inches(2.38), pw - Inches(0.36),
                     bullets, colour=MUTED, size=10)

    # Network effect note
    _rect(slide, Inches(0.55), Inches(5.5), Inches(12.23), Inches(0.6),
          RGBColor(0xFF, 0xFB, 0xEA))
    _box(slide, Inches(0.75), Inches(5.57), Inches(12), Inches(0.4),
         "Network effect: More researchers → more institutions → more researchers. A compounding growth loop.",
         font_size=11, colour=GOLD, wrap=True)


def _slide_10_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Risks & Challenges", 10)
    _footer(slide)
    _section_intro(slide, "Risks and Challenges",
                   "Known Risks and How We Address Them",
                   "Transparency on risks is part of responsible platform development.",
                   top=Inches(0.85))

    risks = [
        ("Competition",
         "Global AI tools (ChatGPT, Perplexity, Claude) serve general audiences. "
         "Mitigation: deep research-specific features and Indian market pricing that global tools cannot match."),
        ("Data Privacy",
         "Researchers share sensitive academic work. Mitigation: strong privacy practices, no training on user data, "
         "Supabase-backed secure storage with access controls."),
        ("Quality Control",
         "AI outputs must meet academic standards. Mitigation: multi-model verification, "
         "temperature-controlled generation, and structured prompting for consistent results."),
        ("Scaling Challenges",
         "Rapid user growth increases AI compute costs. Mitigation: quota-based usage tiers, "
         "fallback model routing, and usage-linked subscription revenue to balance costs."),
    ]
    cw, ch = Inches(5.95), Inches(1.35)
    positions = [
        (Inches(0.55), Inches(2.15)),
        (Inches(6.75), Inches(2.15)),
        (Inches(0.55), Inches(3.6)),
        (Inches(6.75), Inches(3.6)),
    ]
    for (t, b), (px, py) in zip(risks, positions):
        _card(slide, px, py, cw, ch, t, b)

    _rect(slide, Inches(0.55), Inches(5.1), Inches(12.23), Inches(0.55),
          RGBColor(0xF5, 0xF5, 0xF0))
    _box(slide, Inches(0.75), Inches(5.17), Inches(12), Inches(0.35),
         "Trust and Safety: Citation verification, AI detection, and plagiarism checking reinforce academic integrity.",
         font_size=11, colour=MUTED, wrap=True)


def _slide_11_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Future Opportunities", 11)
    _footer(slide)
    _section_intro(slide, "Future Opportunities",
                   "Where Dynamo AI Is Headed",
                   top=Inches(0.85))

    opps = [
        ("Multi-Agent AI Research Assistants",
         "Autonomous agents that can conduct end-to-end literature reviews, synthesise findings, and generate reports independently."),
        ("Institutional Research Solutions",
         "Dedicated dashboards and admin tools for universities, research labs, and government departments."),
        ("International Expansion",
         "Localised versions for South Asia, Southeast Asia, Africa, and other emerging research markets."),
        ("Smarter Publishing Workflows",
         "Guided journal submission workflows: formatting, compliance checking, and co-authorship management."),
        ("Data-Driven Platform Intelligence",
         "Aggregated anonymised usage to identify research trends, improve recommendations, and personalise the experience."),
        ("Monetisation Expansion",
         "Annual billing discounts, API access tiers, white-label licences, and premium corporate research packages."),
    ]
    cw, ch = Inches(5.95), Inches(1.3)
    for i, (t, b) in enumerate(opps):
        row = i // 2
        col = i % 2
        px = Inches(0.55) + col * Inches(6.45)
        py = Inches(1.9) + row * Inches(1.42)
        _card(slide, px, py, cw, ch, t, b)


def _slide_12_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, RGBColor(0xFA, 0xFA, 0xF7))

    # Right accent
    _rect(slide, Inches(10.5), Inches(0), Inches(2.83), Inches(7.5),
          RGBColor(0xFF, 0xFB, 0xD0))

    # Top bar
    _box(slide, Inches(0.75), Inches(0.3), Inches(2), Inches(0.35),
         "Dynamo AI", font_size=13, bold=True, colour=DARK)
    _box(slide, Inches(10.5), Inches(0.3), Inches(2.3), Inches(0.35),
         "SUMMARY", font_size=9, colour=MUTED, align=PP_ALIGN.RIGHT)

    # Logo badge
    _rect(slide, Inches(0.75), Inches(1.3), Inches(0.7), Inches(0.7), YELLOW)
    _box(slide, Inches(0.75), Inches(1.3), Inches(0.7), Inches(0.7),
         "D", font_size=26, bold=True, colour=DARK, align=PP_ALIGN.CENTER)

    # Headline
    _box(slide, Inches(0.75), Inches(2.1), Inches(9.5), Inches(0.65),
         "Dynamo AI — India's", font_size=36, bold=True, colour=DARK)
    _box(slide, Inches(0.75), Inches(2.75), Inches(9.5), Inches(0.65),
         "Research Operating System", font_size=36, bold=True, colour=GOLD)

    # Tagline
    _box(slide, Inches(0.75), Inches(3.55), Inches(9.3), Inches(0.65),
         "One platform that replaces the fragmented research toolchain for India's\n"
         "1.5 million PhD scholars, professors, and research professionals.",
         font_size=13, colour=MUTED, wrap=True)

    # Three stat boxes
    boxes_data = [
        ("SaaS Subscription", "Free · Plus · Pro", "Model"),
        ("From ₹399/mo",      "Affordable for Indian users", "Price"),
        ("India's Best",      "Research AI platform", "Vision"),
    ]
    bw, bh = Inches(3.7), Inches(0.85)
    for i, (val, sub, lbl) in enumerate(boxes_data):
        bx = Inches(0.75) + i * Inches(3.9)
        bg = RGBColor(0xFF, 0xFB, 0xEA) if i == 2 else BGCARD
        bc = GOLD if i == 2 else DARK
        _rect(slide, bx, Inches(4.45), bw, bh, bg, LINE)
        _box(slide, bx + Inches(0.15), Inches(4.5), Inches(1), Inches(0.22),
             lbl.upper(), font_size=8, bold=True,
             colour=RGBColor(0x9A, 0x7E, 0x30) if i == 2 else MUTED)
        _box(slide, bx + Inches(0.15), Inches(4.74), bw - Inches(0.3), Inches(0.3),
             val, font_size=14, bold=True, colour=bc)
        _box(slide, bx + Inches(0.15), Inches(5.06), bw - Inches(0.3), Inches(0.2),
             sub, font_size=9, colour=RGBColor(0x9A, 0x7E, 0x30) if i == 2 else MUTED)

    # Footer
    _rect(slide, Inches(0.55), Inches(6.8), Inches(12.23), Pt(1), LINE)
    _box(slide, Inches(0.55), Inches(6.85), Inches(6), Inches(0.28),
         "CONFIDENTIAL · NOT FOR DISTRIBUTION", font_size=7.5, colour=MUTED)
    _box(slide, Inches(7), Inches(6.85), Inches(5.78), Inches(0.28),
         "app.dynamoai.in", font_size=7.5, colour=GOLD, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────
# MAIN EXPORT FUNCTION
# ─────────────────────────────────────────────────

def build_investor_pptx() -> bytes:
    """Build and return the PPTX bytes."""
    prs = Presentation()
    prs.slide_width  = Emu(12192000)   # 13.333" = 12192000 EMU  (widescreen 16:9)
    prs.slide_height = Emu(6858000)    # 7.5"    = 6858000 EMU

    _slide_01_cover(prs)
    _slide_02_problem(prs)
    _slide_03_solution(prs)
    _slide_04_business_model(prs)
    _slide_05_customer_segment(prs)
    _slide_06_service_offering(prs)
    _slide_07_revenue_model(prs)
    _slide_08_value_prop(prs)
    _slide_09_growth(prs)
    _slide_10_risks(prs)
    _slide_11_future(prs)
    _slide_12_summary(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    out_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "dynamo_ai_investor_deck.pptx")
    data = build_investor_pptx()
    with open(out_path, "wb") as f:
        f.write(data)
    print(f"Saved {len(data):,} bytes → {out_path}")
