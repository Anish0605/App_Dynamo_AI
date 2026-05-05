# presentation_engine.py — Dynamo AI
# Builds real PPT from AI JSON schema.
# Supports both legacy "content/chart" types and new smart slide types.

import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from fastapi.responses import StreamingResponse

# --------------------------------------------------
# COLOUR THEMES
# --------------------------------------------------

COLOR_THEMES = {
    "academic": {
        "bg":             (0x1e, 0x29, 0x3b),   # slate-800
        "accent":         (0xfb, 0xbf, 0x24),   # yellow-400
        "text_primary":   (0xff, 0xff, 0xff),
        "text_secondary": (0x94, 0xa3, 0xb8),   # slate-400
        "card_bg":        (0x0f, 0x17, 0x2a),   # darker slate
        "title_size": 36, "body_size": 18,
    },
    "business": {
        "bg":             (0xff, 0xff, 0xff),
        "accent":         (0x1e, 0x40, 0xaf),   # blue-800
        "text_primary":   (0x1e, 0x29, 0x3b),
        "text_secondary": (0x64, 0x74, 0x8b),
        "card_bg":        (0xf1, 0xf5, 0xf9),
        "title_size": 36, "body_size": 18,
    },
    "pitch": {
        "bg":             (0x09, 0x09, 0x0b),
        "accent":         (0xa7, 0x8b, 0xfa),   # violet-400
        "text_primary":   (0xff, 0xff, 0xff),
        "text_secondary": (0x71, 0x71, 0x7a),
        "card_bg":        (0x18, 0x18, 0x1b),
        "title_size": 36, "body_size": 18,
    },
    "minimal": {
        "bg":             (0xff, 0xff, 0xff),
        "accent":         (0x37, 0x41, 0x51),   # gray-700
        "text_primary":   (0x11, 0x18, 0x27),
        "text_secondary": (0x9c, 0xa3, 0xaf),
        "card_bg":        (0xf9, 0xfa, 0xfb),
        "title_size": 36, "body_size": 18,
    },
}

# Legacy theme compatibility
_LEGACY_THEMES = {
    "light":     {"bg": None, "title_size": 40, "body_size": 20},
    "dark":      {"bg": None, "title_size": 40, "body_size": 20},
    "executive": {"bg": None, "title_size": 44, "body_size": 22},
}

# Slide dimensions (widescreen 10×7.5 in)
W = 10.0
H = 7.5

# --------------------------------------------------
# HELPERS
# --------------------------------------------------

def _rgb(t):
    return RGBColor(t[0], t[1], t[2])

def _fill_bg(slide, color_tuple):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_tuple)

def _add_tb(slide, left, top, width, height, text,
            size, color, bold=False, italic=False,
            align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box and return it."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return tb

def _add_bullets_tb(slide, left, top, width, height,
                    bullets, size, color, marker_color,
                    indent=0.25):
    """Add a text box with bullet points."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        # Bullet marker
        mr = p.add_run()
        mr.text = "▸  "
        mr.font.size  = Pt(size - 2)
        mr.font.color.rgb = _rgb(marker_color)
        mr.font.bold  = True
        # Text
        r = p.add_run()
        r.text = bullet
        r.font.size  = Pt(size)
        r.font.color.rgb = _rgb(color)
    return tb

def _accent_bar(slide, left, top, width, height, color):
    """Draw a solid rectangle accent bar."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape

def _slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank


# --------------------------------------------------
# SMART SLIDE RENDERERS
# --------------------------------------------------

def _render_title(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])

    # Accent bar top
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    # Centred title
    _add_tb(slide, 0.8, 1.8, 8.4, 2.0,
            s.get("heading", ""),
            th["title_size"], th["text_primary"],
            bold=True, align=PP_ALIGN.CENTER)

    # Thin divider
    _accent_bar(slide, 3.5, 4.0, 3.0, 0.04, th["accent"])

    # Subheading
    sub = s.get("subheading", "")
    if sub:
        _add_tb(slide, 0.8, 4.2, 8.4, 0.8,
                sub, 18, th["text_secondary"],
                align=PP_ALIGN.CENTER)

    # Footer note
    _add_tb(slide, 0.8, 6.8, 8.4, 0.4,
            "Generated by Dynamo AI", 10,
            th["text_secondary"],
            align=PP_ALIGN.CENTER)


def _render_thesis(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    # Label
    _add_tb(slide, 0.6, 0.6, 4.0, 0.4,
            "RESEARCH QUESTION", 9,
            th["accent"], bold=True)

    # Heading
    _add_tb(slide, 0.6, 1.1, 8.8, 0.8,
            s.get("heading", ""), th["title_size"] - 4,
            th["text_primary"], bold=True)

    # Thesis callout box (simulated with accent bar + text)
    _accent_bar(slide, 0.6, 2.2, 0.06, 2.8, th["accent"])
    _add_tb(slide, 0.9, 2.2, 8.5, 3.0,
            s.get("thesis", ""), th["body_size"] + 2,
            th["text_primary"], italic=True)


def _render_background(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    _add_tb(slide, 0.6, 0.4, 8.8, 0.7,
            s.get("heading", ""), th["title_size"] - 4,
            th["text_primary"], bold=True)
    _accent_bar(slide, 0.6, 1.2, 1.2, 0.04, th["accent"])

    _add_bullets_tb(slide, 0.6, 1.5, 8.8, 5.0,
                    s.get("bullets", []),
                    th["body_size"], th["text_primary"],
                    th["accent"])


def _render_evidence(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    _add_tb(slide, 0.6, 0.4, 8.8, 0.7,
            s.get("heading", ""), th["title_size"] - 4,
            th["text_primary"], bold=True)
    _accent_bar(slide, 0.6, 1.2, 1.2, 0.04, th["accent"])

    _add_bullets_tb(slide, 0.6, 1.5, 8.8, 4.5,
                    s.get("bullets", []),
                    th["body_size"], th["text_primary"],
                    th["accent"])

    citation = s.get("citation", "")
    if citation:
        # Citation footer line
        _accent_bar(slide, 0.6, 6.5, 8.8, 0.02, th["text_secondary"])
        _add_tb(slide, 0.6, 6.6, 8.8, 0.5,
                f"📚  {citation}", 9,
                th["text_secondary"], italic=True)


def _render_chart(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    _add_tb(slide, 0.6, 0.4, 8.8, 0.7,
            s.get("heading", ""), th["title_size"] - 4,
            th["text_primary"], bold=True)

    chart_info = s.get("chart", {})
    labels = chart_info.get("labels", ["A", "B", "C", "D"])
    values = chart_info.get("values", [40, 60, 50, 75])

    chart_data = ChartData()
    chart_data.categories = labels
    chart_data.add_series("Value", values)

    try:
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8), Inches(1.4),
            Inches(8.4), Inches(4.8),
            chart_data
        ).chart
        # Style the chart
        chart.has_legend = False
        plot = chart.plots[0]
        for series in plot.series:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(th["accent"])
    except Exception:
        _add_tb(slide, 0.6, 2.0, 8.8, 1.0,
                f"Chart: {labels}", th["body_size"],
                th["text_secondary"])


def _render_comparison(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    _add_tb(slide, 0.6, 0.4, 8.8, 0.7,
            s.get("heading", ""), th["title_size"] - 4,
            th["text_primary"], bold=True)
    _accent_bar(slide, 0.6, 1.2, 1.2, 0.04, th["accent"])

    left_col  = s.get("left",  {})
    right_col = s.get("right", {})

    # Divider line between columns
    _accent_bar(slide, 4.95, 1.3, 0.04, 5.8, th["text_secondary"])

    for col, x_off in [(left_col, 0.6), (right_col, 5.2)]:
        _add_tb(slide, x_off, 1.4, 4.1, 0.5,
                col.get("label", ""), 15,
                th["accent"], bold=True)
        _add_bullets_tb(slide, x_off, 2.1, 4.1, 4.5,
                        col.get("points", []),
                        th["body_size"] - 2,
                        th["text_primary"],
                        th["accent"])


def _render_quote(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    _add_tb(slide, 0.6, 0.4, 8.8, 0.6,
            s.get("heading", "Expert Perspective"), 16,
            th["text_secondary"], bold=True)

    # Opening quote mark
    _add_tb(slide, 0.4, 0.9, 1.0, 1.2,
            "\u201c", 72, th["accent"], bold=True)

    _add_tb(slide, 0.8, 1.6, 8.4, 3.0,
            s.get("quote", ""), th["body_size"] + 4,
            th["text_primary"], italic=True,
            align=PP_ALIGN.CENTER)

    citation = s.get("citation", "")
    if citation:
        _add_tb(slide, 0.8, 5.2, 8.4, 0.6,
                f"— {citation}", 12,
                th["text_secondary"],
                align=PP_ALIGN.CENTER)


def _render_conclusion(prs, s, th):
    slide = _slide_blank(prs)
    _fill_bg(slide, th["bg"])
    _accent_bar(slide, 0, 0, W, 0.06, th["accent"])

    _add_tb(slide, 0.6, 0.4, 8.8, 0.7,
            s.get("heading", ""), th["title_size"] - 4,
            th["text_primary"], bold=True)
    _accent_bar(slide, 0.6, 1.2, 1.2, 0.04, th["accent"])

    bullets = s.get("bullets", [])
    top = 1.5
    for i, bullet in enumerate(bullets, 1):
        # Numbered badge
        badge = slide.shapes.add_shape(
            1,
            Inches(0.6), Inches(top),
            Inches(0.38), Inches(0.38)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb(th["accent"])
        badge.line.fill.background()
        _add_tb(slide, 0.6, top - 0.03, 0.38, 0.44,
                str(i), 12, (0, 0, 0),
                bold=True, align=PP_ALIGN.CENTER)
        _add_tb(slide, 1.15, top + 0.02, 8.2, 0.6,
                bullet, th["body_size"],
                th["text_primary"])
        top += 0.9


# Slide type dispatcher
_RENDERERS = {
    "title":      _render_title,
    "thesis":     _render_thesis,
    "background": _render_background,
    "evidence":   _render_evidence,
    "chart":      _render_chart,
    "comparison": _render_comparison,
    "quote":      _render_quote,
    "conclusion": _render_conclusion,
}


# --------------------------------------------------
# SMART BUILD (new pipeline)
# --------------------------------------------------

def build_smart_presentation(payload: dict):
    """
    Renders a fully-themed PPTX from a deck outline dict.
    payload = { title, style, slides: [ {type, ...} ] }
    """
    prs   = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    style = payload.get("style", "academic")
    th    = COLOR_THEMES.get(style, COLOR_THEMES["academic"])

    for s in payload.get("slides", []):
        stype    = s.get("type", "background")
        renderer = _RENDERERS.get(stype, _render_background)
        renderer(prs, s, th)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="DynamoAI_Deck.pptx"'}
    )


# --------------------------------------------------
# LEGACY BUILD (kept for /generate-ppt-smart)
# --------------------------------------------------

THEMES = {
    "light":     {"bg": None, "title_size": 40, "body_size": 20},
    "dark":      {"bg": None, "title_size": 40, "body_size": 20},
    "executive": {"bg": None, "title_size": 44, "body_size": 22},
}


def build_presentation(payload: dict):
    """
    Legacy builder — title + content/chart slides only.
    Preserved for /generate-ppt-smart backwards compatibility.
    """
    prs   = Presentation()
    theme = payload.get("theme", "light")
    style = THEMES.get(theme, THEMES["light"])

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = payload.get("title", "Dynamo AI Presentation")
    slide.placeholders[1].text = "Generated by Dynamo AI"

    for s in payload.get("slides", []):
        if s.get("type") == "content":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = s.get("heading", "")
            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in s.get("bullets", []):
                p = body.add_paragraph()
                p.text = bullet
                p.font.size = Pt(style["body_size"])
                p.level = 1

        elif s.get("type") == "chart":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = s.get("heading", "")
            chart_info = s.get("chart", {})
            labels = chart_info.get("labels", [])
            values = chart_info.get("values", [])
            chart_data = ChartData()
            chart_data.categories = labels
            chart_data.add_series("Series", values)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(1.5),
                Inches(8), Inches(4),
                chart_data
            )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="DynamoAI_Presentation.pptx"'}
    )
