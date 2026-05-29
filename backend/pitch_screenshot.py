"""
pitch_screenshot.py — Screenshots all 12 investor-deck slides using Playwright,
then packages them as:
  1. A PDF  (pixel-perfect, for sharing / presenting)
  2. A PPTX (image-per-slide, editable in Google Slides / PowerPoint)
"""

import asyncio
import io
import os
from pathlib import Path
from typing import List

# ─── SLIDE ROUTES ────────────────────────────────────────────────────────────
SLIDE_ROUTES = [
    "investor-deck/Inv01Cover",
    "investor-deck/Inv02Problem",
    "investor-deck/Inv03Solution",
    "investor-deck/Inv04BusinessModel",
    "investor-deck/Inv05CustomerSegment",
    "investor-deck/Inv06ServiceOffering",
    "investor-deck/Inv07RevenueModel",
    "investor-deck/Inv08ValueProp",
    "investor-deck/Inv09GrowthStrategy",
    "investor-deck/Inv10Risks",
    "investor-deck/Inv11FutureOpportunities",
    "investor-deck/Inv12Summary",
]

SLIDE_W = 1280
SLIDE_H = 720
PITCH_BASE = "http://localhost:21548/__pitch/preview"


# ─── SCREENSHOT ──────────────────────────────────────────────────────────────

async def _screenshot_slide(page, route: str, base_url: str) -> bytes:
    url = f"{base_url}/{route}"
    await page.goto(url, wait_until="networkidle", timeout=30000)
    # Give fonts / animations a moment to settle
    await asyncio.sleep(0.6)
    return await page.screenshot(type="png", clip={"x": 0, "y": 0, "width": SLIDE_W, "height": SLIDE_H})


async def capture_slides(base_url: str = PITCH_BASE) -> List[bytes]:
    from playwright.async_api import async_playwright
    images: List[bytes] = []
    import shutil
    chromium_path = (
        shutil.which("chromium")
        or "/nix/store/qa9cnw4v5xkxyip6mb9kxqfq1z4x2dx1-chromium-138.0.7204.100/bin/chromium"
    )
    async with async_playwright() as p:
        browser = await p.chromium.launch(
            executable_path=chromium_path,
            args=["--no-sandbox", "--disable-setuid-sandbox", "--disable-dev-shm-usage",
                  "--disable-gpu", "--single-process"],
        )
        context = await browser.new_context(
            viewport={"width": SLIDE_W, "height": SLIDE_H},
            device_scale_factor=1,
        )
        page = await context.new_page()
        # Suppress console noise
        page.on("console", lambda _: None)

        for route in SLIDE_ROUTES:
            try:
                png = await _screenshot_slide(page, route, base_url)
                images.append(png)
            except Exception as e:
                # Fallback: grey placeholder
                from PIL import Image, ImageDraw
                img = Image.new("RGB", (SLIDE_W, SLIDE_H), (250, 250, 247))
                d = ImageDraw.Draw(img)
                d.text((40, 40), f"Slide: {route}\n{e}", fill=(100, 100, 100))
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                images.append(buf.getvalue())

        await browser.close()
    return images


# ─── PDF ─────────────────────────────────────────────────────────────────────

def images_to_pdf(images: List[bytes]) -> bytes:
    from PIL import Image
    pil_images = []
    for png in images:
        img = Image.open(io.BytesIO(png)).convert("RGB")
        pil_images.append(img)

    buf = io.BytesIO()
    pil_images[0].save(
        buf,
        format="PDF",
        save_all=True,
        append_images=pil_images[1:],
        resolution=150,
    )
    return buf.getvalue()


# ─── PPTX (image-per-slide) ──────────────────────────────────────────────────

def images_to_pptx(images: List[bytes]) -> bytes:
    from pptx import Presentation
    from pptx.util import Emu, Inches

    prs = Presentation()
    prs.slide_width  = Emu(12192000)   # 16:9 widescreen width
    prs.slide_height = Emu(6858000)    # 16:9 widescreen height

    blank_layout = prs.slide_layouts[6]

    for png in images:
        slide = prs.slides.add_slide(blank_layout)
        img_stream = io.BytesIO(png)
        # Fill entire slide with the screenshot
        slide.shapes.add_picture(
            img_stream,
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── PUBLIC API ──────────────────────────────────────────────────────────────

async def build_pdf(base_url: str = PITCH_BASE) -> bytes:
    images = await capture_slides(base_url)
    return images_to_pdf(images)


async def build_image_pptx(base_url: str = PITCH_BASE) -> bytes:
    images = await capture_slides(base_url)
    return images_to_pptx(images)


# ─── CLI ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    mode = sys.argv[1] if len(sys.argv) > 1 else "pdf"
    base = sys.argv[2] if len(sys.argv) > 2 else PITCH_BASE

    async def main():
        print(f"Capturing {len(SLIDE_ROUTES)} slides from {base} …")
        images = await capture_slides(base)
        print(f"  → got {len(images)} screenshots")

        out_dir = Path(__file__).parent.parent / "frontend"

        if mode in ("pdf", "both"):
            pdf_data = images_to_pdf(images)
            out = out_dir / "dynamo_ai_investor_deck.pdf"
            out.write_bytes(pdf_data)
            print(f"  PDF  → {out}  ({len(pdf_data):,} bytes)")

        if mode in ("pptx", "both"):
            pptx_data = images_to_pptx(images)
            out = out_dir / "dynamo_ai_investor_deck_visual.pptx"
            out.write_bytes(pptx_data)
            print(f"  PPTX → {out}  ({len(pptx_data):,} bytes)")

    asyncio.run(main())
