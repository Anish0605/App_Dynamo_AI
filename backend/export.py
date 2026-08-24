# export.py — Dynamo AI (FINAL, SAFE)

import io
import re
import html
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
from fastapi.responses import StreamingResponse

# --------------------------------------------------
# HISTORY NORMALIZER
# --------------------------------------------------

def normalize_history(history):
    clean = []

    if not isinstance(history, list):
        return clean

    for m in history:
        if not isinstance(m, dict):
            continue

        role = m.get("role")
        content = m.get("content")

        if role in ("user", "assistant") and content:
            clean.append({
                "role": role,
                "content": str(content)
            })

    return clean

# --------------------------------------------------
# WORD EXPORT
# --------------------------------------------------

def word(history):
    history = normalize_history(history)

    doc = Document()
    doc.add_heading("Dynamo AI Research Report", 0)

    for m in history:
        role = "User" if m["role"] == "user" else "Dynamo AI"
        doc.add_heading(role, level=1)
        doc.add_paragraph(m["content"])

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=DynamoAI_Report.docx"}
    )

# --------------------------------------------------
# MARKDOWN → WORD EXPORT (editable .docx for a full paper/report)
# --------------------------------------------------

def _add_markdown_runs(paragraph, text):
    """Parse inline **bold** and *italic* markdown into docx runs."""
    tokens = re.split(r'(\*\*.+?\*\*|\*.+?\*)', text)
    for token in tokens:
        if not token:
            continue
        if token.startswith("**") and token.endswith("**") and len(token) > 4:
            paragraph.add_run(token[2:-2]).bold = True
        elif token.startswith("*") and token.endswith("*") and len(token) > 2:
            paragraph.add_run(token[1:-1]).italic = True
        else:
            paragraph.add_run(token)


def markdown_to_docx(markdown_text: str, title: str = "Research Paper"):
    """
    Convert a markdown-formatted paper/report into a properly structured,
    editable .docx — real Word headings, bold/italic runs, and bullet/numbered
    lists — instead of a flat text dump, so the user can open and edit it
    locally in Word/Google Docs.
    """
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    lines = (markdown_text or "").replace("\r\n", "\n").split("\n")

    for raw_line in lines:
        stripped = raw_line.strip()

        if not stripped:
            continue

        # Skip markdown horizontal rules
        if re.match(r'^(-{3,}|_{3,}|\*{3,})$', stripped):
            continue

        heading_match = re.match(r'^(#{1,6})\s+(.*)', stripped)
        if heading_match:
            level = len(heading_match.group(1))
            text = heading_match.group(2).strip().strip("*").strip()
            docx_level = max(0, min(level - 1, 4))  # md H1 -> docx Title, H2 -> Heading 1, etc.
            h = doc.add_heading(level=docx_level)
            _add_markdown_runs(h, text)
            continue

        bullet_match = re.match(r'^[-*•]\s+(.*)', stripped)
        if bullet_match:
            p = doc.add_paragraph(style="List Bullet")
            _add_markdown_runs(p, bullet_match.group(1))
            continue

        numbered_match = re.match(r'^\d+[\.\)]\s+(.*)', stripped)
        if numbered_match:
            p = doc.add_paragraph(style="List Number")
            _add_markdown_runs(p, numbered_match.group(1))
            continue

        p = doc.add_paragraph()
        _add_markdown_runs(p, stripped)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    safe_title = re.sub(r'[^\w\s-]', '', title or "Research_Paper").strip()[:60] or "Research_Paper"
    filename = re.sub(r'\s+', '_', safe_title) + ".docx"

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


# --------------------------------------------------
# POWERPOINT EXPORT
# --------------------------------------------------

def ppt(history):
    history = normalize_history(history)

    prs = Presentation()

    for m in history[-5:]:
        if m["role"] == "assistant":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Research Insight"
            slide.placeholders[1].text = m["content"][:700]

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=DynamoAI_Report.pptx"}
    )

# --------------------------------------------------
# PDF EXPORT
# --------------------------------------------------

def pdf(history):
    history = normalize_history(history)

    buf = io.BytesIO()
    pdf_doc = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()

    story = [
        Paragraph("Dynamo AI Intelligence Report", styles["Title"]),
        Spacer(1, 12)
    ]

    for m in history:
        role = "User:" if m["role"] == "user" else "Dynamo AI:"
        safe_text = html.escape(m["content"])

        story.append(
            Paragraph(f"<b>{role}</b> {safe_text}", styles["Normal"])
        )
        story.append(Spacer(1, 12))

    pdf_doc.build(story)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=DynamoAI_Report.pdf"}
    )
